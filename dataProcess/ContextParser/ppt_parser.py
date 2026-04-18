#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTX 解析器 - 符合 ContextParser 統一架構

此模組提供 PPTX 檔案解析功能，遵循 parser_design.md 定義的四大核心功能：
- preview: 讀取 PPTX 文件內容，返回純文字字串
- extract: 將文字內容分解為最小語義單元（unit_paras），識別結構字符，計算縮排層級
- segment: 將 unit_paras 合併成較大的語義段落（segments）
- chunk: 為每個 segment 提取關鍵字，生成最終的 chunk
"""

import os
import sys
import json
import logging
import re
import requests
import base64
import io
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
if PROJECT_ROOT not in sys.path:
    sys.path.insert(0, PROJECT_ROOT)

DATA_PROCESS_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if DATA_PROCESS_DIR not in sys.path:
    sys.path.insert(0, DATA_PROCESS_DIR)

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    MSO_SHAPE_TYPE = None
    PYTHON_PPTX_AVAILABLE = False

try:
    from package import LOGger
except Exception:
    class _FallbackCallableLogger:
        def __init__(self, name: str = 'ppt_parser'):
            self._logger = logging.getLogger(name)
            if not self._logger.handlers:
                self._logger.setLevel(logging.INFO)
                handler = logging.StreamHandler()
                handler.setFormatter(logging.Formatter('[%(levelname)s] %(message)s'))
                self._logger.addHandler(handler)

        def __call__(self, msg, *args, **kwargs):
            self._logger.info(msg)

    class _FallbackLOGger:
        FAIL = ''
        WARNING = ''
        OKBLUE = ''
        OKCYAN = ''

        @staticmethod
        def addloger(logfile: str = ''):
            return _FallbackCallableLogger('ppt_parser')

        @staticmethod
        def exception_process(x, logfile: str = '', **kwargs):
            logging.getLogger('ppt_parser').exception(x)

        @staticmethod
        def load_json(file_path: str):
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return json.load(f)
            except Exception:
                return {}

        @staticmethod
        def myDebuger(*args, **kwargs):
            return {}

    LOGger = _FallbackLOGger()

try:
    from dataProcess.ContextParser.context_parser import parse_keywords_from_text
    from dataProcess.ContextParser.context_parser import _merge_heading_tags_by_orders
    from dataProcess.ContextParser.context_parser import _build_heading_tags_by_order
    from dataProcess.ContextParser.context_parser import dedup_multi_prompts_by_llm
    from dataProcess.ContextParser.context_parser import filter_meaningless_tags
    from dataProcess.ContextParser.context_parser import analyze_images_via_batch_common
except Exception:
    try:
        from ContextParser.context_parser import parse_keywords_from_text
        from ContextParser.context_parser import _merge_heading_tags_by_orders
        from ContextParser.context_parser import _build_heading_tags_by_order
        from ContextParser.context_parser import dedup_multi_prompts_by_llm
        from ContextParser.context_parser import filter_meaningless_tags
        from ContextParser.context_parser import analyze_images_via_batch_common
    except Exception:
        try:
            from .context_parser import parse_keywords_from_text  # type: ignore
            from .context_parser import _merge_heading_tags_by_orders  # type: ignore
            from .context_parser import _build_heading_tags_by_order  # type: ignore
            from .context_parser import dedup_multi_prompts_by_llm  # type: ignore
            from .context_parser import filter_meaningless_tags  # type: ignore
            from .context_parser import analyze_images_via_batch_common  # type: ignore
        except Exception:
            parse_keywords_from_text = None
            _merge_heading_tags_by_orders = None
            _build_heading_tags_by_order = None
            dedup_multi_prompts_by_llm = None
            filter_meaningless_tags = None
            analyze_images_via_batch_common = None

if filter_meaningless_tags is None:
    def filter_meaningless_tags(tags):
        if not isinstance(tags, list):
            return []
        return [t for t in tags if t]

m_fn = os.path.basename(__file__).replace('.py', '')
m_debug = LOGger.myDebuger(stamps=[m_fn])
m_logfile = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'log', 'ppt_parser_%t.log')

def _setup_logger() -> logging.Logger:
    """設定日誌"""
    log_dir = os.path.join(os.path.dirname(__file__), 'log')
    os.makedirs(log_dir, exist_ok=True)
    
    logger = LOGger.addloger(logfile=os.path.join(log_dir, 'context_parser_%t.log'))
    logger.error = lambda x,*args,colora=LOGger.FAIL,**kwargs: logger(x,*args,**kwargs, colora=colora)
    logger.warning = lambda x,*args,colora=LOGger.WARNING,**kwargs: logger(x,*args,**kwargs, colora=colora)
    logger.info = lambda x,*args,**kwargs: logger(x,*args,**kwargs)
    logger.debug = lambda x,*args,colora=LOGger.OKBLUE,**kwargs: logger(x,*args,**kwargs, colora=colora)
    logger.exception = lambda x,logfile='',**kwargs: LOGger.exception_process(x,logfile=logfile,**kwargs)
    logger.summary = lambda x,*args,colora=LOGger.OKCYAN,**kwargs: logger(x,*args,**kwargs, colora=colora)
    return logger

m_logger = _setup_logger()
m_config_file = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'config.json')
m_config = LOGger.load_json(m_config_file) if os.path.exists(m_config_file) else {}

# ============================================================
# 輔助函數
# ============================================================

def _table_to_markdown(table_data: List[List[str]]) -> str:
    """
    將表格數據轉換為 Markdown 格式
    
    Args:
        table_data: 表格數據，每個元素是一行（列表）
    
    Returns:
        Markdown 格式的表格字符串
    """
    if not table_data or len(table_data) == 0:
        return ""
    
    # 確定最大列數
    max_cols = max(len(row) for row in table_data) if table_data else 0
    if max_cols == 0:
        return ""
    
    # 標準化所有行的列數（不足的補空字符串）
    normalized_data = []
    for row in table_data:
        normalized_row = list(row) if isinstance(row, list) else [str(row)]
        # 補齊列數
        while len(normalized_row) < max_cols:
            normalized_row.append("")
        normalized_data.append(normalized_row[:max_cols])
    
    # 轉換為 Markdown 表格
    markdown_lines = []
    
    # 表頭（第一行）
    if normalized_data:
        header_row = normalized_data[0]
        escaped_header = [str(cell).replace("|", "\\|").replace("\n", "<br>") for cell in header_row]
        markdown_lines.append("| " + " | ".join(escaped_header) + " |")
        
        # 分隔行
        separator = "| " + " | ".join(["---"] * max_cols) + " |"
        markdown_lines.append(separator)
        
        # 數據行（從第二行開始）
        for row in normalized_data[1:]:
            escaped_row = [str(cell).replace("|", "\\|").replace("\n", "<br>") for cell in row]
            markdown_lines.append("| " + " | ".join(escaped_row) + " |")
    
    return "\n".join(markdown_lines)


def _table_to_text(table_data: List[List[str]]) -> str:
    """
    將表格數據轉換為純文字格式
    
    Args:
        table_data: 表格數據
    
    Returns:
        純文字格式的表格字符串
    """
    if not table_data:
        return ""
    
    lines = []
    for row in table_data:
        row_text = " | ".join(str(cell) for cell in row)
        lines.append(row_text)
    
    return "\n".join(lines)


def _format_table_by_format(table_data: List[List[str]], format_type: str, table_id: int) -> str:
    """
    依照格式輸出表格內容
    
    Args:
        table_data: 表格資料
        format_type: 'markdown', 'text', 'placeholder'
        table_id: 表格 ID
    
    Returns:
        str: 表格文字或佔位符
    """
    if format_type == 'placeholder':
        return _make_table_placeholder(table_id)
    if format_type == 'markdown':
        return _table_to_markdown(table_data)
    return _table_to_text(table_data)

#
# PPTX 物件佔位符（表格 / 圖像）
#
TABLE_PLACEHOLDER = "[TABLE_PLACEHOLDER_{table_id}]"
IMAGE_PLACEHOLDER = "[IMAGE_PLACEHOLDER_{image_id}]"
PPT_TABLE_PLACEHOLDER = TABLE_PLACEHOLDER
PPT_IMAGE_PLACEHOLDER = IMAGE_PLACEHOLDER


def _make_table_placeholder(table_id: int) -> str:
    """建立表格佔位符字串"""
    return PPT_TABLE_PLACEHOLDER.format(table_id=table_id)


def _make_image_placeholder(image_id: int) -> str:
    """建立圖像佔位符字串"""
    return PPT_IMAGE_PLACEHOLDER.format(image_id=image_id)

def _guess_image_mime(
        *,
        ext: str = None,
        mime: str = None,
        filters=None,
        default_ext: str = "png",
    ):
    """
    統一推斷圖片副檔名與 MIME

    Priority:
    1) 明確指定 mime
    2) 明確指定 ext
    3) pptx stream Filter
    4) fallback default

    Returns:
        (ext, mime)
    """

    # 1) mime 優先
    if mime:
        mime = mime.lower()
        if mime == "image/jpeg":
            return "jpg", "image/jpeg"
        if mime == "image/png":
            return "png", "image/png"
        if mime in ("image/jp2", "image/jpeg2000"):
            return "jp2", mime

    # 2) ext 次之
    if ext:
        ext = ext.lower().lstrip(".")
        if ext in ("jpg", "jpeg"):
            return "jpg", "image/jpeg"
        if ext == "png":
            return "png", "image/png"
        if ext in ("jp2", "jpx"):
            return "jp2", "image/jp2"

    # 3) stream Filter
    if filters:
        f = str(filters)
        if "DCTDecode" in f:
            return "jpg", "image/jpeg"
        if "JPXDecode" in f:
            return "jp2", "image/jp2"
        if "FlateDecode" in f:
            # 通常是 raw bitmap，保守用 png
            return "png", "image/png"

    # 4) fallback
    if default_ext == "jpg":
        return "jpg", "image/jpeg"
    return "png", "image/png"

def _pack_image_bytes(image_bytes: bytes, ext: str, mime: str, to_base64: bool) -> Dict[str, Any]:
    """統一封裝輸出格式（ext/mime + base64 或 bytes）。"""
    data: Dict[str, Any] = {"ext": ext, "mime": mime}
    if to_base64:
        data["base64"] = base64.b64encode(image_bytes).decode("utf-8")
    else:
        data["bytes"] = image_bytes
    return data

def _extract_image_data_from_ppt(
        prs: Presentation,
        slide,
        shape,
        *,
        to_base64: bool = True,
        render_dpi: int = 200,
    ) -> Dict[str, Any]:
    """
    從 PPTX shape 提取圖像資料
    
    Args:
        prs: Presentation 物件
        slide: Slide 物件
        shape: Shape 物件（必須是 PICTURE 類型）
        to_base64: 是否轉換為 base64
        render_dpi: 渲染 DPI（未使用，保留以保持接口一致性）
    
    Returns:
        Dict[str, Any]: 包含圖像資料的字典
    """
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        m_logger.warning("[_extract_image_data_from_ppt] shape 不是 PICTURE 類型")
        return {}
    
    try:
        # 獲取圖像
        image = shape.image
        image_bytes = image.blob
        
        # 獲取圖像格式
        ext = image.ext  # 例如: 'png', 'jpeg', 'jpg'
        mime = image.content_type  # 例如: 'image/png', 'image/jpeg'
        
        # 推斷 ext 和 mime
        ext, mime = _guess_image_mime(ext=ext, mime=mime, default_ext="png")
        
        return _pack_image_bytes(image_bytes, ext, mime, to_base64)
        
    except Exception as e:
        m_logger.warning(f"[_extract_image_data_from_ppt] 提取圖像失敗: {e}", colora=LOGger.WARNING)
        return {}


# ============================================================
# 1. Preview（預覽）
# ============================================================
def preview(
        file_path: str,
        include_tables: bool = True,
        include_metadata: bool = False,
        separator: str = '\n\n',
        **kwargs
    ) -> str:
    """
    預覽 PPTX 文件內容
    
    Args:
        file_path: PPTX 檔案路徑
        include_tables: 是否包含表格（預設 True）
        include_metadata: 是否包含文件元數據（預設 False）
        separator: 段落分隔符號（預設 '\n\n'）
        **kwargs: 其他參數
    
    Returns:
        str: 純文字內容
    """
    if not PYTHON_PPTX_AVAILABLE:
        raise ImportError("python-pptx 未安裝，請執行: pip install python-pptx")
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"檔案不存在: {file_path}")
    
    m_logger.info(f"[preview] 開始預覽 PPTX 檔案: {file_path}")
    
    content_lines = []
    
    try:
        prs = Presentation(file_path)
        m_logger.debug(f"[preview] PPTX 總投影片數: {len(prs.slides)}", colora=LOGger.OKBLUE)
        
        # 提取元數據（可選）
        if include_metadata:
            metadata_lines = _extract_metadata_text_pptx(prs)
            if metadata_lines:
                content_lines.extend(metadata_lines)
                content_lines.append('')  # 空行分隔
        
        # 遍歷投影片
        for slide_idx, slide in enumerate(prs.slides, start=1):
            m_logger.debug(f"[preview] 處理第 {slide_idx} 張投影片", colora=LOGger.OKBLUE)
            
            slide_texts = []
            
            # 提取標題
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_texts.append(slide.shapes.title.text.strip())
            
            # 提取文字框內容
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text and text not in slide_texts:  # 避免重複標題
                        slide_texts.append(text)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE and include_tables:
                    # 提取表格
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)
                    if table_data:
                        table_text = _table_to_text(table_data)
                        slide_texts.append(f"\n[表格]\n{table_text}")
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # 圖像使用佔位符
                    slide_texts.append("[圖像]")
            
            if slide_texts:
                content_lines.append('\n'.join(slide_texts))
        
        m_logger.info(f"[preview] 預覽完成，共提取 {len(content_lines)} 個內容塊", colora=LOGger.OKCYAN)
        return separator.join(content_lines)
        
    except Exception as e:
        m_logger.error(f"[preview] 預覽失敗: {e}", colora=LOGger.FAIL)
        raise ValueError(f"無法讀取 PPTX 檔案: {e}")


def parse(
        file_path: str,
        include_tables: bool = True,
        include_metadata: bool = False,
        separator: str = '\n\n',
        **kwargs
    ) -> str:
    """
    解析 PPTX 文件內容（引用 preview）
    
    Args:
        file_path: PPTX 檔案路徑
        include_tables: 是否包含表格（預設 True）
        include_metadata: 是否包含文件元數據（預設 False）
        separator: 段落分隔符號（預設 '\n\n'）
        **kwargs: 其他參數
    
    Returns:
        str: 文件內容
    """
    return preview(
        file_path=file_path,
        include_tables=include_tables,
        include_metadata=include_metadata,
        separator=separator,
        **kwargs
    )


def process(
        file_path: str,
        include_tables: bool = True,
        table_format: str = 'markdown',
        segment_mode: str = 'structure',
        suitable_char_count: int = 100,
        max_char_multiplier: float = 1.5,
        min_char_multiplier: float = 0.5,
        separator: str = '\n',
        use_advanced_structure: bool = True,
        consider_page_boundary: bool = True,
        llm_provider: str = 'remote',
        llm_model: str = 'remote8b',
        llm_base_url: str = None,
        extract_kw_lbd: int = 10,
        include_images: bool = True,
        image_llm_provider: str = 'openai',
        image_llm_model: str = 'gpt4o_chat',
        image_placeholder: bool = True,
        **kwargs
    ) -> Dict[str, Any]:
    """
    一條龍處理流程：extract → segment → chunk（省略 preview 以節省流程）
    
    Args:
        file_path: 檔案路徑
        include_tables: 是否包含表格
        table_format: 表格格式（'markdown', 'text', 'html'）
        segment_mode: 分段模式（'char_count', 'structure'）
        suitable_char_count: 合適的字數
        max_char_multiplier: 最大字數倍數
        min_char_multiplier: 最小字數倍數
        separator: 段落之間的連接符號
        use_advanced_structure: 是否使用進階結構分析
        consider_page_boundary: 是否考慮頁面邊界
        llm_provider: LLM 提供者
        llm_model: LLM 模型
        llm_base_url: LLM API 的 base URL
        extract_kw_lbd: 提取關鍵字的最小字數閾值
        include_images: 是否包含圖像（預設 True）
        image_llm_provider: 圖像 LLM 提供者
        image_llm_model: 圖像 LLM 模型
        image_placeholder: 是否啟用圖像佔位符功能
        **kwargs: 其他參數
    
    Returns:
        Dict[str, Any]: 包含以下鍵的字典：
            - chunks: 最終的 chunk 結果列表
            - unit_paras: extract 階段的結果
            - segments: segment 階段的結果
            - metadata: 元數據
            - stats: 統計資訊
    """
    m_logger.info(f"開始處理檔案: {file_path}")
    
    # 階段 1: Extract
    m_logger.info("階段 1/3: Extract")
    extract_result = extract(
        file_path=file_path,
        include_tables=include_tables,
        table_format=table_format,
        include_images=include_images,
        image_placeholder=image_placeholder,
        **kwargs
    )
    unit_paras = extract_result['unit_paras']
    metadata = extract_result.get('metadata', {})
    m_logger.info(f"  Extract 完成，段落數: {len(unit_paras)}")
    
    # 階段 2: Segment
    m_logger.info("階段 2/3: Segment")
    segments = segment(
        unit_paras=unit_paras,
        table_format=table_format,
        include_images=include_images,
        suitable_char_count=suitable_char_count,
        segment_mode=segment_mode,
        max_char_multiplier=max_char_multiplier,
        min_char_multiplier=min_char_multiplier,
        separator=separator,
        use_advanced_structure=use_advanced_structure,
        metadata=metadata,
        consider_page_boundary=consider_page_boundary,
        llm_provider=llm_provider,
        llm_model=llm_model,
        llm_base_url=llm_base_url,
        image_llm_provider=image_llm_provider,
        image_llm_model=image_llm_model,
        **kwargs
    )
    m_logger.info(f"  Segment 完成，段落數: {len(segments)}")
    
    # 階段 3: Chunk
    m_logger.info("階段 3/3: Chunk")
    chunks = chunk(
        segments=segments,
        separator=separator,
        extract_kw_lbd=extract_kw_lbd,
        llm_provider=llm_provider,
        llm_model=llm_model,
        llm_base_url=llm_base_url,
        metadata=metadata,
        **kwargs
    )
    m_logger.info(f"  Chunk 完成，段落數: {len(chunks)}")
    
    # 統計資訊（content_length 從 unit_paras 計算，因已省略 preview）
    content_length = sum(len(p.get('unit_text', '')) for p in unit_paras)
    stats = {
        'file_path': file_path,
        'content_length': content_length,
        'unit_paras_count': len(unit_paras),
        'segments_count': len(segments),
        'chunks_count': len(chunks),
        'reduction_rate': f"{(1 - len(segments) / len(unit_paras)) * 100:.1f}%" if unit_paras else "N/A",
        'total_keywords': sum(len(c.get('multi_prompts', [])) for c in chunks),
        'chunks_with_keywords': sum(1 for c in chunks if c.get('multi_prompts', []))
    }
    
    m_logger.summary(f"處理完成: {file_path}")
    m_logger.summary(f"  段落數變化: {len(unit_paras)} → {len(segments)} → {len(chunks)}")
    m_logger.summary(f"  總關鍵字數: {stats['total_keywords']}")
    
    return {
        'chunks': chunks,
        'unit_paras': unit_paras,
        'segments': segments,
        'metadata': metadata,
        'stats': stats
    }


# ============================================================
# 2. Extract（提取）
# ============================================================

def _extract_metadata_text_pptx(prs: Presentation) -> List[str]:
    """使用 python-pptx 提取元數據文字"""
    try:
        lines = ["=== PPTX 元數據 ==="]
        core_props = prs.core_properties
        if core_props.title:
            lines.append(f"標題: {core_props.title}")
        if core_props.author:
            lines.append(f"作者: {core_props.author}")
        if core_props.subject:
            lines.append(f"主題: {core_props.subject}")
        if core_props.comments:
            lines.append(f"備註: {core_props.comments}")
        return lines
    except Exception as e:
        m_logger.debug(f"[_extract_metadata_text_pptx] 提取元數據失敗: {e}", colora=LOGger.OKBLUE)
        return []

def extract(
        file_path: str = None,
        text: str = None,
        include_tables: bool = True,
        table_format: str = 'markdown',
        to_base64: bool = True,
        use_bookmarks: bool = True,
        use_layout_analysis: bool = True,
        page_margin_left: float = 72.0,
        indent_unit_pt: float = 36.0,
        include_images: bool = True,
        image_placeholder: bool = True,
        **kwargs
    ) -> Dict[str, Any]:
    """
    提取 PPTX 文件的結構化資訊
    
    Args:
        file_path: PPTX 檔案路徑（與 text 二選一）
        text: 純文字內容（與 file_path 二選一，用於已讀取的內容）
        include_tables: 是否包含表格（預設 True）
        table_format: 表格格式（'markdown', 'text', 'placeholder'）
        use_bookmarks: 是否使用書籤識別標題（預設 True，PPTX 中對應 slide 標題）
        use_layout_analysis: 是否使用佈局分析（預設 True）
        page_margin_left: 頁面左邊距（pt，用於計算縮排）
        indent_unit_pt: 縮排單位（pt，預設 36pt = 1 級）
        include_images: 是否包含圖像
        image_placeholder: 是否使用圖像佔位符
        **kwargs: 其他參數
    
    Returns:
        Dict[str, Any]: 包含 unit_paras 和 metadata 的字典
    """
    if not PYTHON_PPTX_AVAILABLE:
        raise ImportError("python-pptx 未安裝，請執行: pip install python-pptx")
    
    if not file_path and not text:
        raise ValueError("必須提供 file_path 或 text 其中之一")
    
    # 如果提供 text，使用簡單解析
    if text and not file_path:
        return _extract_from_text(text, **kwargs)
    
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"檔案不存在: {file_path}")
    
    m_logger.info(f"[extract] 開始提取 PPTX 檔案: {file_path}")
    m_logger.debug(f"[extract] 參數：include_tables={include_tables}, include_images={include_images}, image_placeholder={image_placeholder}, table_format={table_format}, use_bookmarks={use_bookmarks}, use_layout_analysis={use_layout_analysis}", colora=LOGger.OKBLUE)
    
    # 提取書籤和元數據
    bookmarks = []
    metadata = {}
    bookmarks, metadata = _extract_bookmarks_and_metadata(file_path)
    m_logger.debug(f"[extract] 書籤數量: {len(bookmarks)}, 有元數據: {bool(metadata)}", colora=LOGger.OKBLUE)
    
    # 使用 python-pptx 提取文字和表格
    metadata.setdefault('tables', [])
    metadata.setdefault('images', [])
    unit_paras = []
    order = 0
    table_count = 0
    image_count = 0
    heading_count = 0
    
    try:
        prs = Presentation(file_path)
        m_logger.debug(f"[extract] PPTX 總投影片數: {len(prs.slides)}", colora=LOGger.OKBLUE)
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            m_logger.debug(f"[extract] 處理第 {slide_idx} 張投影片", colora=LOGger.OKBLUE)
            
            # 提取標題
            slide_title = None
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()
                heading_count += 1
                
                # 標題作為 heading
                unit_para = {
                    'unit_text': slide_title,
                    'indent_level': 0,  # 標題層級為 0
                    'order': order,
                    'structure_chars': [{
                        'type': 'heading',
                        'level': 1,
                        'slide_title': slide_title,
                        'slide': slide_idx
                    }],
                    'page_number': slide_idx
                }
                unit_paras.append(unit_para)
                order += 1
                m_logger.debug(f"[extract] 段落 {order-1}: 投影片標題 (Level 1), slide={slide_idx}, text_preview={slide_title[:50]}...", colora=LOGger.OKBLUE)
            
            # 提取文字框內容
            for shape in slide.shapes:
                # 跳過標題 shape（已經處理過）
                if shape == slide.shapes.title:
                    continue
                
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue
                    
                    # 計算縮排層級（根據文字框位置和段落層級）
                    indent_level = _calculate_indent_from_shape(shape, slide, page_margin_left, indent_unit_pt)
                    
                    # 檢查是否匹配書籤標題
                    bookmark_info = _match_bookmark(text, bookmarks, slide_idx) if use_bookmarks else None
                    
                    structure_chars = []
                    if bookmark_info:
                        heading_count += 1
                        structure_chars.append({
                            'type': 'heading',
                            'level': bookmark_info['level'],
                            'bookmark_title': bookmark_info['title'],
                            'slide': slide_idx
                        })
                        indent_level = bookmark_info['level'] - 1
                        m_logger.debug(f"[extract] 段落 {order}: 書籤標題 (Level {bookmark_info['level']}), indent_level={indent_level}, slide={slide_idx}, text_preview={text[:50]}...", colora=LOGger.OKBLUE)
                    else:
                        m_logger.debug(f"[extract] 段落 {order}: 普通段落, indent_level={indent_level}, slide={slide_idx}, text_preview={text[:50]}...", colora=LOGger.OKBLUE)
                    
                    unit_para = {
                        'unit_text': text,
                        'indent_level': indent_level,
                        'order': order,
                        'structure_chars': structure_chars,
                        'page_number': slide_idx
                    }
                    unit_paras.append(unit_para)
                    order += 1
                
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE and include_tables:
                    # 提取表格
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text.strip())
                        table_data.append(row_data)
                    
                    if table_data:
                        table_id = table_count
                        placeholder = _make_table_placeholder(table_id)
                        table_text = _format_table_by_format(table_data, table_format, table_id)
                        
                        # 構建表格資訊字典
                        table_info = {
                            'id': table_id,
                            'placeholder': placeholder,
                            'slide': slide_idx,
                            'rows': len(table_data),
                            'columns': len(table_data[0]) if table_data else 0,
                            'data': table_data,
                            'markdown': _table_to_markdown(table_data),
                            'text': _table_to_text(table_data)
                        }
                        metadata['tables'].append(table_info)
                        
                        unit_para = {
                            'unit_text': table_text or placeholder,
                            'indent_level': 0,
                            'order': order,
                            'structure_chars': [{
                                'type': 'table',
                                'table_id': table_id,
                                'slide': slide_idx,
                                'rows': len(table_data),
                                'columns': len(table_data[0]) if table_data else 0
                            }],
                            'page_number': slide_idx
                        }
                        unit_paras.append(unit_para)
                        order += 1
                        table_count += 1
                
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # 提取圖像
                    if include_images or image_placeholder:
                        image_id = image_count
                        placeholder = _make_image_placeholder(image_id)
                        
                        image_info = {
                            'id': image_id,
                            'placeholder': placeholder,
                            'slide': slide_idx
                        }
                        
                        # 如果啟用圖像提取，則提取圖像資料
                        if include_images:
                            image_data = _extract_image_data_from_ppt(prs, slide, shape, to_base64=to_base64)
                            if image_data:
                                image_info.update(image_data)
                        
                        metadata['images'].append(image_info)
                        
                        if image_placeholder:
                            unit_para = {
                                'unit_text': placeholder,
                                'indent_level': 0,
                                'order': order,
                                'structure_chars': [{
                                    'type': 'image',
                                    'image_id': image_id,
                                    'slide': slide_idx
                                }],
                                'page_number': slide_idx
                            }
                            unit_paras.append(unit_para)
                            order += 1
                        
                        image_count += 1
        
        # 更新元數據統計
        metadata['total_tables'] = table_count
        metadata['total_images'] = image_count
        metadata['image_count'] = image_count
        metadata['has_images'] = image_count > 0
        metadata['total_paragraphs'] = order
        metadata['has_headings'] = heading_count > 0
        metadata['heading_count'] = heading_count
        metadata['file_path'] = file_path
        
        m_logger.info(f"[extract] 提取完成: 總段落數={order}, 標題數={heading_count}, 表格數={table_count}, 圖像數={image_count}", colora=LOGger.OKCYAN)
        
        return {
            'unit_paras': unit_paras,
            'metadata': metadata
        }
        
    except Exception as e:
        m_logger.error(f"[extract] 提取失敗: {e}", colora=LOGger.FAIL)
        raise


def _extract_from_text(text: str, **kwargs) -> Dict[str, Any]:
    """從純文字提取（備用方案）"""
    lines = text.split('\n')
    unit_paras = []
    
    for i, line in enumerate(lines):
        line_text = line.strip()
        if not line_text:
            continue
        
        unit_para = {
            'unit_text': line_text,
            'indent_level': 0,
            'order': i,
            'structure_chars': []
        }
        unit_paras.append(unit_para)
    
    return {
        'unit_paras': unit_paras,
        'metadata': {}
    }


def _extract_bookmarks_and_metadata(file_path: str) -> Tuple[List[Dict], Dict]:
    """
    使用 PPTX 提取書籤和元數據
    
    Returns:
        Tuple[List[Dict], Dict]: (bookmarks, metadata)
    """
    bookmarks = []
    metadata = {}
    
    try:
        prs = Presentation(file_path)
        
        # 提取元數據
        core_props = prs.core_properties
        if core_props:
            metadata = {
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'keywords': core_props.keywords.split(',') if core_props.keywords else [],
                'created': str(core_props.created) if core_props.created else '',
                'modified': str(core_props.modified) if core_props.modified else '',
                'total_slides': len(prs.slides)
            }
        else:
            metadata = {'total_slides': len(prs.slides)}
        
        # 提取書籤（PPTX 中書籤對應 slide 標題）
        for slide_idx, slide in enumerate(prs.slides, start=1):
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()
                bookmarks.append({
                    'level': 1,  # PPTX 中 slide 標題通常為第 1 層
                    'title': title,
                    'slide': slide_idx
                })
        
        m_logger.debug(f"[_extract_bookmarks_and_metadata] 提取到 {len(bookmarks)} 個書籤", colora=LOGger.OKBLUE)
        
        metadata['has_bookmarks'] = len(bookmarks) > 0
        metadata['bookmark_count'] = len(bookmarks)
        
    except Exception as e:
        m_logger.warning(f"[_extract_bookmarks_and_metadata] 提取失敗: {e}", colora=LOGger.WARNING)
    
    return bookmarks, metadata


def _match_bookmark(line_text: str, bookmarks: List[Dict], slide_num: int) -> Optional[Dict]:
    """
    匹配文字與書籤
    
    Args:
        line_text: 文字內容
        bookmarks: 書籤列表
        slide_num: 當前投影片編號
    
    Returns:
        Optional[Dict]: 匹配的書籤資訊，如果沒有匹配則返回 None
    """
    if not bookmarks:
        return None
    
    # 清理文字以便比對
    clean_text = line_text.strip().lower()
    
    for bookmark in bookmarks:
        bookmark_title = bookmark['title'].strip().lower()
        bookmark_slide = bookmark['slide']
        
        # 檢查投影片編號是否匹配（允許前後 1 張的誤差）
        if abs(bookmark_slide - slide_num) <= 1:
            # 檢查文字是否匹配
            if clean_text == bookmark_title or bookmark_title in clean_text or clean_text in bookmark_title:
                return bookmark
    
    return None


def _calculate_indent_from_shape(shape, slide, page_margin_left: float, indent_unit_pt: float) -> int:
    """
    從 shape 位置計算縮排層級
    
    Args:
        shape: PPTX shape 物件
        slide: PPTX slide 物件
        page_margin_left: 頁面左邊距（pt）
        indent_unit_pt: 縮排單位（pt）
    
    Returns:
        int: 縮排層級
    """
    try:
        # 獲取 shape 的左邊界位置（以 EMU 為單位，1 inch = 914400 EMU）
        left_emu = shape.left
        left_pt = left_emu / 12700.0  # 轉換為 pt (1 pt = 12700 EMU)
        
        # 計算縮排
        indent_pt = left_pt - page_margin_left
        indent_level = int(indent_pt / indent_unit_pt)
        
        # 如果 shape 有文字框，檢查段落層級
        if shape.has_text_frame:
            # 檢查第一個段落的層級
            if shape.text_frame.paragraphs:
                first_para = shape.text_frame.paragraphs[0]
                if hasattr(first_para, 'level') and first_para.level is not None:
                    # PPTX 的 level 從 0 開始，我們也從 0 開始
                    indent_level = max(indent_level, first_para.level)
        
        return max(0, indent_level)
    except Exception:
        return 0

def display_metadata_image(metadata):
    metadata_display = {}
    for key, value in metadata.items():
        if key == 'images' and isinstance(value, list):
            # 只顯示圖像數量、格式和大小
            images_summary = []
            for img in value:
                if isinstance(img, dict):
                    img_info = {}
                    if 'base64' in img:
                        base64_str = img['base64']
                        img_info['base64'] = f"<base64: {len(base64_str)} chars>"
                    if 'mime' in img:
                        img_info['mime'] = img['mime']
                    images_summary.append(img_info)
                else:
                    images_summary.append(img)
            metadata_display[key] = images_summary
        else:
            metadata_display[key] = value
    return metadata_display

# ============================================================
# 3. Segment（分段）
# ============================================================

def segment(
        unit_paras: List[Dict[str, Any]],
        suitable_char_count: int = 500,
        segment_mode: str = 'structure',
        max_char_multiplier: float = 1.5,
        min_char_multiplier: float = 0.5,
        separator: str = '\n',
        use_advanced_structure: bool = True,
        consider_page_boundary: bool = True,
        metadata: Dict[str, Any] = None,
        table_format: str = 'markdown',
        enable_image_llm: bool = True,
        llm_provider: str = 'remote',
        llm_model: str = 'remote8b',
        llm_base_url: str = None,
        image_context_window: int = 200,
        max_images_per_batch: int = 50,
        image_prompt_template: str = None,
        image_llm_provider: str = 'openai',
        image_llm_model: str = 'gpt4o_chat',
        **kwargs
    ) -> List[Dict[str, Any]]:
    """
        將 unit_paras 合併成 segments
        
        Args:
            unit_paras: extract 返回的 unit_paras 列表
            suitable_char_count: 目標字數（預設 500）
            segment_mode: 分段模式（'structure', 'hybrid', 'size'）
            max_char_multiplier: 最大字數倍數（預設 1.5）
            min_char_multiplier: 最小字數倍數（預設 0.5）
            separator: 段落分隔符號（預設 '\n'）
            use_advanced_structure: 是否使用進階 structure 模式（預設 True）
                - True: 使用移植自 md_parser 的 _segment_with_structure（包含特殊類型往上合併邏輯；特殊類型含表格/圖像）
                - False: 使用舊版 _segment_by_structure
            metadata: 從 extract() 過來的元數據(會有圖像資訊)
            enable_image_llm: 是否啟用圖像 LLM
            consider_page_boundary: 是否考慮頁面邊界（預設 True）
            llm_provider: LLM 提供者
            llm_model: LLM 模型
            llm_base_url: LLM API 的 base URL
            image_context_window: 上下文範圍（字元數，預設 200）
            max_images_per_batch: 每批次最大圖像數（預設 50）
            image_prompt_template: 自訂 prompt 模板
            image_llm_provider: 圖像 LLM 提供者
            image_llm_model: 圖像 LLM 模型
            **kwargs: 其他參數
        
        Returns:
            List[Dict[str, Any]]: segments 列表
    """
    if not unit_paras:
        m_logger.warning("[segment] unit_paras 為空，返回空列表")
        return []

    heading_tags_by_order = _build_heading_tags_by_order(unit_paras) if _build_heading_tags_by_order else None
    
    m_logger.info(f"[segment] 開始分段: 輸入段落數={len(unit_paras)}, mode={segment_mode}, suitable_char_count={suitable_char_count}", colora=LOGger.OKCYAN)
    m_logger.debug(f"[segment] 分段參數: max_char_multiplier={max_char_multiplier}, min_char_multiplier={min_char_multiplier}, separator={repr(separator)}, consider_page_boundary={consider_page_boundary}", colora=LOGger.OKBLUE)
    
    if segment_mode == 'structure':
        if use_advanced_structure:
            m_logger.debug(f"[segment] 使用 structure 模式（進階版）: 基於 PPTX 結構（書籤、縮排）進行分段", colora=LOGger.OKBLUE)
            result = _segment_with_structure(
                unit_paras,
                separator,
                suitable_char_count,
                consider_page_boundary,
                heading_tags_by_order
            )
        else:
            m_logger.debug(f"[segment] 使用 structure 模式（基礎版）: 基於字數和縮排進行分段", colora=LOGger.OKBLUE)
            result = _segment_by_size(unit_paras, suitable_char_count, max_char_multiplier, min_char_multiplier, separator, heading_tags_by_order)
    elif segment_mode == 'hybrid':
        m_logger.warning(f"[segment] hybrid 模式尚未實作，改用 structure 模式", colora=LOGger.WARNING)
        result = _segment_with_structure(unit_paras, separator, suitable_char_count, consider_page_boundary, heading_tags_by_order) if use_advanced_structure else _segment_by_size(unit_paras, suitable_char_count, max_char_multiplier, min_char_multiplier, separator, heading_tags_by_order)
    else:  # 'size'
        m_logger.debug(f"[segment] 使用 size 模式: 純粹基於字數進行分段", colora=LOGger.OKBLUE)
        result = _segment_by_size(unit_paras, suitable_char_count, max_char_multiplier, min_char_multiplier, separator, heading_tags_by_order)
    
    m_logger.info(f"[segment] 分段完成: 輸出段落數={len(result)}, 減少率={(1 - len(result) / len(unit_paras)) * 100:.1f}%", colora=LOGger.OKCYAN)

    # 圖像分析（可選）
    m_logger.info(f"[segment] 圖像分析參數: enable_image_llm={enable_image_llm}, metadata.get('has_images')={metadata.get('has_images')}")
    if enable_image_llm and metadata and metadata.get('has_images'):
        m_logger.info("[segment] 開始圖像分析")
        image_analysis_result = _analyze_images_via_batch(
            segments=result,
            metadata=metadata,
            llm_provider=image_llm_provider,
            llm_model=image_llm_model,
            llm_base_url=llm_base_url,
            enable_image_llm=enable_image_llm,
            image_context_window=image_context_window,
            max_images_per_batch=max_images_per_batch,
            image_prompt_template=image_prompt_template,
        )
        m_logger.info(f"[segment] 圖像分析完成，共 {image_analysis_result.get('total_analyzed', 0)} 個結果")

    return result




def _segment_with_structure(
        filtered_paras: List[Dict[str, Any]], 
        separator: str, 
        suitable_char_count: int,
        consider_page_boundary: bool = True,
        heading_tags_by_order: Optional[Dict[int, List[str]]] = None
    ) -> List[Dict[str, Any]]:
    """
    按照結構分段（移植自 pdf_parser，適配 PPTX 結構）
    
    邏輯：
    - 遇到標題（slide 標題）時開始新段落
    - 合併段落時考慮 indent_level
    - 表格往上找最近的標題並合併
    - 考慮 slide 邊界（類似 PDF 的頁面邊界）
    - 後處理：確保每個 segment 都以標題開頭
    
    Args:
        filtered_paras: 段落列表
        separator: 段落之間的連接符號
        suitable_char_count: 目標字數
        consider_page_boundary: 是否考慮 slide 邊界
    
    Returns:
        List[Dict[str, Any]]: 分段後的段落列表
    """
    if not filtered_paras:
        return []
    
    segments = []
    current_segment = {
        'unit_texts': [],
        'indent_levels': [],
        'orders': [],
        'page_numbers': []
    }
    current_char_count = 0
    # 記錄已經被合併到 heading segment 的 orders，避免重複處理
    merged_to_heading_orders = set()
    
    m_logger.debug(f"[_segment_with_structure] 開始分段，suitable_char_count={suitable_char_count}, separator長度={len(separator)}", colora=LOGger.OKBLUE)
    
    for para in filtered_paras:
        unit_text = para.get('unit_text', '').strip()
        if not unit_text:
            continue
        
        indent_level = para.get('indent_level', 0)
        order = para.get('order', 0)
        structure_chars = para.get('structure_chars', [])
        page_number = para.get('page_number', 0)
        
        # 如果這個 order 已經被合併，跳過
        if order in merged_to_heading_orders:
            m_logger.debug(f"[_segment_with_structure] order={order} 已經被合併到 heading segment，跳過", colora=LOGger.OKBLUE)
            continue
        
        # 檢查是否有標題結構字符
        has_heading = any(sc.get('type') == 'heading' for sc in structure_chars)
        # 檢查是否是表格類型
        is_table = any(sc.get('type') == 'table' for sc in structure_chars)
        is_special_type = is_table
        para_char_count = len(unit_text)
        
        # 檢查是否跨 slide
        crosses_page = False
        if consider_page_boundary and current_segment['page_numbers']:
            current_page = current_segment['page_numbers'][-1]
            if page_number != current_page:
                crosses_page = True
        
        m_logger.debug(f"[_segment_with_structure] 處理 order={order}, indent_level={indent_level}, has_heading={has_heading}, is_table={is_table}, slide={page_number}, crosses_page={crosses_page}, para_char_count={para_char_count}, current_char_count={current_char_count}", colora=LOGger.OKBLUE)
        
        # 如果是表格類型，往上找最近的 heading
        if is_special_type:
            # 檢查當前段落是否有 heading
            current_has_heading = False
            if current_segment['unit_texts']:
                # 檢查當前段落中是否有 heading（通過檢查 orders 對應的 paragraphs）
                for prev_order in current_segment['orders']:
                    for prev_para in filtered_paras:
                        if prev_para.get('order') == prev_order:
                            prev_structure_chars = prev_para.get('structure_chars', [])
                            if any(sc.get('type') == 'heading' for sc in prev_structure_chars):
                                current_has_heading = True
                                break
                    if current_has_heading:
                        break
            
            if current_has_heading:
                # 當前段落有 heading，強制合併到當前段落
                m_logger.debug(f"[_segment_with_structure] 表格類型，當前段落有 heading，強制合併到當前段落", colora=LOGger.OKBLUE)
                current_segment['unit_texts'].append(unit_text)
                current_segment['indent_levels'].append(indent_level)
                current_segment['orders'].append(order)
                current_segment['page_numbers'].append(page_number)
                current_char_count = current_char_count + len(separator) + para_char_count
                continue
            else:
                # 當前段落沒有 heading，往上找最近的 heading
                found_heading_segment_idx = None
                heading_last_order = None
                for seg_idx in range(len(segments) - 1, -1, -1):
                    # 檢查該 segment 是否有 heading（通過檢查保存的 _has_heading 標記）
                    seg_has_heading = segments[seg_idx].get('_has_heading', False)
                    if seg_has_heading:
                        found_heading_segment_idx = seg_idx
                        heading_last_order = segments[seg_idx].get('_last_order', None)
                        break
                
                if found_heading_segment_idx is not None:
                    # 找到最近的 heading segment，合併中間的 unit_text
                    m_logger.debug(f"[_segment_with_structure] 表格類型，找到最近的 heading segment (index={found_heading_segment_idx})，heading_last_order={heading_last_order}, current_order={order}", colora=LOGger.OKBLUE)
                    
                    # 收集中間的 unit_text
                    middle_texts = []
                    middle_indent_levels = []
                    middle_orders = []
                    if heading_last_order is not None:
                        for mid_para in filtered_paras:
                            mid_order = mid_para.get('order', 0)
                            # 檢查是否在 heading_last_order 和 current_order 之間
                            if heading_last_order < mid_order < order:
                                mid_text = mid_para.get('unit_text', '').strip()
                                if mid_text:
                                    # 檢查這個 para 是否已經在 segments 中
                                    already_in_segments = False
                                    for seg in segments:
                                        seg_last_order = seg.get('_last_order', None)
                                        if seg_last_order is not None and mid_order <= seg_last_order:
                                            already_in_segments = True
                                            break
                                    
                                    # 檢查這個 para 是否在 current_segment 中
                                    in_current_segment = mid_order in current_segment['orders']
                                    
                                    if not already_in_segments:
                                        middle_texts.append(mid_text)
                                        middle_indent_levels.append(mid_para.get('indent_level', 0))
                                        middle_orders.append(mid_order)
                                        merged_to_heading_orders.add(mid_order)
                                        
                                        # 如果這個 order 在 current_segment 中，需要從 current_segment 中移除
                                        if in_current_segment:
                                            m_logger.debug(f"[_segment_with_structure] 發現中間 unit_text (order={mid_order}) 在 current_segment 中，從 current_segment 中移除", colora=LOGger.OKBLUE)
                                            # 找到並移除
                                            idx_to_remove = current_segment['orders'].index(mid_order)
                                            current_segment['unit_texts'].pop(idx_to_remove)
                                            current_segment['indent_levels'].pop(idx_to_remove)
                                            current_segment['orders'].pop(idx_to_remove)
                                            current_segment['page_numbers'].pop(idx_to_remove)
                                            # 更新字數
                                            current_char_count = sum(len(t) for t in current_segment['unit_texts']) + max(0, len(current_segment['unit_texts']) - 1) * len(separator)
                                        
                                        m_logger.debug(f"[_segment_with_structure] 發現中間 unit_text: order={mid_order}, text={mid_text[:50]}...", colora=LOGger.OKBLUE)
                    
                    # 合併中間的 unit_text 和當前的表到 heading segment
                    target_segment = segments[found_heading_segment_idx]
                    target_text = target_segment.get('unit_text', '')
                    
                    # 先合併中間的 unit_text
                    if middle_texts:
                        middle_combined = separator.join(middle_texts)
                        target_text = target_text + separator + middle_combined
                        m_logger.debug(f"[_segment_with_structure] 合併了 {len(middle_texts)} 個中間 unit_text 到 heading segment", colora=LOGger.OKBLUE)
                    
                    # 再合併當前的表格
                    new_text = target_text + separator + unit_text
                    segments[found_heading_segment_idx]['unit_text'] = new_text
                    
                    # 更新 indent_level（取最小的）
                    all_indent_levels = [target_segment.get('indent_level', 0)] + middle_indent_levels + [indent_level]
                    segments[found_heading_segment_idx]['indent_level'] = min(all_indent_levels)
                    
                    # 更新 _last_order
                    segments[found_heading_segment_idx]['_last_order'] = order
                    
                    # 標記當前的表格已經被合併
                    merged_to_heading_orders.add(order)
                    
                    m_logger.debug(f"[_segment_with_structure] ✓ 表格及中間 {len(middle_texts)} 個 unit_text 強制合併成功到 heading segment", colora=LOGger.OKBLUE)
                    continue
                else:
                    # 沒找到 heading，按照原來的邏輯處理（合併到當前段落或開始新段落）
                    m_logger.debug(f"[_segment_with_structure] 表格類型，沒找到 heading，按照原邏輯處理", colora=LOGger.OKBLUE)
        
        # 如果遇到標題，且當前段落不為空，檢查是否可以合併
        if has_heading and current_segment['unit_texts']:
            current_min_indent = min(current_segment['indent_levels'])
            
            # 如果考慮 slide 邊界且跨 slide，強制分段
            if crosses_page:
                m_logger.debug(f"[_segment_with_structure] 遇到標題且跨 slide，強制開始新段落", colora=LOGger.OKBLUE)
                segments.append({
                    'unit_text': separator.join(current_segment['unit_texts']),
                    'indent_level': min(current_segment['indent_levels']) if current_segment['indent_levels'] else 0,
                    'order': len(segments),
                    'orders': list(current_segment['orders']),
                    '_has_heading': any(any(sc.get('type') == 'heading' for sc in para.get('structure_chars', [])) for para in filtered_paras if para.get('order') in current_segment['orders']),
                    '_first_order': current_segment['orders'][0] if current_segment['orders'] else None
                })
                current_segment = {
                    'unit_texts': [unit_text],
                    'indent_levels': [indent_level],
                    'orders': [order],
                    'page_numbers': [page_number]
                }
                current_char_count = para_char_count
                continue
            
            # 檢查標題層級
            if indent_level > current_min_indent:
                new_char_count = current_char_count + len(separator) + para_char_count
                m_logger.warning(f"[_segment_with_structure] 遇到標題(indent_level={indent_level} > current_min_indent={current_min_indent})，檢查是否可以合併")
                m_logger.debug(f"[_segment_with_structure] 計算合併後字數: current_char_count={current_char_count} + separator({len(separator)}) + para_char_count({para_char_count}) = {new_char_count}", colora=LOGger.OKBLUE)
                if new_char_count <= suitable_char_count * 1.5:
                    m_logger.debug(f"[_segment_with_structure] ✓ 合併標題成功: order={order} 合併到當前段落（標題更深層級）", colora=LOGger.OKBLUE)
                    current_segment['unit_texts'].append(unit_text)
                    current_segment['indent_levels'].append(indent_level)
                    current_segment['orders'].append(order)
                    current_segment['page_numbers'].append(page_number)
                    current_char_count = new_char_count
                    continue
                else:
                    m_logger.warning(f"[_segment_with_structure] ✗ 標題字數太多，無法合併: new_char_count({new_char_count}) > suitable_char_count*1.5({suitable_char_count * 1.5})")
            elif indent_level == current_min_indent:
                # 相同層級的標題，如果字數允許也可以合併
                new_char_count = current_char_count + len(separator) + para_char_count
                m_logger.warning(f"[_segment_with_structure] 遇到標題(indent_level={indent_level} == current_min_indent={current_min_indent})，檢查是否可以合併")
                m_logger.debug(f"[_segment_with_structure] 計算合併後字數: current_char_count={current_char_count} + separator({len(separator)}) + para_char_count({para_char_count}) = {new_char_count}", colora=LOGger.OKBLUE)
                if new_char_count <= suitable_char_count * 1.5:
                    m_logger.debug(f"[_segment_with_structure] ✓ 合併標題成功: order={order} 合併到當前段落（標題相同層級）", colora=LOGger.OKBLUE)
                    current_segment['unit_texts'].append(unit_text)
                    current_segment['indent_levels'].append(indent_level)
                    current_segment['orders'].append(order)
                    current_segment['page_numbers'].append(page_number)
                    current_char_count = new_char_count
                    continue
                else:
                    m_logger.debug(f"[_segment_with_structure] ✗ 標題字數太多，無法合併: new_char_count({new_char_count}) > suitable_char_count*1.5({suitable_char_count * 1.5})", colora=LOGger.OKBLUE)
            
            # 開始新段落
            m_logger.debug(f"[_segment_with_structure] 遇到標題，開始新段落。indent_level={indent_level}, current_min_indent={current_min_indent}, 當前段落字數={current_char_count}", colora=LOGger.OKBLUE)
            segments.append({
                'unit_text': separator.join(current_segment['unit_texts']),
                'indent_level': min(current_segment['indent_levels']) if current_segment['indent_levels'] else 0,
                'order': len(segments),
                'orders': list(current_segment['orders']),
                '_has_heading': any(any(sc.get('type') == 'heading' for sc in para.get('structure_chars', [])) for para in filtered_paras if para.get('order') in current_segment['orders']),
                '_first_order': current_segment['orders'][0] if current_segment['orders'] else None
            })
            current_segment = {
                'unit_texts': [unit_text],
                'indent_levels': [indent_level],
                'orders': [order],
                'page_numbers': [page_number]
            }
            current_char_count = para_char_count
        else:
            # 檢查是否可以合併
            can_merge = True
            
            # 如果當前段落為空，直接添加
            if not current_segment['unit_texts']:
                m_logger.debug(f"[_segment_with_structure] 當前段落為空，直接添加 order={order}", colora=LOGger.OKBLUE)
                current_segment['unit_texts'].append(unit_text)
                current_segment['indent_levels'].append(indent_level)
                current_segment['orders'].append(order)
                current_segment['page_numbers'].append(page_number)
                current_char_count = para_char_count
            else:
                # 檢查是否跨 slide
                if crosses_page:
                    m_logger.warning(f"[_segment_with_structure] 跨 slide 邊界，開始新段落")
                    can_merge = False
                else:
                    current_min_indent = min(current_segment['indent_levels'])
                    
                    if indent_level > current_min_indent:
                        new_char_count = current_char_count + len(separator) + para_char_count
                        if new_char_count <= suitable_char_count * 1.5:
                            m_logger.debug(f"[_segment_with_structure] ✓ 合併成功: order={order} 合併到當前段落（更深層級）", colora=LOGger.OKBLUE)
                            current_segment['unit_texts'].append(unit_text)
                            current_segment['indent_levels'].append(indent_level)
                            current_segment['orders'].append(order)
                            current_segment['page_numbers'].append(page_number)
                            current_char_count = new_char_count
                        else:
                            can_merge = False
                    elif indent_level == current_min_indent:
                        new_char_count = current_char_count + len(separator) + para_char_count
                        if new_char_count <= suitable_char_count * 1.5:
                            m_logger.debug(f"[_segment_with_structure] ✓ 合併成功: order={order} 合併到當前段落（相同層級）", colora=LOGger.OKBLUE)
                            current_segment['unit_texts'].append(unit_text)
                            current_segment['indent_levels'].append(indent_level)
                            current_segment['orders'].append(order)
                            current_segment['page_numbers'].append(page_number)
                            current_char_count = new_char_count
                        else:
                            can_merge = False
                    else:
                        m_logger.warning(f"[_segment_with_structure] ✗ indent_level更淺，開始新段落")
                        can_merge = False
                
                if not can_merge:
                    m_logger.debug(f"[_segment_with_structure] 開始新段落，當前段落字數={current_char_count}", colora=LOGger.OKBLUE)
                    segments.append({
                        'unit_text': separator.join(current_segment['unit_texts']),
                        'indent_level': min(current_segment['indent_levels']) if current_segment['indent_levels'] else 0,
                        'order': len(segments),
                        'orders': list(current_segment['orders']),
                        '_has_heading': any(any(sc.get('type') == 'heading' for sc in para.get('structure_chars', [])) for para in filtered_paras if para.get('order') in current_segment['orders']),
                        '_last_order': max(current_segment['orders']) if current_segment['orders'] else order,
                        '_first_order': current_segment['orders'][0] if current_segment['orders'] else None
                    })
                    current_segment = {
                        'unit_texts': [unit_text],
                        'indent_levels': [indent_level],
                        'orders': [order],
                        'page_numbers': [page_number]
                    }
                    current_char_count = para_char_count
    
    # 添加最後一個段落
    if current_segment['unit_texts']:
        m_logger.debug(f"[_segment_with_structure] 添加最後一個段落，字數={current_char_count}", colora=LOGger.OKBLUE)
        segments.append({
            'unit_text': separator.join(current_segment['unit_texts']),
            'indent_level': min(current_segment['indent_levels']) if current_segment['indent_levels'] else 0,
            'order': len(segments),
            'orders': list(current_segment['orders']),
            '_has_heading': any(any(sc.get('type') == 'heading' for sc in para.get('structure_chars', [])) for para in filtered_paras if para.get('order') in current_segment['orders']),
            '_last_order': max(current_segment['orders']) if current_segment['orders'] else 0,
            '_first_order': current_segment['orders'][0] if current_segment['orders'] else None
        })
    
    m_logger.debug(f"[_segment_with_structure] 分段完成，共 {len(segments)} 個段落", colora=LOGger.OKBLUE)
    
    # 後處理：確保每個 segment 都以標題類型開頭（優先級高於字數限制）
    m_logger.debug(f"[_segment_with_structure] 開始後處理，確保每個 segment 都以標題類型開頭", colora=LOGger.OKBLUE)
    
    # 檢查是否有任何標題（通過檢查 _has_heading 標記）
    has_any_heading = any(seg.get('_has_heading', False) for seg in segments)
    
    if not has_any_heading:
        m_logger.debug(f"[_segment_with_structure] 沒有標題類型的 segment，跳過後處理", colora=LOGger.OKBLUE)
        # 清理臨時標記
        for seg in segments:
            if _merge_heading_tags_by_orders:
                seg['tags'] = _merge_heading_tags_by_orders(seg.get('orders', []), heading_tags_by_order)
            else:
                seg['tags'] = seg.get('tags', []) or []
            seg['tags'] = filter_meaningless_tags(seg['tags'])
            seg.pop('orders', None)
            seg.pop('_has_heading', None)
            seg.pop('_last_order', None)
            seg.pop('_first_order', None)
        return segments
    
    # 定義輔助函數：檢查 segment 是否以標題開頭（通過 _first_order）
    def segment_starts_with_heading(segment):
        """檢查 segment 是否以標題開頭（通過檢查第一個 para 的 structure_chars）"""
        first_order = segment.get('_first_order')
        if first_order is None:
            return False
        for para in filtered_paras:
            if para.get('order') == first_order:
                structure_chars = para.get('structure_chars', [])
                return any(sc.get('type') == 'heading' for sc in structure_chars)
        return False
    
    # 高效後處理：從前往後遍歷，將非標題 segment 合併到上一個 segment
    adjusted_segments = []
    
    for seg_idx, segment in enumerate(segments):
        unit_text = segment.get('unit_text', '')
        
        if segment_starts_with_heading(segment):
            # 以標題開頭的 segment，直接保留
            adjusted_segments.append(segment.copy())
            m_logger.debug(f"[_segment_with_structure] Segment {seg_idx} 以標題開頭，保留", colora=LOGger.OKBLUE)
        else:
            # 不以標題開頭的 segment，合併到上一個 segment
            if adjusted_segments:
                # 合併到上一個 segment
                last_segment = adjusted_segments[-1]
                last_text = last_segment.get('unit_text', '')
                combined_text = last_text + separator + unit_text if last_text else unit_text
                adjusted_segments[-1]['unit_text'] = combined_text
                # 更新 indent_level（取最小值）
                adjusted_segments[-1]['indent_level'] = min(
                    last_segment.get('indent_level', 0),
                    segment.get('indent_level', 0)
                )
                last_segment_orders = last_segment.get('orders', []) or []
                segment_orders = segment.get('orders', []) or []
                adjusted_segments[-1]['orders'] = [*last_segment_orders, *segment_orders]
                m_logger.debug(f"[_segment_with_structure] Segment {seg_idx} 不以標題開頭，合併到上一個 segment (index={len(adjusted_segments)-1})", colora=LOGger.OKBLUE)
            else:
                # 第一個 segment 不以標題開頭，也保留（但這種情況應該很少見）
                adjusted_segments.append(segment.copy())
                m_logger.debug(f"[_segment_with_structure] Segment {seg_idx} 是第一個且不以標題開頭，仍然保留", colora=LOGger.OKBLUE)
    
    # 清理臨時標記並重新編號
    for idx, seg in enumerate(adjusted_segments):
        seg['order'] = idx
        if _merge_heading_tags_by_orders:
            seg['tags'] = _merge_heading_tags_by_orders(seg.get('orders', []), heading_tags_by_order)
        else:
            seg['tags'] = seg.get('tags', []) or []
        seg['tags'] = filter_meaningless_tags(seg['tags'])
        seg.pop('orders', None)
        seg.pop('_has_heading', None)
        seg.pop('_last_order', None)
        seg.pop('_first_order', None)
    
    m_logger.debug(f"[_segment_with_structure] 後處理完成，從 {len(segments)} 個 segments 調整為 {len(adjusted_segments)} 個 segments", colora=LOGger.OKBLUE)
    
    return adjusted_segments


def _segment_by_size(
    unit_paras: List[Dict[str, Any]],
    suitable_char_count: int,
    max_char_multiplier: float,
    min_char_multiplier: float,
    separator: str,
    heading_tags_by_order: Optional[Dict[int, List[str]]] = None
) -> List[Dict[str, Any]]:
    """
    基於字數的分段策略（簡單版本）
    
    Args:
        unit_paras: 段落列表
        suitable_char_count: 目標字數
        max_char_multiplier: 最大字數倍數
        min_char_multiplier: 最小字數倍數
        separator: 分隔符
    
    Returns:
        List[Dict[str, Any]]: segments 列表
    """
    segments = []
    current_text_parts = []
    current_min_indent = None
    current_char_count = 0
    current_orders = []
    
    max_chars = int(suitable_char_count * max_char_multiplier)
    min_chars = int(suitable_char_count * min_char_multiplier)
    
    m_logger.debug(f"[_segment_by_size] 分段閾值: min_chars={min_chars}, max_chars={max_chars}, suitable_char_count={suitable_char_count}", colora=LOGger.OKBLUE)
    
    for unit in unit_paras:
        unit_text = unit['unit_text']
        indent_level = unit['indent_level']
        unit_order = unit.get('order', 0)
        
        if not current_text_parts:
            current_text_parts.append(unit_text)
            current_char_count = len(unit_text)
            current_min_indent = indent_level
            current_orders = [unit_order]
            continue
        
        # 檢查是否應該分段
        if current_char_count + len(unit_text) > max_chars and current_char_count >= min_chars:
            # 輸出當前 segment
            segments.append({
                'unit_text': separator.join(current_text_parts),
                'indent_level': current_min_indent if current_min_indent is not None else 0,
                'order': len(segments),
                'orders': list(current_orders)
            })
            current_text_parts = [unit_text]
            current_char_count = len(unit_text)
            current_min_indent = indent_level
            current_orders = [unit_order]
        else:
            # 合併
            current_text_parts.append(unit_text)
            current_char_count += len(unit_text)
            current_min_indent = indent_level if current_min_indent is None else min(current_min_indent, indent_level)
            current_orders.append(unit_order)
    
    # 添加最後一個 segment
    if current_text_parts:
        segments.append({
            'unit_text': separator.join(current_text_parts),
            'indent_level': current_min_indent if current_min_indent is not None else 0,
            'order': len(segments),
            'orders': list(current_orders)
        })
    
    for seg in segments:
        if _merge_heading_tags_by_orders:
            seg['tags'] = _merge_heading_tags_by_orders(seg.get('orders', []), heading_tags_by_order)
        else:
            seg['tags'] = seg.get('tags', []) or []
        seg['tags'] = filter_meaningless_tags(seg['tags'])
        seg.pop('orders', None)
    
    return segments

# ============================================================
# 4. Chunk（分塊）
# ============================================================

def chunk(
        segments: List[Dict[str, Any]],
        extract_kw_lbd: int = 10,
        llm_provider: str = 'remote',
        llm_model: str = 'remote8b',
        llm_base_url: str = None,
        metadata: Dict[str, Any] = None,
        **kwargs
    ) -> List[Dict[str, Any]]:
    """
    為 segments 提取關鍵字，生成 chunks
    
    Args:
        segments: segment 返回的 segments 列表
        separator: 段落之間的連接符號
        extract_kw_lbd: 提取關鍵字的最小字數閾值
        llm_provider: LLM 提供者
        llm_model: LLM 模型
        llm_base_url: LLM API 的 base URL
        metadata: extract 返回的 metadata
        **kwargs: 其他參數
    
    Returns:
        List[Dict[str, Any]]: chunks 列表
    """
    if not segments:
        return []
    
    # 獲取 LLM API base URL
    if not llm_base_url:
        llm_base_url = m_config.get('llm', {}).get('base_url', 'http://10.1.3.127:7017')
    pptx_metadata = metadata or {}

    table_map = {str(t.get('id')): t for t in pptx_metadata.get('tables', []) or [] if t.get('id') is not None}
    image_map = {str(i.get('id')): i for i in pptx_metadata.get('images', []) or [] if i.get('id') is not None}
    table_pattern = re.compile(r"\[(?:PPT_TABLE|TABLE_PLACEHOLDER)_(\d+)\]")
    image_pattern = re.compile(r"\[(?:PPT_IMAGE|IMAGE_PLACEHOLDER)_(\d+)\]")
    
    # 載入關鍵字提取提示詞配置
    prompt_file = os.path.join(os.path.dirname(os.path.dirname(os.path.dirname(__file__))), 'prompt', 'extract_keywords.json')
    try:
        with open(prompt_file, 'r', encoding='utf-8') as f:
            kw_config = json.load(f)
        system_prompt = kw_config.get('system_prompt', '')
        user_prompt_template = kw_config.get('user_prompt_template', '')
        generation_config = kw_config.get('generation_config', {})
    except Exception as e:
        LOGger.exception_process(e, logfile='', stamps=['ppt_parser','chunk'])
        m_logger.error(f"無法載入關鍵字提取提示詞配置: {e}，將跳過關鍵字提取")
        # 如果無法載入配置，返回不包含關鍵字的結果
        results = []
        for idx, seg in enumerate(segments):
            unit_text = seg.get('unit_text', '').strip()
            chunk_item = {
                'call_prompt': unit_text,
                'indent_level': seg.get('indent_level', 0),
                'order': seg.get('order', idx),
                'multi_prompts': [t for t in (seg.get('tags', []) or []) if t],
                'tags': seg.get('tags', [])
            }
            _add_pptx_metadata(chunk_item, pptx_metadata)
            results.append(chunk_item)
        return results
    
    # 準備需要提取關鍵字的段落（字數超過閾值）
    segments_to_extract = []
    segment_indices = []
    
    for idx, seg in enumerate(segments):
        unit_text = seg.get('unit_text', '').strip()
        if len(unit_text) >= extract_kw_lbd:
            segments_to_extract.append(unit_text)
            segment_indices.append(idx)
    
    m_logger.info(f"[chunk] 準備提取關鍵字: 總段落數={len(segments)}, 需要提取的段落數={len(segments_to_extract)}", colora=LOGger.OKCYAN)
    # 批量提取關鍵字
    keywords_list = []
    if segments_to_extract:
        m_logger.info(f"批次提取 {len(segments_to_extract)} 個段落的關鍵字...")
        
        # 構建所有 prompt
        prompts = []
        for content in segments_to_extract:
            user_prompt = user_prompt_template.format(content=content)
            # 合併 system_prompt 和 user_prompt（如果 LLM API 需要）
            full_prompt = f"{system_prompt}\n\n{user_prompt}" if system_prompt else user_prompt
            prompts.append(full_prompt)
        
        # 批量調用 LLM
        try:
            batch_chat_url = f"{llm_base_url.rstrip('/')}/chat/batch"
            payload = {
                'prompts': prompts,
                'provider': llm_provider,
                'model': llm_model,
                'max_tokens': generation_config.get('max_new_tokens', 150),
                'temperature': generation_config.get('temperature', 0.3),
                'system_prompt': system_prompt if system_prompt else None,
                'parallel': True,
                'max_batch_size': 190  # 與 docx_parser 一致
            }
            
            response = requests.post(batch_chat_url, json=payload, timeout=300)
            response.raise_for_status()
            
            if response.status_code == 200:
                batch_result = response.json()
                batch_items = batch_result.get('results', [])
                
                # 解析每個結果（使用 parse_keywords_from_text 進行清理和解析）
                max_keywords = generation_config.get('max_keywords', 2)
                min_keyword_length = generation_config.get('min_keyword_length', 1)
                
                for item_idx, item in enumerate(batch_items):
                    keywords = []
                    if item.get('error'):
                        m_logger.warning(f"提取關鍵字時發生錯誤: {item.get('error')}")
                    else:
                        result_obj = item.get('result', {})
                        output = result_obj.get('output', '').strip()
                        
                        # 使用 parse_keywords_from_text 解析關鍵字（支援多種格式）
                        if parse_keywords_from_text:
                            keywords = parse_keywords_from_text(
                                output,
                                max_keywords=max_keywords,
                                min_keyword_length=min_keyword_length
                            )
                        else:
                            # 備用解析方法
                            keywords = _parse_keywords_fallback(output, max_keywords)
                    
                    # 去重：保持順序的唯一化
                    keywords = list(dict.fromkeys(keywords))[:max_keywords]
                    keywords_list.append(keywords)
            else:
                m_logger.warning(f"LLM API 返回錯誤狀態碼: {response.status_code}")
                keywords_list = [[] for _ in segments_to_extract]
                
        except Exception as e:
            m_logger.error(f"批量提取關鍵字失敗: {e}")
            keywords_list = [[] for _ in segments_to_extract]
    else:
        m_logger.warning("沒有段落需要提取關鍵字（字數不足）")
    
    # 構建結果
    results = []
    kw_idx = 0
    pptx_metadata = metadata or {}
    
    for idx, seg in enumerate(segments):
        unit_text = seg.get('unit_text', '').strip()
        
        # 如果這個段落需要提取關鍵字
        if idx in segment_indices:
            keywords = keywords_list[kw_idx] if kw_idx < len(keywords_list) else []
            kw_idx += 1
            keywords = filter_meaningless_tags(keywords)
        else:
            keywords = []
        
        tags = seg.get('tags', []) or []
        merged_prompts = [p for p in list(dict.fromkeys([*keywords, *tags])) if p]
        chunk_item = {
            'call_prompt': unit_text,
            'indent_level': seg.get('indent_level', 0),
            'order': seg.get('order', idx),
            'multi_prompts': merged_prompts,
            'tags': tags
        }
        _add_pptx_metadata(chunk_item, pptx_metadata)
        results.append(chunk_item)

    # LLM 批量去重 multi_prompts（在嚴格去重後）
    if dedup_multi_prompts_by_llm:
        try:
            deduped_lists, metas = dedup_multi_prompts_by_llm(
                [r.get('multi_prompts', []) for r in results],
                config=m_config,
                llm_base_url=llm_base_url,
                llm_provider=llm_provider,
                llm_model=llm_model,
                logger=m_logger,
                return_meta=True
            )
            for r, mp, meta in zip(results, deduped_lists, metas):
                r['multi_prompts'] = mp
                r['multi_prompts_dedup'] = meta
        except Exception as e:
            m_logger.warning(f"multi_prompts 去重失敗，保留原結果: {e}")
            for r in results:
                r['multi_prompts_dedup'] = {'success': False, 'reason': 'llm_failed'}
    else:
        for r in results:
            r['multi_prompts_dedup'] = {'success': False, 'reason': 'dedup_not_available'}

    m_logger.info(f"[chunk] Chunk 完成，共 {len(results)} 個 chunks", colora=LOGger.OKCYAN)
    return results


def _add_pptx_metadata(chunk: Dict[str, Any], metadata: Dict[str, Any]) -> None:
    """
    為 chunk 添加 PPTX 特定元數據

    Args:
        chunk: 要添加元數據的 chunk
        metadata: PPTX 的 metadata
    """
    if metadata.get('title'):
        chunk['ppt_title'] = metadata['title']
    if metadata.get('author'):
        chunk['ppt_author'] = metadata['author']
    if metadata.get('subject'):
        chunk['ppt_subject'] = metadata['subject']


def _parse_keywords_fallback(output: str, max_keywords: int = 5) -> List[str]:
    """備用關鍵字解析方法（當 parse_keywords_from_text 不可用時）"""
    keywords = []
    
    # 嘗試直接解析 JSON
    try:
        parsed = json.loads(output)
        if isinstance(parsed, list):
            keywords = [str(kw) for kw in parsed if kw]
            return keywords[:max_keywords]
    except:
        pass
    
    # 嘗試從文字中提取 JSON 陣列
    match = re.search(r'\[.*?\]', output, re.DOTALL)
    if match:
        try:
            parsed = json.loads(match.group(0))
            if isinstance(parsed, list):
                keywords = [str(kw) for kw in parsed if kw]
                return keywords[:max_keywords]
        except:
            pass
    
    # 嘗試提取引號包圍的詞
    keywords = re.findall(r'["\']([^"\']+)["\']', output)
    if keywords:
        return keywords[:max_keywords]
    
    return []
# ============================================================
# 主程序（用於測試）
# ============================================================



def _analyze_images_via_batch(
        segments: List[Dict[str, Any]],
        metadata: Dict[str, Any],
        llm_provider: str = 'openai',
        llm_model: str = 'gpt4o_chat',
        llm_base_url: str = None,
        enable_image_llm: bool = True,
        image_context_window: int = 200,
        max_images_per_batch: int = 50,
        image_prompt_template: str = None,
        **kwargs
    ) -> Dict[str, Any]:
    """
    Analyze images via batch LLM using placeholders and context.
    """
    if analyze_images_via_batch_common is None:
        m_logger.warning("[_analyze_images_via_batch] shared helper unavailable")
        return {}
    return analyze_images_via_batch_common(
        segments=segments,
        metadata=metadata,
        llm_provider=llm_provider,
        llm_model=llm_model,
        llm_base_url=llm_base_url,
        enable_image_llm=enable_image_llm,
        image_context_window=image_context_window,
        max_images_per_batch=max_images_per_batch,
        image_prompt_template=image_prompt_template,
        placeholder_pattern=r'\[(?:PPT_IMAGE|IMAGE_PLACEHOLDER)_(\d+)\]',
        placeholder_replacements=['[IMAGE_PLACEHOLDER_{image_id}]', '[PPT_IMAGE_PLACEHOLDER_{image_id}]', '[PPT_IMAGE_{image_id}]'],
        text_keys=['unit_text', 'call_prompt', 'text'],
        config=m_config,
        logger=m_logger
    )


if __name__ == '__main__':
    # 簡單測試
    import sys
    
    if len(sys.argv) < 2:
        print("用法: python ppt_parser.py <ppt_file_path>")
        sys.exit(1)
    
    ppt_file = sys.argv[1]
    
    print("=== Preview ===")
    content = preview(ppt_file)
    print(f"內容長度: {len(content)} 字元")
    print(f"前 500 字元:\n{content[:500]}\n")
    
    print("=== Extract ===")
    result = extract(ppt_file)
    unit_paras = result['unit_paras']
    metadata = result['metadata']
    print(f"段落數: {len(unit_paras)}")
    print(f"元數據: {metadata}")
    print(f"前 3 個段落:")
    for i, para in enumerate(unit_paras[:3]):
        print(f"  [{i}] Level {para['indent_level']}: {para['unit_text'][:60]}...")
    
    print("\n=== Segment ===")
    segments = segment(unit_paras, suitable_char_count=500)
    print(f"Segment 數: {len(segments)}")
    print(f"前 3 個 segments:")
    for i, seg in enumerate(segments[:3]):
        print(f"  [{i}] Level {seg['indent_level']}: {seg['unit_text'][:60]}...")
    
    print("\n=== Chunk ===")
    chunks = chunk(segments, metadata=metadata)
    print(f"Chunk 數: {len(chunks)}")
    print(f"前 3 個 chunks:")
    for i, c in enumerate(chunks[:3]):
        print(f"  [{i}] Keywords: {c['multi_prompts']}, Chunk: {c['call_prompt'][:60]}...")
