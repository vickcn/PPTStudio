#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any, Optional

from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt


API_BASE = "http://10.1.3.127:6414"
DEFAULT_DPI = 96
TIMEOUT = 3600
ML_HOME_PREFIX = "C:\\ML_HOME\\"
ML_HOME_UNC_PREFIX = r"\\10.1.3.127\ml_home\\"


if __package__ in {None, ""}:
    from ppt_stdio import (
        add_image,
        add_shape,
        add_text,
        get_info,
        new,
        open_presentation,
        save,
        set_slide_background_color,
        set_slide_background_image,
        _apply_text_frame_layout,
    )
    from layout_engine import order_shapes_for_build
else:
    from .ppt_stdio import (
        add_image,
        add_shape,
        add_text,
        get_info,
        new,
        open_presentation,
        save,
        set_slide_background_color,
        set_slide_background_image,
        _apply_text_frame_layout,
    )
    from .layout_engine import order_shapes_for_build


def read_json(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8-sig"))


def json_slide_to_api_slide(value: Any) -> int:
    return max(to_int(value, default=1) - 1, 0)


def resolve_ml_home_path(value: str) -> Optional[Path]:
    candidate = Path(value)
    if candidate.exists():
        return candidate

    normalized = str(value)
    if normalized.startswith(ML_HOME_PREFIX):
        alt = Path(ML_HOME_UNC_PREFIX + normalized[len(ML_HOME_PREFIX):].lstrip("\\/"))
        if alt.exists():
            return alt
    return None


def emu_to_px(inches: float, dpi: int = DEFAULT_DPI) -> int:
    return int(round(inches * dpi))


def to_int(value: Any, default: int = 0) -> int:
    try:
        return int(round(float(value)))
    except Exception:
        return default


def pick_first(seq: list[Any], default: Any = None) -> Any:
    return seq[0] if seq else default


RUN_STYLE_KEYS = (
    "font_name",
    "font_size_pt",
    "effective_font_size_pt",
    "bold",
    "italic",
    "font_color",
    "latin_font_name",
    "east_asian_font_name",
    "complex_script_font_name",
)


def _paragraphs_joined_text(paragraphs: list[dict[str, Any]]) -> str:
    return "\n".join(str(paragraph.get("text") or "") for paragraph in paragraphs).replace("\x0b", "\n")


def _shape_text_matches_paragraphs(shape: dict[str, Any], paragraphs: list[dict[str, Any]]) -> bool:
    shape_text = str(shape.get("text") or "").replace("\x0b", "\n").strip()
    para_text = _paragraphs_joined_text(paragraphs).strip()
    return shape_text == para_text


def _merge_run_style(text_run: dict[str, Any], style_run: dict[str, Any]) -> dict[str, Any]:
    merged = dict(text_run)
    for key in RUN_STYLE_KEYS:
        if style_run.get(key) is not None and merged.get(key) is None:
            merged[key] = style_run.get(key)
    return merged


def _merge_paragraph_styles(
    text_paragraphs: list[dict[str, Any]],
    style_paragraphs: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    merged_paragraphs: list[dict[str, Any]] = []
    for idx, paragraph in enumerate(text_paragraphs):
        merged = dict(paragraph)
        style_para = style_paragraphs[idx] if idx < len(style_paragraphs) else {}
        text_runs = paragraph.get("runs") or []
        style_runs = style_para.get("runs") or []
        if text_runs:
            merged_runs: list[dict[str, Any]] = []
            for run_idx, run in enumerate(text_runs):
                style_run = style_runs[run_idx] if run_idx < len(style_runs) else (style_runs[0] if style_runs else {})
                merged_runs.append(_merge_run_style(run, style_run))
            merged["runs"] = merged_runs
        merged_paragraphs.append(merged)
    return merged_paragraphs


def _paragraphs_from_shape_text(shape: dict[str, Any], style_paragraphs: list[dict[str, Any]]) -> list[dict[str, Any]]:
    shape_text = str(shape.get("text") or "").replace("\x0b", "\n")
    lines = shape_text.split("\n")
    shape_paragraphs = shape.get("paragraphs") if isinstance(shape.get("paragraphs"), list) else []
    result: list[dict[str, Any]] = []
    for idx, line in enumerate(lines):
        if idx < len(shape_paragraphs):
            paragraph = dict(shape_paragraphs[idx])
            paragraph["text"] = line
            runs = paragraph.get("runs") or []
            if runs:
                paragraph["runs"] = [dict(runs[0])]
                paragraph["runs"][0]["text"] = line
            result.append(paragraph)
            continue
        if idx < len(style_paragraphs):
            paragraph = dict(style_paragraphs[idx])
            paragraph["text"] = line
            runs = paragraph.get("runs") or [{"text": line}]
            paragraph["runs"] = [dict(runs[0])]
            paragraph["runs"][0]["text"] = line
            result.append(paragraph)
            continue
        result.append({"text": line, "runs": [{"text": line}]})
    return result


def first_run_style(shape: dict[str, Any]) -> dict[str, Any]:
    paragraphs = text_paragraphs(shape)
    for paragraph in paragraphs:
        for run in paragraph.get("runs", []):
            if any(run.get(key) is not None for key in ("font_name", "font_size_pt", "bold", "italic")):
                return run
    return {}


def infer_text_style(shape: dict[str, Any]) -> dict[str, Any]:
    run = first_run_style(shape)
    font_summary = shape.get("font_summary") or {}
    font_names = font_summary.get("font_names") or []
    font_sizes = font_summary.get("font_sizes_pt") or []

    font_size = run.get("font_size_pt")
    if font_size is None:
        font_size = pick_first(font_sizes)
    font_size = to_int(font_size, default=18)

    font_name = run.get("font_name") or pick_first(font_names)
    bold = run.get("bold")
    italic = run.get("italic")

    return {
        "font_size": font_size,
        "font_name": font_name,
        "bold": bool(bold) if bold is not None else False,
        "italic": bool(italic) if italic is not None else False,
    }


LIST_LINE_RE = re.compile(r"^\d+\.\s")


def is_list_style_text(shape: dict[str, Any]) -> bool:
    text = str(shape.get("text") or "").replace("\x0b", "\n")
    if "\n" in text:
        return True
    paragraphs = text_paragraphs(shape)
    if len(paragraphs) > 1:
        return True
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if LIST_LINE_RE.match(stripped):
            return True
        if stripped[0] in "-•·●○":
            return True
    return False


def infer_text_frame_layout(shape: dict[str, Any]) -> dict[str, Any]:
    text_frame = shape.get("text_frame")
    layout: dict[str, Any] = {}
    if isinstance(text_frame, dict):
        if "word_wrap" in text_frame:
            layout["word_wrap"] = text_frame.get("word_wrap")
        if "auto_fit" in text_frame:
            layout["auto_fit"] = text_frame.get("auto_fit")
    if is_list_style_text(shape):
        layout["word_wrap"] = True
    return layout


def text_paragraphs(shape: dict[str, Any]) -> list[dict[str, Any]]:
    shape_text = str(shape.get("text") or "").replace("\x0b", "\n")
    shape_paragraphs = shape.get("paragraphs") if isinstance(shape.get("paragraphs"), list) else []
    font_detail = shape.get("font_detail") or {}
    fd_paragraphs = font_detail.get("paragraphs") if isinstance(font_detail.get("paragraphs"), list) else []

    if shape_paragraphs:
        if _shape_text_matches_paragraphs(shape, shape_paragraphs):
            if fd_paragraphs:
                return _merge_paragraph_styles(shape_paragraphs, fd_paragraphs)
            return shape_paragraphs
        if shape_text.strip():
            return _paragraphs_from_shape_text(shape, fd_paragraphs)

    if fd_paragraphs and shape_text.strip():
        fd_text = _paragraphs_joined_text(fd_paragraphs).strip()
        if fd_text != shape_text.strip():
            return _paragraphs_from_shape_text(shape, fd_paragraphs)

    if fd_paragraphs:
        return fd_paragraphs
    return shape_paragraphs


def has_run_level_text_style(shape: dict[str, Any]) -> bool:
    paragraphs = text_paragraphs(shape)
    runs: list[dict[str, Any]] = []
    for paragraph in paragraphs:
        runs.extend(paragraph.get("runs", []) or [])
    if len(paragraphs) > 1 or len(runs) > 1:
        return True

    style_keys = ("font_name", "font_size_pt", "effective_font_size_pt", "bold", "italic", "font_color")
    return any(any(run.get(key) is not None for key in style_keys) for run in runs)


def apply_run_style(pptx_run: Any, run_spec: dict[str, Any], default_style: dict[str, Any]) -> None:
    font_size = run_spec.get("font_size_pt")
    if font_size is None:
        font_size = run_spec.get("effective_font_size_pt")
    if font_size is None:
        font_size = default_style["font_size"]

    font_name = run_spec.get("font_name") or default_style.get("font_name")
    bold = run_spec.get("bold")
    italic = run_spec.get("italic")

    pptx_run.font.size = Pt(to_int(font_size, default=default_style["font_size"]))
    pptx_run.font.bold = default_style["bold"] if bold is None else bool(bold)
    pptx_run.font.italic = default_style["italic"] if italic is None else bool(italic)
    if font_name:
        pptx_run.font.name = str(font_name)

    set_run_typeface(pptx_run, "latin", run_spec.get("latin_font_name") or font_name)
    set_run_typeface(pptx_run, "ea", run_spec.get("east_asian_font_name") or font_name)
    set_run_typeface(pptx_run, "cs", run_spec.get("complex_script_font_name") or font_name)

    font_color = run_spec.get("font_color")
    if isinstance(font_color, list) and len(font_color) == 3:
        pptx_run.font.color.rgb = RGBColor(
            to_int(font_color[0], 0),
            to_int(font_color[1], 0),
            to_int(font_color[2], 0),
        )


def set_run_typeface(pptx_run: Any, tag: str, font_name: Any) -> None:
    if not font_name:
        return

    rpr = pptx_run._r.get_or_add_rPr()
    child = rpr.find(qn(f"a:{tag}"))
    if child is None:
        child = OxmlElement(f"a:{tag}")
        rpr.append(child)
    child.set("typeface", str(font_name))


def find_shape_by_id(build_doc: Any, slide_index: int, shape_id: int) -> Any | None:
    for shape in build_doc.prs.slides[slide_index].shapes:
        if getattr(shape, "shape_id", None) == shape_id:
            return shape
    return None


def restore_shape_geometry(shape_obj: Any, payload: dict[str, Any]) -> None:
    width = to_int(payload.get("width"), 0)
    height = to_int(payload.get("height"), 0)
    if width <= 0 or height <= 0:
        return
    shape_obj.left = Emu(payload["left"])
    shape_obj.top = Emu(payload["top"])
    shape_obj.width = Emu(width)
    shape_obj.height = Emu(height)


def apply_rich_text(shape_obj: Any, shape: dict[str, Any], default_style: dict[str, Any]) -> None:
    paragraphs = text_paragraphs(shape)
    if not paragraphs or not getattr(shape_obj, "has_text_frame", False):
        return

    text_frame = shape_obj.text_frame
    text_frame.clear()

    paragraph_count = 0
    for paragraph_spec in paragraphs:
        para_text = str(paragraph_spec.get("text") or "").replace("\x0b", "\n")
        lines = para_text.split("\n")
        runs = paragraph_spec.get("runs", []) or []

        for line_text in lines:
            if paragraph_count == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            paragraph_count += 1

            if len(lines) == 1 and runs:
                for run_spec in runs:
                    pptx_run = paragraph.add_run()
                    pptx_run.text = str(run_spec.get("text") or "")
                    apply_run_style(pptx_run, run_spec, default_style)
                continue

            if not line_text:
                continue

            style_spec = runs[0] if runs else {}
            pptx_run = paragraph.add_run()
            pptx_run.text = line_text
            apply_run_style(pptx_run, style_spec, default_style)


def resolve_exported_media_map(layout: dict[str, Any]) -> dict[str, str]:
    result: dict[str, str] = {}
    for item in layout.get("exported_media", []) or []:
        zip_path = item.get("zip_path")
        export_path = item.get("export_path")
        if zip_path and export_path:
            result[str(zip_path)] = str(export_path)
    return result


def resolve_image_path(shape: dict[str, Any], layout: dict[str, Any], media_map: dict[str, str]) -> Optional[str]:
    for key in ("export_image_path", "image_path"):
        candidate = shape.get(key)
        if candidate:
            resolved = resolve_ml_home_path(str(candidate))
            if resolved:
                return str(resolved)

    image_zip_path = shape.get("image_zip_path")
    if image_zip_path and image_zip_path in media_map:
        resolved = resolve_ml_home_path(media_map[image_zip_path])
        if resolved:
            return str(resolved)

    exported_dir = layout.get("exported_media_dir")
    if exported_dir and image_zip_path:
        resolved_dir = resolve_ml_home_path(str(exported_dir))
        if resolved_dir:
            candidate = resolved_dir / Path(str(image_zip_path)).name
            if candidate.exists():
                return str(candidate)

    return None


def to_optional_float(value: Any) -> Optional[float]:
    if value is None:
        return None
    try:
        return float(value)
    except Exception:
        return None


def maybe_set_background(slide_index: int, slide: dict[str, Any], layout: dict[str, Any], build_doc: Any) -> None:
    background = slide.get("background") or {}
    background_api = slide.get("background_api") or {}

    bg = background_api or background or {}
    if (
        background_api.get("background_type") == "inherit"
        or background_api.get("mode") == "inherit"
        or background.get("mode") == "inherit"
    ):
        bg = layout.get("slide_master_background") or background or background_api or {}

    color_rgb = background.get("color_rgb")
    if color_rgb is None:
        color_rgb = bg.get("color_rgb")
    if isinstance(color_rgb, list) and len(color_rgb) == 3:
        set_slide_background_color(build_doc, slide_index, tuple(to_int(v) for v in color_rgb))
        return

    image_path = background.get("image_export_path") or background.get("image_path_hint")
    if not image_path:
        image_path = bg.get("image_export_path") or bg.get("image_path_hint")
    if image_path:
        resolved = resolve_ml_home_path(str(image_path))
        if resolved:
            set_slide_background_image(build_doc, slide_index, str(resolved))
            return

    image_ref = bg.get("image_ref")
    if image_ref:
        resolved = resolve_ml_home_path(str(image_ref))
        if resolved:
            set_slide_background_image(build_doc, slide_index, str(resolved))
            return


def rgb_triplet(value: Any) -> Optional[tuple[int, int, int]]:
    if isinstance(value, (list, tuple)) and len(value) == 3:
        return tuple(int(v) for v in value)
    return None


def infer_empty_shape_type(shape: dict[str, Any]) -> str:
    preset = str(shape.get("preset_geometry") or "")
    name = str(shape.get("name") or "")
    if preset == "ellipse" or name.startswith("Oval"):
        return "oval"
    return "rectangle"


def should_rebuild_empty_shape(shape: dict[str, Any]) -> bool:
    name = str(shape.get("name") or "")
    preset = str(shape.get("preset_geometry") or "")
    if rgb_triplet(shape.get("fill_color")) is not None:
        return True
    if shape.get("fill_none"):
        return True
    if name.startswith("Oval"):
        return True
    if name == "Rectangle 23":
        return True
    if preset == "line" or name == "Shape 4":
        return True
    return False


def add_shape_from_json(slide_index: int, shape: dict[str, Any], build_doc: Any) -> None:
    geometry = shape.get("geometry") or {}
    left = geometry.get("left_emu")
    top = geometry.get("top_emu")
    width = geometry.get("width_emu")
    height = geometry.get("height_emu")
    text = shape.get("text")

    if left is None or top is None or width is None or height is None or text is None:
        return

    if text == "" and not should_rebuild_empty_shape(shape):
        return

    payload: dict[str, Any] = {
        "slide_index": slide_index,
        "left": to_int(left),
        "top": to_int(top),
        "width": to_int(width),
        "height": to_int(height),
    }
    fill_color = rgb_triplet(shape.get("fill_color"))
    line_color = rgb_triplet(shape.get("line_color"))
    line_width = shape.get("line_width")
    fill_none = bool(shape.get("fill_none"))
    if line_width is not None:
        line_width = to_int(line_width)

    if text != "":
        style = infer_text_style(shape)
        text_frame = infer_text_frame_layout(shape)
        added = add_text(
            build_doc,
            slide_index=slide_index,
            text=str(text),
            left=payload["left"],
            top=payload["top"],
            width=payload["width"],
            height=payload["height"],
            font_size=style["font_size"],
            bold=style["bold"],
            italic=style["italic"],
            font_name=style["font_name"],
            fill_color=fill_color,
            line_color=line_color,
            line_width=line_width,
            word_wrap=text_frame.get("word_wrap"),
            auto_fit=text_frame.get("auto_fit"),
        )
        if has_run_level_text_style(shape):
            shape_obj = find_shape_by_id(build_doc, slide_index, int(added["shape_id"]))
            if shape_obj is not None:
                apply_rich_text(shape_obj, shape, style)
                if getattr(shape_obj, "has_text_frame", False):
                    _apply_text_frame_layout(
                        shape_obj.text_frame,
                        word_wrap=text_frame.get("word_wrap"),
                        auto_fit=text_frame.get("auto_fit"),
                    )
                restore_shape_geometry(shape_obj, payload)
    else:
        shape_height = payload["height"]
        if shape_height < 1:
            shape_height = line_width if line_width else 12700
        add_shape(
            build_doc,
            slide_index=slide_index,
            shape_type=infer_empty_shape_type(shape),
            left=payload["left"],
            top=payload["top"],
            width=payload["width"],
            height=shape_height,
            fill_color=fill_color,
            line_color=line_color,
            line_width=line_width,
            fill_none=fill_none,
        )


def build_ppt_from_json(layout: dict[str, Any], output_path: Path, dpi: int = DEFAULT_DPI) -> dict[str, Any]:
    slide_size = layout.get("slide_size") or {}
    width_in = float(slide_size.get("width_in") or 10.0)
    height_in = float(slide_size.get("height_in") or 7.5)
    slide_count = int(layout.get("slide_count") or len(layout.get("slides", [])))

    layout["build_file_path"] = str(output_path)
    new(
        str(output_path),
        plank_page_num=slide_count,
        plank_page_width=emu_to_px(width_in, dpi=dpi),
        plank_page_height=emu_to_px(height_in, dpi=dpi),
        dpi=dpi,
    )
    build_doc = open_presentation(str(output_path))

    exported_media_map = resolve_exported_media_map(layout)

    for slide in layout.get("slides", []):
        slide_index = json_slide_to_api_slide(slide["slide_index"])
        maybe_set_background(slide_index, slide, layout, build_doc)

        for shape in order_shapes_for_build(slide.get("shapes", [])):
            if (shape.get("shape_type") or "").lower() == "picture" or shape.get("xml_tag") == "pic":
                image_path = resolve_image_path(shape, layout, exported_media_map)
                if not image_path:
                    print(
                        f"[warn] missing exported image for slide {slide_index}, shape_id={shape.get('shape_id')}",
                        file=sys.stderr,
                    )
                    continue

                geometry = shape.get("geometry") or {}
                crop = shape.get("crop") or {}
                add_image(
                    build_doc,
                    slide_index=slide_index,
                    image_path=image_path,
                    left=to_int(geometry.get("left_emu")),
                    top=to_int(geometry.get("top_emu")),
                    width=to_int(geometry.get("width_emu")),
                    height=to_int(geometry.get("height_emu")),
                    keep_aspect_ratio=False,
                    rotation=to_optional_float(shape.get("rotation_deg")),
                    crop_left=to_optional_float(crop.get("left")),
                    crop_right=to_optional_float(crop.get("right")),
                    crop_top=to_optional_float(crop.get("top")),
                    crop_bottom=to_optional_float(crop.get("bottom")),
                )
                continue

            add_shape_from_json(slide_index, shape, build_doc)

    save(build_doc, str(output_path))
    return {
        "ok": True,
        "message": "success",
        "data": get_info(build_doc),
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Rebuild a PPTX from parsed JSON layout data.")
    parser.add_argument("json_path", help="Input layout JSON path")
    parser.add_argument(
        "-o",
        "--output",
        help="Output PPTX path; default is <json_stem>_rebuild.pptx beside the JSON",
        default=None,
    )
    parser.add_argument("--dpi", type=int, default=DEFAULT_DPI, help="Canvas DPI used for ppt/new")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    json_path = Path(args.json_path).resolve()
    if not json_path.exists():
        raise SystemExit(f"JSON file not found: {json_path}")

    layout = read_json(json_path)
    output_path = Path(args.output).resolve() if args.output else json_path.with_name(f"{json_path.stem}_rebuild.pptx")

    info = build_ppt_from_json(layout, output_path, dpi=args.dpi)
    print(json.dumps(info, ensure_ascii=False, indent=2))
    print(str(output_path))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
