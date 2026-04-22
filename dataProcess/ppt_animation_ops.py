# -*- coding: utf-8 -*-
"""
PPT 動畫核心操作模組（OOXML / timing XML）

目標
- 讀取指定頁所有動畫
- 讀取指定 shape 的動畫
- 新增 / 更新 / 刪除 / 清空 shape 動畫
- 重排投影片動畫順序
- 讀取 / 設定 / 清除投影片切換動畫

說明
- python-pptx 目前沒有穩定的高階動畫 API，因此本模組採用 OOXML 直接操作。
- 投影片切換動畫（transition）相對穩定；物件動畫（timing）採「可讀、可維護、可再編輯」為優先，
  以常見結構建立基本 entrance 效果。
- 本模組設計風格對齊 ppt_stdio.py：helper + 結構化 dict 回傳 + 對 PPTDocument 友善。
"""

from __future__ import annotations

from copy import deepcopy
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
except ImportError as exc:  # pragma: no cover
    raise ImportError("ppt_animation_ops 需要 python-pptx") from exc


_TRIGGER_TO_NODETYPE = {
    "on_click": "clickEffect",
    "with_previous": "withEffect",
    "after_previous": "afterEffect",
}

_NODETYPE_TO_TRIGGER = {v: k for k, v in _TRIGGER_TO_NODETYPE.items()}

_EFFECT_TO_PRESET = {
    "appear": {"preset_class": "entr", "preset_id": "1", "transition": "in", "filter": None},
    "fade": {"preset_class": "entr", "preset_id": "10", "transition": "in", "filter": "fade"},
    "wipe_left": {"preset_class": "entr", "preset_id": "2", "transition": "in", "filter": "wipe(l)"},
    "wipe_right": {"preset_class": "entr", "preset_id": "2", "transition": "in", "filter": "wipe(r)"},
    "wipe_up": {"preset_class": "entr", "preset_id": "2", "transition": "in", "filter": "wipe(u)"},
    "wipe_down": {"preset_class": "entr", "preset_id": "2", "transition": "in", "filter": "wipe(d)"},
}

_PRESET_TO_EFFECT = {
    (v["preset_class"], v["preset_id"], v["filter"]): k
    for k, v in _EFFECT_TO_PRESET.items()
}

_TRANSITION_TO_TAG = {
    "none": None,
    "fade": "p:fade",
    "push": "p:push",
    "wipe": "p:wipe",
    "cut": "p:cut",
}

_TAG_TO_TRANSITION = {
    "fade": "fade",
    "push": "push",
    "wipe": "wipe",
    "cut": "cut",
}


def _ns_pml() -> Dict[str, str]:
    return {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }


def _local_name(tag: str) -> str:
    return str(tag).split("}", 1)[-1]


def _validate_slide_index(document: Any, slide_index: int) -> None:
    prs = document.prs
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"slide_index 超出範圍: {slide_index}, 總頁數={len(prs.slides)}")


def _get_shape(document: Any, slide_index: int, shape_id: Optional[int] = None, shape_index: Optional[int] = None):
    if hasattr(document, "_get_shape"):
        return document._get_shape(slide_index=slide_index, shape_id=shape_id, shape_index=shape_index)

    _validate_slide_index(document, slide_index)
    slide = document.prs.slides[slide_index]
    if shape_id is None and shape_index is None:
        raise ValueError("shape_id 與 shape_index 至少需提供一個")
    if shape_id is not None:
        for idx, shape in enumerate(slide.shapes):
            if getattr(shape, "shape_id", None) == shape_id:
                return shape, idx
        raise ValueError(f"找不到 shape_id={shape_id}")
    assert shape_index is not None
    if shape_index < 0 or shape_index >= len(slide.shapes):
        raise IndexError(f"shape_index 超出範圍: {shape_index}, shape_count={len(slide.shapes)}")
    return slide.shapes[shape_index], shape_index


def _get_shape_cnvpr_id(shape: Any) -> Optional[int]:
    try:
        cnv = shape._element.xpath(".//*[local-name()='cNvPr'][1]")
        if cnv:
            raw = cnv[0].get("id")
            return int(raw) if raw is not None else None
    except Exception:
        return None
    return None


def _get_slide_timing_node(slide: Any) -> Any:
    try:
        return slide._element.find(qn("p:timing"))
    except Exception:
        return None


def _ensure_slide_timing_node(slide: Any) -> Any:
    timing = _get_slide_timing_node(slide)
    if timing is not None:
        return timing

    timing = OxmlElement("p:timing")
    sld = slide._element
    ext_lst = sld.find(qn("p:extLst"))
    if ext_lst is not None:
        idx = list(sld).index(ext_lst)
        sld.insert(idx, timing)
    else:
        sld.append(timing)
    return timing


def _find_first_child(parent: Any, child_qname: str) -> Any:
    return parent.find(qn(child_qname))


def _get_or_create_child(parent: Any, child_qname: str) -> Any:
    child = _find_first_child(parent, child_qname)
    if child is not None:
        return child
    child = OxmlElement(child_qname)
    parent.append(child)
    return child


def _get_or_create_main_sequence(slide: Any) -> Any:
    timing = _ensure_slide_timing_node(slide)
    tn_lst = _get_or_create_child(timing, "p:tnLst")

    root_par = _find_first_child(tn_lst, "p:par")
    if root_par is None:
        root_par = OxmlElement("p:par")
        tn_lst.append(root_par)

    root_ctn = _find_first_child(root_par, "p:cTn")
    if root_ctn is None:
        root_ctn = OxmlElement("p:cTn")
        root_ctn.set("id", "1")
        root_ctn.set("dur", "indefinite")
        root_ctn.set("restart", "never")
        root_ctn.set("nodeType", "tmRoot")
        root_par.append(root_ctn)

    root_child = _get_or_create_child(root_ctn, "p:childTnLst")

    seq = _find_first_child(root_child, "p:seq")
    if seq is None:
        seq = OxmlElement("p:seq")
        seq.set("concurrent", "1")
        seq.set("nextAc", "seek")
        root_child.append(seq)

    main_ctn = _find_first_child(seq, "p:cTn")
    if main_ctn is None:
        main_ctn = OxmlElement("p:cTn")
        main_ctn.set("id", "2")
        main_ctn.set("dur", "indefinite")
        main_ctn.set("nodeType", "mainSeq")
        seq.insert(0, main_ctn)

    _get_or_create_child(main_ctn, "p:childTnLst")

    prev_cond_lst = _get_or_create_child(seq, "p:prevCondLst")
    if _find_first_child(prev_cond_lst, "p:cond") is None:
        cond = OxmlElement("p:cond")
        cond.set("evt", "onPrev")
        cond.set("delay", "0")
        prev_cond_lst.append(cond)

    next_cond_lst = _get_or_create_child(seq, "p:nextCondLst")
    if _find_first_child(next_cond_lst, "p:cond") is None:
        cond = OxmlElement("p:cond")
        cond.set("evt", "onNext")
        cond.set("delay", "0")
        next_cond_lst.append(cond)

    return seq


def _get_main_sequence_child_list(slide: Any) -> Any:
    seq = _get_or_create_main_sequence(slide)
    main_ctn = _find_first_child(seq, "p:cTn")
    if main_ctn is None:
        raise RuntimeError("main sequence cTn 建立失敗")
    return _get_or_create_child(main_ctn, "p:childTnLst")


def _iter_ctn_ids(root: Any) -> List[int]:
    ids: List[int] = []
    try:
        for node in root.xpath(".//*[local-name()='cTn']"):
            raw = node.get("id")
            if raw is None:
                continue
            try:
                ids.append(int(raw))
            except Exception:
                continue
    except Exception:
        pass
    return ids


def _next_tn_id(slide: Any) -> int:
    ids = _iter_ctn_ids(slide._element)
    return (max(ids) + 1) if ids else 1


def _normalize_animation_effect_type(effect_type: str) -> str:
    normalized = str(effect_type or "fade").strip().lower()
    if normalized not in _EFFECT_TO_PRESET:
        raise ValueError(f"不支援的 effect_type: {effect_type}")
    return normalized


def _normalize_animation_trigger(trigger: str) -> str:
    normalized = str(trigger or "on_click").strip().lower()
    if normalized not in _TRIGGER_TO_NODETYPE:
        raise ValueError("trigger 必須是 on_click / with_previous / after_previous")
    return normalized


def _get_target_spid_from_effect(effect_node: Any) -> Optional[int]:
    try:
        sp_nodes = effect_node.xpath(".//*[local-name()='spTgt']")
        if sp_nodes:
            raw = sp_nodes[0].get("spid")
            return int(raw) if raw is not None else None
    except Exception:
        return None
    return None


def _serialize_animation_node(node: Any) -> Dict[str, Any]:
    notes: List[str] = []
    effect_node = None
    container_ctn = None

    if _local_name(node.tag) == "par":
        container_ctn = _find_first_child(node, "p:cTn")
        if container_ctn is not None:
            child_lst = _find_first_child(container_ctn, "p:childTnLst")
            if child_lst is not None:
                for child in list(child_lst):
                    if _local_name(child.tag) in {"animEffect", "set", "anim", "cmd"}:
                        effect_node = child
                        break
    elif _local_name(node.tag) in {"animEffect", "set", "anim", "cmd"}:
        effect_node = node

    if effect_node is None:
        return {
            "shape_cnvpr_id": None,
            "effect_type": "unknown",
            "trigger": "unknown",
            "duration_ms": None,
            "delay_ms": 0,
            "source": "timing_xml",
            "notes": ["無法辨識動畫節點結構。"],
        }

    effect_local = _local_name(effect_node.tag)
    effect_type = "unknown"
    trigger = "unknown"
    duration_ms = None
    delay_ms = 0
    preset_class = None
    preset_id = None
    spid = _get_target_spid_from_effect(effect_node)

    if container_ctn is not None:
        trigger = _NODETYPE_TO_TRIGGER.get(container_ctn.get("nodeType"), "unknown")
        st_cond_lst = _find_first_child(container_ctn, "p:stCondLst")
        if st_cond_lst is not None:
            cond = _find_first_child(st_cond_lst, "p:cond")
            if cond is not None:
                try:
                    delay_ms = int(cond.get("delay", "0"))
                except Exception:
                    delay_ms = 0

    ctn_candidates = effect_node.xpath(".//*[local-name()='cTn'][1]")
    if ctn_candidates:
        ctn = ctn_candidates[0]
        try:
            duration_ms = int(ctn.get("dur")) if ctn.get("dur") not in (None, "indefinite") else None
        except Exception:
            duration_ms = None
        preset_class = ctn.get("presetClass")
        preset_id = ctn.get("presetID")

    if effect_local == "animEffect":
        filter_val = effect_node.get("filter")
        effect_type = _PRESET_TO_EFFECT.get((preset_class, preset_id, filter_val), "fade" if filter_val == "fade" else "unknown")
        if effect_type == "unknown":
            notes.append(f"未完整對應的 animEffect: presetClass={preset_class}, presetID={preset_id}, filter={filter_val}")
    elif effect_local == "set":
        effect_type = "appear"
    else:
        effect_type = effect_local
        notes.append(f"目前僅完整支援 animEffect/set，讀到 {effect_local}。")

    return {
        "shape_cnvpr_id": spid,
        "effect_type": effect_type,
        "trigger": trigger,
        "duration_ms": duration_ms,
        "delay_ms": delay_ms,
        "source": "timing_xml",
        "notes": notes,
    }


def _list_animation_container_nodes(slide: Any) -> List[Any]:
    seq_child_lst = _get_main_sequence_child_list(slide)
    containers: List[Any] = []
    for child in list(seq_child_lst):
        if _local_name(child.tag) == "par":
            containers.append(child)
    return containers


def _find_animation_nodes_for_shape(slide: Any, shape_cnvpr_id: int) -> List[Any]:
    matched: List[Any] = []
    for node in _list_animation_container_nodes(slide):
        info = _serialize_animation_node(node)
        if info.get("shape_cnvpr_id") == shape_cnvpr_id:
            matched.append(node)
    return matched


def _build_target_element(shape_cnvpr_id: int) -> Any:
    tgt_el = OxmlElement("p:tgtEl")
    sp_tgt = OxmlElement("p:spTgt")
    sp_tgt.set("spid", str(int(shape_cnvpr_id)))
    tgt_el.append(sp_tgt)
    return tgt_el


def _build_animation_effect_node(
    shape_cnvpr_id: int,
    effect_type: str = "fade",
    trigger: str = "on_click",
    duration_ms: int = 500,
    delay_ms: int = 0,
    id_seed: int = 1,
) -> Any:
    effect_type = _normalize_animation_effect_type(effect_type)
    trigger = _normalize_animation_trigger(trigger)
    if duration_ms <= 0:
        raise ValueError("duration_ms 必須 > 0")
    if delay_ms < 0:
        raise ValueError("delay_ms 不可小於 0")

    outer_par = OxmlElement("p:par")
    outer_ctn = OxmlElement("p:cTn")
    outer_ctn.set("id", str(int(id_seed)))
    outer_ctn.set("fill", "hold")
    outer_ctn.set("nodeType", _TRIGGER_TO_NODETYPE[trigger])
    outer_ctn.set("dur", "indefinite")
    outer_par.append(outer_ctn)

    st_cond_lst = OxmlElement("p:stCondLst")
    cond = OxmlElement("p:cond")
    if trigger == "on_click":
        cond.set("evt", "onClick")
    else:
        cond.set("evt", "onBegin")
    cond.set("delay", str(int(delay_ms)))
    st_cond_lst.append(cond)
    outer_ctn.append(st_cond_lst)

    child_tn_lst = OxmlElement("p:childTnLst")
    outer_ctn.append(child_tn_lst)

    preset = _EFFECT_TO_PRESET[effect_type]
    inner_id = int(id_seed) + 1

    if effect_type == "appear":
        effect_node = OxmlElement("p:set")
        c_bhvr = OxmlElement("p:cBhvr")
        effect_node.append(c_bhvr)

        inner_ctn = OxmlElement("p:cTn")
        inner_ctn.set("id", str(inner_id))
        inner_ctn.set("dur", str(int(duration_ms)))
        inner_ctn.set("fill", "hold")
        inner_ctn.set("presetClass", preset["preset_class"])
        inner_ctn.set("presetID", preset["preset_id"])
        c_bhvr.append(inner_ctn)
        c_bhvr.append(_build_target_element(shape_cnvpr_id))

        attr_name_lst = OxmlElement("p:attrNameLst")
        attr_name = OxmlElement("p:attrName")
        attr_name.text = "style.visibility"
        attr_name_lst.append(attr_name)
        c_bhvr.append(attr_name_lst)

        to_node = OxmlElement("p:to")
        str_val = OxmlElement("p:strVal")
        str_val.set("val", "visible")
        to_node.append(str_val)
        effect_node.append(to_node)
    else:
        effect_node = OxmlElement("p:animEffect")
        effect_node.set("transition", str(preset["transition"]))
        if preset["filter"]:
            effect_node.set("filter", str(preset["filter"]))
        c_bhvr = OxmlElement("p:cBhvr")
        effect_node.append(c_bhvr)

        inner_ctn = OxmlElement("p:cTn")
        inner_ctn.set("id", str(inner_id))
        inner_ctn.set("dur", str(int(duration_ms)))
        inner_ctn.set("fill", "hold")
        inner_ctn.set("presetClass", preset["preset_class"])
        inner_ctn.set("presetID", preset["preset_id"])
        c_bhvr.append(inner_ctn)
        c_bhvr.append(_build_target_element(shape_cnvpr_id))

    child_tn_lst.append(effect_node)
    return outer_par


def _remove_animation_node(node: Any) -> None:
    parent = node.getparent()
    if parent is None:
        raise RuntimeError("動畫節點沒有 parent，無法刪除")
    parent.remove(node)


def _shape_map_for_slide(document: Any, slide_index: int) -> Dict[int, Dict[str, Any]]:
    slide = document.prs.slides[slide_index]
    mapping: Dict[int, Dict[str, Any]] = {}
    for idx, shape in enumerate(slide.shapes):
        cnvpr_id = _get_shape_cnvpr_id(shape)
        if cnvpr_id is None:
            continue
        mapping[cnvpr_id] = {
            "shape_id": getattr(shape, "shape_id", None),
            "shape_index": idx,
            "name": getattr(shape, "name", None),
            "shape_type": str(getattr(shape, "shape_type", "")),
        }
    return mapping


def get_slide_animations(document: Any, slide_index: int) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    slide = document.prs.slides[slide_index]
    shape_map = _shape_map_for_slide(document, slide_index)

    animations: List[Dict[str, Any]] = []
    for animation_index, node in enumerate(_list_animation_container_nodes(slide)):
        info = _serialize_animation_node(node)
        cnvpr_id = info.get("shape_cnvpr_id")
        resolved = shape_map.get(cnvpr_id or -1, {})
        animations.append(
            {
                "animation_index": animation_index,
                "shape_cnvpr_id": cnvpr_id,
                "shape_id": resolved.get("shape_id"),
                "shape_index": resolved.get("shape_index"),
                "shape_name": resolved.get("name"),
                "shape_type": resolved.get("shape_type"),
                "effect_type": info.get("effect_type"),
                "trigger": info.get("trigger"),
                "duration_ms": info.get("duration_ms"),
                "delay_ms": info.get("delay_ms"),
                "source": info.get("source"),
                "notes": info.get("notes", []),
            }
        )

    return {
        "slide_index": slide_index,
        "animation_count": len(animations),
        "animations": animations,
        "notes": [],
    }


def get_shape_animations(
    document: Any,
    slide_index: int,
    shape_id: Optional[int] = None,
    shape_index: Optional[int] = None,
) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    shape, resolved_shape_index = _get_shape(document, slide_index, shape_id=shape_id, shape_index=shape_index)
    cnvpr_id = _get_shape_cnvpr_id(shape)
    if cnvpr_id is None:
        raise RuntimeError("找不到目標 shape 的 cNvPr id，無法定位動畫")

    slide = document.prs.slides[slide_index]
    nodes = _find_animation_nodes_for_shape(slide, cnvpr_id)

    animations: List[Dict[str, Any]] = []
    for local_idx, node in enumerate(nodes):
        info = _serialize_animation_node(node)
        animations.append(
            {
                "animation_index": local_idx,
                "shape_cnvpr_id": cnvpr_id,
                "shape_id": getattr(shape, "shape_id", None),
                "shape_index": resolved_shape_index,
                "shape_name": getattr(shape, "name", None),
                "shape_type": str(getattr(shape, "shape_type", "")),
                "effect_type": info.get("effect_type"),
                "trigger": info.get("trigger"),
                "duration_ms": info.get("duration_ms"),
                "delay_ms": info.get("delay_ms"),
                "source": info.get("source"),
                "notes": info.get("notes", []),
            }
        )

    return {
        "slide_index": slide_index,
        "shape_cnvpr_id": cnvpr_id,
        "shape_id": getattr(shape, "shape_id", None),
        "shape_index": resolved_shape_index,
        "animation_count": len(animations),
        "animations": animations,
        "notes": [],
    }


def add_shape_animation(
    document: Any,
    slide_index: int,
    shape_id: Optional[int] = None,
    shape_index: Optional[int] = None,
    effect_type: str = "fade",
    trigger: str = "on_click",
    duration_ms: int = 500,
    delay_ms: int = 0,
) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    shape, resolved_shape_index = _get_shape(document, slide_index, shape_id=shape_id, shape_index=shape_index)
    cnvpr_id = _get_shape_cnvpr_id(shape)
    if cnvpr_id is None:
        raise RuntimeError("找不到目標 shape 的 cNvPr id，無法建立動畫")

    slide = document.prs.slides[slide_index]
    seq_child_lst = _get_main_sequence_child_list(slide)
    id_seed = _next_tn_id(slide)
    effect_node = _build_animation_effect_node(
        shape_cnvpr_id=cnvpr_id,
        effect_type=effect_type,
        trigger=trigger,
        duration_ms=duration_ms,
        delay_ms=delay_ms,
        id_seed=id_seed,
    )
    seq_child_lst.append(effect_node)

    info = _serialize_animation_node(effect_node)
    return {
        "slide_index": slide_index,
        "shape_cnvpr_id": cnvpr_id,
        "shape_id": getattr(shape, "shape_id", None),
        "shape_index": resolved_shape_index,
        "added": True,
        "effect_type": info.get("effect_type"),
        "trigger": info.get("trigger"),
        "duration_ms": info.get("duration_ms"),
        "delay_ms": info.get("delay_ms"),
        "notes": info.get("notes", []),
    }


def update_shape_animation(
    document: Any,
    slide_index: int,
    animation_index: int,
    shape_id: Optional[int] = None,
    shape_index: Optional[int] = None,
    effect_type: Optional[str] = None,
    trigger: Optional[str] = None,
    duration_ms: Optional[int] = None,
    delay_ms: Optional[int] = None,
) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    shape, resolved_shape_index = _get_shape(document, slide_index, shape_id=shape_id, shape_index=shape_index)
    cnvpr_id = _get_shape_cnvpr_id(shape)
    if cnvpr_id is None:
        raise RuntimeError("找不到目標 shape 的 cNvPr id，無法更新動畫")

    slide = document.prs.slides[slide_index]
    nodes = _find_animation_nodes_for_shape(slide, cnvpr_id)
    if animation_index < 0 or animation_index >= len(nodes):
        raise IndexError(f"animation_index 超出範圍: {animation_index}, animation_count={len(nodes)}")

    target = nodes[animation_index]
    current = _serialize_animation_node(target)
    new_effect_type = effect_type or current.get("effect_type") or "fade"
    new_trigger = trigger or current.get("trigger") or "on_click"
    new_duration = int(duration_ms if duration_ms is not None else (current.get("duration_ms") or 500))
    new_delay = int(delay_ms if delay_ms is not None else (current.get("delay_ms") or 0))

    rebuilt = _build_animation_effect_node(
        shape_cnvpr_id=cnvpr_id,
        effect_type=new_effect_type,
        trigger=new_trigger,
        duration_ms=new_duration,
        delay_ms=new_delay,
        id_seed=_next_tn_id(slide),
    )

    parent = target.getparent()
    if parent is None:
        raise RuntimeError("動畫節點沒有 parent，無法更新")
    insert_idx = list(parent).index(target)
    parent.remove(target)
    parent.insert(insert_idx, rebuilt)

    info = _serialize_animation_node(rebuilt)
    return {
        "slide_index": slide_index,
        "shape_cnvpr_id": cnvpr_id,
        "shape_id": getattr(shape, "shape_id", None),
        "shape_index": resolved_shape_index,
        "updated": True,
        "animation_index": animation_index,
        "effect_type": info.get("effect_type"),
        "trigger": info.get("trigger"),
        "duration_ms": info.get("duration_ms"),
        "delay_ms": info.get("delay_ms"),
        "notes": info.get("notes", []),
    }


def delete_shape_animation(
    document: Any,
    slide_index: int,
    animation_index: int,
    shape_id: Optional[int] = None,
    shape_index: Optional[int] = None,
) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    shape, resolved_shape_index = _get_shape(document, slide_index, shape_id=shape_id, shape_index=shape_index)
    cnvpr_id = _get_shape_cnvpr_id(shape)
    if cnvpr_id is None:
        raise RuntimeError("找不到目標 shape 的 cNvPr id，無法刪除動畫")

    slide = document.prs.slides[slide_index]
    nodes = _find_animation_nodes_for_shape(slide, cnvpr_id)
    if animation_index < 0 or animation_index >= len(nodes):
        raise IndexError(f"animation_index 超出範圍: {animation_index}, animation_count={len(nodes)}")

    target = nodes[animation_index]
    info = _serialize_animation_node(target)
    _remove_animation_node(target)

    return {
        "slide_index": slide_index,
        "shape_cnvpr_id": cnvpr_id,
        "shape_id": getattr(shape, "shape_id", None),
        "shape_index": resolved_shape_index,
        "deleted": True,
        "animation_index": animation_index,
        "effect_type": info.get("effect_type"),
        "trigger": info.get("trigger"),
        "notes": info.get("notes", []),
    }


def clear_shape_animations(
    document: Any,
    slide_index: int,
    shape_id: Optional[int] = None,
    shape_index: Optional[int] = None,
) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    shape, resolved_shape_index = _get_shape(document, slide_index, shape_id=shape_id, shape_index=shape_index)
    cnvpr_id = _get_shape_cnvpr_id(shape)
    if cnvpr_id is None:
        raise RuntimeError("找不到目標 shape 的 cNvPr id，無法清空動畫")

    slide = document.prs.slides[slide_index]
    nodes = _find_animation_nodes_for_shape(slide, cnvpr_id)
    count = len(nodes)
    for node in list(nodes):
        _remove_animation_node(node)

    return {
        "slide_index": slide_index,
        "shape_cnvpr_id": cnvpr_id,
        "shape_id": getattr(shape, "shape_id", None),
        "shape_index": resolved_shape_index,
        "cleared": True,
        "removed_count": count,
        "notes": [],
    }


def clear_slide_animations(document: Any, slide_index: int) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    slide = document.prs.slides[slide_index]
    nodes = _list_animation_container_nodes(slide)
    count = len(nodes)
    for node in list(nodes):
        _remove_animation_node(node)
    return {
        "slide_index": slide_index,
        "cleared": True,
        "removed_count": count,
        "notes": [],
    }


def reorder_slide_animations(document: Any, slide_index: int, new_order: List[int]) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    slide = document.prs.slides[slide_index]
    seq_child_lst = _get_main_sequence_child_list(slide)
    current_nodes = _list_animation_container_nodes(slide)
    count = len(current_nodes)
    if len(new_order) != count:
        raise ValueError(f"new_order 長度必須等於目前動畫數量 {count}")
    if sorted(new_order) != list(range(count)):
        raise ValueError("new_order 必須是 0 到 animation_count-1 的完整排列")

    others = [child for child in list(seq_child_lst) if _local_name(child.tag) != "par"]
    for child in list(seq_child_lst):
        seq_child_lst.remove(child)
    for idx in new_order:
        seq_child_lst.append(current_nodes[idx])
    for child in others:
        seq_child_lst.append(child)

    return {
        "slide_index": slide_index,
        "animation_count": count,
        "new_order": new_order,
        "notes": [],
    }


def get_slide_transition(document: Any, slide_index: int) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    slide = document.prs.slides[slide_index]
    transition = slide._element.find(qn("p:transition"))
    if transition is None:
        return {
            "slide_index": slide_index,
            "has_transition": False,
            "transition_type": "none",
            "duration_ms": None,
            "advance_on_click": True,
            "advance_after_ms": None,
            "notes": [],
        }

    transition_type = "none"
    for child in list(transition):
        local = _local_name(child.tag)
        if local in _TAG_TO_TRANSITION:
            transition_type = _TAG_TO_TRANSITION[local]
            break

    duration_ms = None
    try:
        if transition.get("spd"):
            spd = transition.get("spd")
            duration_ms = {"slow": 2000, "med": 1000, "fast": 500}.get(spd)
    except Exception:
        duration_ms = None

    advance_on_click = transition.get("advClick", "1") != "0"
    advance_after_ms = None
    if transition.get("advTm") is not None:
        try:
            advance_after_ms = int(transition.get("advTm"))
        except Exception:
            advance_after_ms = None

    return {
        "slide_index": slide_index,
        "has_transition": True,
        "transition_type": transition_type,
        "duration_ms": duration_ms,
        "advance_on_click": advance_on_click,
        "advance_after_ms": advance_after_ms,
        "notes": [],
    }


def _duration_to_spd(duration_ms: Optional[int]) -> Optional[str]:
    if duration_ms is None:
        return None
    if duration_ms <= 650:
        return "fast"
    if duration_ms >= 1500:
        return "slow"
    return "med"


def set_slide_transition(
    document: Any,
    slide_index: int,
    transition_type: str = "fade",
    duration_ms: Optional[int] = None,
    advance_on_click: bool = True,
    advance_after_ms: Optional[int] = None,
) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    normalized = str(transition_type or "none").strip().lower()
    if normalized not in _TRANSITION_TO_TAG:
        raise ValueError(f"不支援的 transition_type: {transition_type}")
    if advance_after_ms is not None and advance_after_ms < 0:
        raise ValueError("advance_after_ms 不可小於 0")

    slide = document.prs.slides[slide_index]
    existing = slide._element.find(qn("p:transition"))
    if existing is not None:
        slide._element.remove(existing)

    if normalized == "none":
        return {
            "slide_index": slide_index,
            "transition_type": "none",
            "duration_ms": None,
            "advance_on_click": True,
            "advance_after_ms": None,
            "notes": [],
        }

    transition = OxmlElement("p:transition")
    transition.set("advClick", "1" if advance_on_click else "0")
    if advance_after_ms is not None:
        transition.set("advTm", str(int(advance_after_ms)))
    spd = _duration_to_spd(duration_ms)
    if spd is not None:
        transition.set("spd", spd)

    transition_child_tag = _TRANSITION_TO_TAG[normalized]
    if transition_child_tag:
        transition.append(OxmlElement(transition_child_tag))

    timing = slide._element.find(qn("p:timing"))
    ext_lst = slide._element.find(qn("p:extLst"))
    if timing is not None:
        idx = list(slide._element).index(timing)
        slide._element.insert(idx, transition)
    elif ext_lst is not None:
        idx = list(slide._element).index(ext_lst)
        slide._element.insert(idx, transition)
    else:
        slide._element.append(transition)

    return {
        "slide_index": slide_index,
        "transition_type": normalized,
        "duration_ms": duration_ms,
        "advance_on_click": advance_on_click,
        "advance_after_ms": advance_after_ms,
        "notes": [],
    }


def clear_slide_transition(document: Any, slide_index: int) -> Dict[str, Any]:
    _validate_slide_index(document, slide_index)
    slide = document.prs.slides[slide_index]
    existing = slide._element.find(qn("p:transition"))
    removed = False
    if existing is not None:
        slide._element.remove(existing)
        removed = True
    return {
        "slide_index": slide_index,
        "cleared": True,
        "removed": removed,
        "notes": [],
    }
