# -*- coding: utf-8 -*-
"""
PPTX 核心操作模組
- 建立新簡報
- 開啟既有簡報
- 儲存簡報
- 新增空白頁
- 新增文字框
- 新增圖片
- 新增表格
- 讀取簡報資訊
- 列出投影片資訊

依賴:
    pip install python-pptx pillow
"""

import os
import json
from typing import Any, Dict, List, Optional, Tuple

# 實作「佈景主題 / 投影片背景」讀取（解析 pptx 內 XML）時請啟用標準庫：
# import zipfile
# from xml.etree import ElementTree as ET

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from copy import deepcopy
    from pptx.enum.shapes import PP_PLACEHOLDER
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

try:
    import math
    import shutil
    import subprocess
    import tempfile
    from pathlib import Path

    import fitz  # PyMuPDF
    import matplotlib.pyplot as plt
    import matplotlib.image as mpimg
    from matplotlib import font_manager
    PYTHON_PPTX_EXTRA_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_EXTRA_AVAILABLE = False

try:
    from PIL import Image
    from PIL import ImageFont
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False


EMU_PER_INCH = 914400
DEFAULT_DPI = 96
DEFAULT_FONT_ZH = "微軟正黑體"
DEFAULT_FONT_EN = "Consolas"

import dataProcess as dp
PYTHON_PPTX_AVAILABLE = dp.ptp.PYTHON_PPTX_AVAILABLE
Presentation = dp.ptp.Presentation
m_logger = dp.ptp.m_logger
LOGger = dp.ptp.LOGger

_CONFIG_CACHE: Optional[Dict[str, Any]] = None




def _ensure_pptx_available():
    if not PYTHON_PPTX_AVAILABLE:
        raise ImportError("python-pptx 未安裝，請先執行: pip install python-pptx")


def _px_to_emu(px: int, dpi: int = DEFAULT_DPI) -> int:
    return int((px / float(dpi)) * EMU_PER_INCH)


def _normalize_file_path(file_path: str) -> str:
    if not file_path:
        raise ValueError("file_path 不可為空")
    file_path = os.path.abspath(file_path)
    if not file_path.lower().endswith(".pptx"):
        file_path += ".pptx"
    return file_path


def _ensure_parent_dir(file_path: str) -> None:
    parent_dir = os.path.dirname(file_path)
    if parent_dir:
        os.makedirs(parent_dir, exist_ok=True)


def _get_blank_layout(prs: Presentation):
    # python-pptx 內建常見 blank layout 通常是 index 6
    # 若模板不同，退回最後一個 layout
    try:
        return prs.slide_layouts[6]
    except Exception:
        return prs.slide_layouts[len(prs.slide_layouts) - 1]


def _remove_all_slides(prs: Presentation) -> None:
    """
    清空簡報中所有投影片
    """
    slide_id_list = prs.slides._sldIdLst
    while len(slide_id_list) > 0:
        r_id = slide_id_list[0].rId
        prs.part.drop_rel(r_id)
        del slide_id_list[0]


def _validate_slide_index(prs: Presentation, slide_index: int) -> None:
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"slide_index 超出範圍: {slide_index}, 總頁數={len(prs.slides)}")


def _rgb_tuple_to_color(rgb: Optional[Tuple[int, int, int]]) -> Optional[RGBColor]:
    if rgb is None:
        return None
    if len(rgb) != 3:
        raise ValueError("rgb 顏色格式必須是 (R, G, B)")
    r, g, b = rgb
    return RGBColor(int(r), int(g), int(b))


def _load_pptstudio_config() -> Dict[str, Any]:
    global _CONFIG_CACHE
    if _CONFIG_CACHE is not None:
        return _CONFIG_CACHE

    config_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "config.json")
    if not os.path.exists(config_path):
        _CONFIG_CACHE = {}
        return _CONFIG_CACHE

    try:
        with open(config_path, "r", encoding="utf-8") as f:
            _CONFIG_CACHE = json.load(f)
    except Exception:
        _CONFIG_CACHE = {}
    return _CONFIG_CACHE


def _contains_cjk(text: str) -> bool:
    for char in text:
        code = ord(char)
        if 0x4E00 <= code <= 0x9FFF:
            return True
        if 0x3400 <= code <= 0x4DBF:
            return True
        if 0x3040 <= code <= 0x30FF:
            return True
        if 0xAC00 <= code <= 0xD7AF:
            return True
    return False


def _get_default_fonts() -> Tuple[str, str]:
    config = _load_pptstudio_config()
    fonts_config = config.get("fonts", {}) if isinstance(config, dict) else {}
    zh_font = str(fonts_config.get("default_zh", DEFAULT_FONT_ZH)).strip() or DEFAULT_FONT_ZH
    en_font = str(fonts_config.get("default_en", DEFAULT_FONT_EN)).strip() or DEFAULT_FONT_EN
    return zh_font, en_font


def _looks_like_font_path(value: str) -> bool:
    lowered = value.lower()
    if lowered.endswith((".ttf", ".otf", ".ttc", ".otc")):
        return True
    if os.path.sep in value:
        return True
    if os.path.altsep and os.path.altsep in value:
        return True
    return False


def _resolve_font_name(font_name: Optional[str]) -> Optional[str]:
    if font_name is None:
        return None

    candidate = str(font_name).strip()
    if not candidate:
        return None

    if not _looks_like_font_path(candidate):
        return candidate

    if not os.path.exists(candidate):
        raise FileNotFoundError(f"找不到字型檔: {candidate}")

    font_path = os.path.abspath(candidate)

    if PIL_AVAILABLE:
        try:
            font = ImageFont.truetype(font_path, size=12)
            family_name = font.getname()[0]
            if family_name:
                return family_name
        except Exception:
            pass

    return os.path.splitext(os.path.basename(font_path))[0]

def _activate_matplotlib_chinese_font() -> None:
    """
    設定 matplotlib 可用的中文字型，避免截圖標題顯示為方塊。
    """
    if not PYTHON_PPTX_EXTRA_AVAILABLE:
        return

    preferred_fonts = [
        "Microsoft JhengHei",
        "Microsoft YaHei",
        "PingFang TC",
        "Noto Sans CJK TC",
        "Noto Sans CJK SC",
        "SimHei",
        "DejaVu Sans",
    ]
    try:
        available_fonts = {f.name for f in font_manager.fontManager.ttflist}
        active_fonts = [font for font in preferred_fonts if font in available_fonts]
        if not active_fonts:
            active_fonts = ["DejaVu Sans"]
    except Exception:
        active_fonts = ["Microsoft JhengHei", "DejaVu Sans"]

    plt.rcParams["font.sans-serif"] = active_fonts
    plt.rcParams["axes.unicode_minus"] = False


def _rgb_from_color(color_obj: Any) -> Optional[List[int]]:
    if color_obj is None:
        return None
    try:
        rgb = color_obj.rgb
    except Exception:
        return None
    if rgb is None:
        return None

    hex_value = str(rgb).strip().lstrip("#")
    if len(hex_value) != 6:
        return None
    try:
        return [int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16)]
    except Exception:
        return None


def _clamp_transparency(value: float) -> float:
    if value < 0.0:
        return 0.0
    if value > 1.0:
        return 1.0
    return value


def _find_fill_color_node(shape: Any) -> Optional[Any]:
    try:
        sp_pr = shape._element.spPr
        if sp_pr is None:
            return None
        solid_fill = sp_pr.find(qn("a:solidFill"))
        if solid_fill is None:
            return None
        for tag in ("a:srgbClr", "a:schemeClr", "a:sysClr", "a:prstClr"):
            color_node = solid_fill.find(qn(tag))
            if color_node is not None:
                return color_node
        return None
    except Exception:
        return None


def _read_fill_transparency_from_xml(shape: Any) -> Optional[float]:
    color_node = _find_fill_color_node(shape)
    if color_node is None:
        return None
    alpha = color_node.find(qn("a:alpha"))
    if alpha is None:
        return None
    alpha_val = alpha.get("val")
    if alpha_val is None:
        return None
    try:
        alpha_100000 = float(alpha_val)
    except Exception:
        return None
    opacity = max(0.0, min(100000.0, alpha_100000)) / 100000.0
    return round(1.0 - opacity, 6)


def _write_fill_transparency_to_xml(shape: Any, fill_transparency: float) -> None:
    transparency = _clamp_transparency(float(fill_transparency))
    alpha_100000 = int(round((1.0 - transparency) * 100000.0))

    sp_pr = shape._element.spPr
    if sp_pr is None:
        raise RuntimeError("shape 缺少 spPr，無法設定 fill 透明度")

    solid_fill = sp_pr.find(qn("a:solidFill"))
    if solid_fill is None:
        solid_fill = OxmlElement("a:solidFill")
        sp_pr.append(solid_fill)

    color_node = None
    for tag in ("a:srgbClr", "a:schemeClr", "a:sysClr", "a:prstClr"):
        color_node = solid_fill.find(qn(tag))
        if color_node is not None:
            break
    if color_node is None:
        color_node = OxmlElement("a:srgbClr")
        color_node.set("val", "FFFFFF")
        solid_fill.append(color_node)

    for child in list(color_node):
        if child.tag == qn("a:alpha"):
            color_node.remove(child)
    alpha = OxmlElement("a:alpha")
    alpha.set("val", str(alpha_100000))
    color_node.append(alpha)


def _dash_style_to_name(dash_style: Any) -> Optional[str]:
    if dash_style is None:
        return "solid"
    name = getattr(dash_style, "name", None)
    if name:
        return str(name).lower()
    return str(dash_style).lower()


def _name_to_dash_style(line_style: Optional[str]) -> Any:
    if line_style is None:
        return None

    normalized = str(line_style).strip().lower()
    if not normalized:
        return None
    if normalized == "solid":
        return None

    candidates: Dict[str, List[str]] = {
        "dash": ["DASH", "SYS_DASH", "LG_DASH"],
        "dot": ["ROUND_DOT", "SQUARE_DOT", "SYS_DOT"],
        "dash_dot": ["DASH_DOT", "SYS_DASH_DOT", "LG_DASH_DOT"],
        "dash_dot_dot": ["DASH_DOT_DOT", "SYS_DASH_DOT_DOT", "LG_DASH_DOT_DOT"],
    }
    if normalized not in candidates:
        raise ValueError("line_style 必須是 solid/dash/dot/dash_dot/dash_dot_dot")

    for enum_name in candidates[normalized]:
        enum_val = getattr(MSO_LINE_DASH_STYLE, enum_name, None)
        if enum_val is not None:
            return enum_val
    raise ValueError(f"目前環境不支援 line_style={normalized}")


class PPTDocument:
    """
    封裝一個 Presentation 物件，提供可操作方法
    """

    def __init__(self, prs: Presentation, file_path: Optional[str] = None):
        self.prs = prs
        self.file_path = file_path

    @property
    def slide_count(self) -> int:
        return len(self.prs.slides)

    def save(self, file_path: Optional[str] = None) -> str:
        save_path = file_path or self.file_path
        if not save_path:
            raise ValueError("沒有可儲存的 file_path，請明確提供")
        save_path = _normalize_file_path(save_path)
        _ensure_parent_dir(save_path)
        self.prs.save(save_path)
        self.file_path = save_path
        return save_path

    def add_blank_slide(self) -> int:
        blank_layout = _get_blank_layout(self.prs)
        self.prs.slides.add_slide(blank_layout)
        return len(self.prs.slides) - 1

    def add_blank_slides(self, page_num: int = 1) -> List[int]:
        if page_num < 1:
            raise ValueError("page_num 至少要 >= 1")
        indices = []
        for _ in range(page_num):
            indices.append(self.add_blank_slide())
        return indices

    def add_textbox(
            self,
            slide_index: int,
            text: str,
            left: int,
            top: int,
            width: int,
            height: int,
            font_size: int = 20,
            bold: bool = False,
            italic: bool = False,
            font_name: Optional[str] = None,
            font_color: Optional[Tuple[int, int, int]] = None,
            align: str = "left",
        ) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)
        slide = self.prs.slides[slide_index]
        resolved_font_name = _resolve_font_name(font_name)
        default_zh_font, default_en_font = _get_default_fonts()
        resolved_default_zh = _resolve_font_name(default_zh_font)
        resolved_default_en = _resolve_font_name(default_en_font)

        shape = slide.shapes.add_textbox(
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height),
        )

        tf = shape.text_frame
        tf.clear()

        p = tf.paragraphs[0]
        p.text = text or ""

        align_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }
        p.alignment = align_map.get((align or "left").lower(), PP_ALIGN.LEFT)

        for run in p.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            if resolved_font_name:
                run.font.name = resolved_font_name
            else:
                default_font = resolved_default_zh if _contains_cjk(run.text or "") else resolved_default_en
                if default_font:
                    run.font.name = default_font
            color = _rgb_tuple_to_color(font_color)
            if color:
                run.font.color.rgb = color

        return {
            "slide_index": slide_index,
            "shape_id": shape.shape_id,
            "text": text,
            "left": left,
            "top": top,
            "width": width,
            "height": height,
        }

    def add_image(
            self,
            slide_index: int,
            image_path: str,
            left: int,
            top: int,
            width: Optional[int] = None,
            height: Optional[int] = None,
            keep_aspect_ratio: bool = True,
        ) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)

        if not image_path:
            raise ValueError("image_path 不可為空")
        image_path = os.path.abspath(image_path)
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"找不到圖片: {image_path}")

        slide = self.prs.slides[slide_index]

        img_width = width
        img_height = height

        if keep_aspect_ratio and PIL_AVAILABLE and (width is None or height is None):
            with Image.open(image_path) as img:
                px_w, px_h = img.size
            ratio = px_w / px_h if px_h else 1.0

            if width is not None and height is None:
                img_width = width
                img_height = int(width / ratio)
            elif height is not None and width is None:
                img_height = height
                img_width = int(height * ratio)

        kwargs = {
            "image_file": image_path,
            "left": Emu(left),
            "top": Emu(top),
        }
        if img_width is not None:
            kwargs["width"] = Emu(img_width)
        if img_height is not None:
            kwargs["height"] = Emu(img_height)

        picture = slide.shapes.add_picture(**kwargs)

        return {
            "slide_index": slide_index,
            "shape_id": picture.shape_id,
            "image_path": image_path,
            "left": left,
            "top": top,
            "width": picture.width,
            "height": picture.height,
        }

    def add_table(
            self,
            slide_index: int,
            rows: int,
            cols: int,
            left: int,
            top: int,
            width: int,
            height: int,
            data: Optional[List[List[Any]]] = None,
            first_row_as_header: bool = False,
            font_size: int = 14,
        ) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)

        if rows < 1 or cols < 1:
            raise ValueError("rows 與 cols 都必須 >= 1")

        slide = self.prs.slides[slide_index]
        graphic_frame = slide.shapes.add_table(
            rows=rows,
            cols=cols,
            left=Emu(left),
            top=Emu(top),
            width=Emu(width),
            height=Emu(height),
        )
        table = graphic_frame.table

        if data:
            for r in range(min(rows, len(data))):
                row_data = data[r]
                for c in range(min(cols, len(row_data))):
                    cell = table.cell(r, c)
                    cell.text = "" if row_data[c] is None else str(row_data[c])

                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)
                            if first_row_as_header and r == 0:
                                run.font.bold = True

        return {
            "slide_index": slide_index,
            "shape_id": graphic_frame.shape_id,
            "rows": rows,
            "cols": cols,
            "left": left,
            "top": top,
            "width": width,
            "height": height,
        }

    def add_shape(
            self,
            slide_index: int,
            shape_type: str,
            left: int,
            top: int,
            width: int,
            height: int,
            text: str = "",
            fill_color: Optional[Tuple[int, int, int]] = None,
            line_color: Optional[Tuple[int, int, int]] = None,
            line_width: Optional[int] = None,
            font_size: int = 18,
            bold: bool = False,
            font_name: Optional[str] = None,
            font_color: Optional[Tuple[int, int, int]] = None,
        ) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)
        slide = self.prs.slides[slide_index]

        shape_map = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "diamond": MSO_SHAPE.DIAMOND,
            "hexagon": MSO_SHAPE.HEXAGON,
            "parallelogram": MSO_SHAPE.PARALLELOGRAM,
            "trapezoid": MSO_SHAPE.TRAPEZOID,
            "chevron": MSO_SHAPE.CHEVRON,
            "right_arrow": MSO_SHAPE.RIGHT_ARROW,
            "left_arrow": MSO_SHAPE.LEFT_ARROW,
            "up_arrow": MSO_SHAPE.UP_ARROW,
            "down_arrow": MSO_SHAPE.DOWN_ARROW,
            "cloud": MSO_SHAPE.CLOUD,
        }

        if shape_type not in shape_map:
            raise ValueError(f"不支援的 shape_type: {shape_type}")

        shape = slide.shapes.add_shape(
            shape_map[shape_type],
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height),
        )

        # fill
        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _rgb_tuple_to_color(fill_color)

        # line
        if line_color is not None:
            shape.line.color.rgb = _rgb_tuple_to_color(line_color)

        if line_width is not None:
            shape.line.width = Emu(line_width)

        # text
        if text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            p.text = text

            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if font_name:
                    run.font.name = font_name
                if font_color is not None:
                    run.font.color.rgb = _rgb_tuple_to_color(font_color)

        return {
            "slide_index": slide_index,
            "shape_id": shape.shape_id,
            "shape_type": shape_type,
            "text": text,
            "left": left,
            "top": top,
            "width": width,
            "height": height,
        }
    
    def add_line(
            self,
            slide_index: int,
            x1: int,
            y1: int,
            x2: int,
            y2: int,
            line_color: Optional[Tuple[int, int, int]] = None,
            line_width: Optional[int] = None,
        ) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)
        slide = self.prs.slides[slide_index]

        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Emu(x1),
            Emu(y1),
            Emu(x2),
            Emu(y2),
        )

        if line_color is not None:
            shape.line.color.rgb = _rgb_tuple_to_color(line_color)

        if line_width is not None:
            shape.line.width = Emu(line_width)

        return {
            "slide_index": slide_index,
            "shape_id": shape.shape_id,
            "type": "line",
            "x1": x1,
            "y1": y1,
            "x2": x2,
            "y2": y2,
        }
    
    def add_arrow(
            self,
            slide_index: int,
            left: int,
            top: int,
            width: int,
            height: int,
            direction: str = "right",
            text: str = "",
            fill_color: Optional[Tuple[int, int, int]] = None,
            line_color: Optional[Tuple[int, int, int]] = None,
            line_width: Optional[int] = None,
            font_size: int = 18,
            bold: bool = False,
            font_name: Optional[str] = None,
            font_color: Optional[Tuple[int, int, int]] = None,
        ) -> Dict[str, Any]:
        direction_map = {
            "right": "right_arrow",
            "left": "left_arrow",
            "up": "up_arrow",
            "down": "down_arrow",
        }

        if direction not in direction_map:
            raise ValueError("direction 必須是 right / left / up / down")

        return self.add_shape(
            slide_index=slide_index,
            shape_type=direction_map[direction],
            left=left,
            top=top,
            width=width,
            height=height,
            text=text,
            fill_color=fill_color,
            line_color=line_color,
            line_width=line_width,
            font_size=font_size,
            bold=bold,
            font_name=font_name,
            font_color=font_color,
        )
    
    def set_slide_background_color(
            self,
            slide_index: int,
            rgb: Tuple[int, int, int],
        ) -> Dict[str, Any]:
        """
        設定指定頁的純色背景
        """
        _validate_slide_index(self.prs, slide_index)

        if len(rgb) != 3:
            raise ValueError("rgb 必須是 (R, G, B)")

        slide = self.prs.slides[slide_index]
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))

        return {
            "slide_index": slide_index,
            "background_type": "color",
            "rgb": [int(rgb[0]), int(rgb[1]), int(rgb[2])],
        }

    def set_slide_background_image(
            self,
            slide_index: int,
            image_path: str,
        ) -> Dict[str, Any]:
        """
        用滿版圖片模擬投影片背景
        """
        _validate_slide_index(self.prs, slide_index)

        if not image_path:
            raise ValueError("image_path 不可為空")

        image_path = os.path.abspath(image_path)
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"找不到圖片: {image_path}")

        slide = self.prs.slides[slide_index]

        picture = slide.shapes.add_picture(
            image_path,
            0,
            0,
            width=self.prs.slide_width,
            height=self.prs.slide_height,
        )

        return {
            "slide_index": slide_index,
            "background_type": "image",
            "image_path": image_path,
            "shape_id": picture.shape_id,
        }

    def set_slides_background_color(
            self,
            slide_indices: List[int],
            rgb: Tuple[int, int, int],
        ) -> List[Dict[str, Any]]:
        results = []
        for slide_index in slide_indices:
            results.append(self.set_slide_background_color(slide_index, rgb))
        return results

    def set_slides_background_image(
            self,
            slide_indices: List[int],
            image_path: str,
        ) -> List[Dict[str, Any]]:
        results = []
        for slide_index in slide_indices:
            results.append(self.set_slide_background_image(slide_index, image_path))
        return results

    def get_info(self) -> Dict[str, Any]:
        return {
            "file_path": self.file_path,
            "slide_count": len(self.prs.slides),
            "slide_width": int(self.prs.slide_width),
            "slide_height": int(self.prs.slide_height),
            "layout_count": len(self.prs.slide_layouts),
        }

    def list_slides(self) -> List[Dict[str, Any]]:
        results = []
        for idx, slide in enumerate(self.prs.slides):
            shape_infos = []
            title_text = None

            for shape in slide.shapes:
                shape_info = {
                    "shape_id": getattr(shape, "shape_id", None),
                    "name": getattr(shape, "name", None),
                    "shape_type": str(getattr(shape, "shape_type", "")),
                    "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
                }
                if getattr(shape, "has_text_frame", False):
                    try:
                        text = shape.text_frame.text.strip()
                    except Exception:
                        text = ""
                    shape_info["text_preview"] = text[:120]
                    if title_text is None and text:
                        title_text = text[:120]
                shape_infos.append(shape_info)

            results.append({
                "slide_index": idx,
                "shape_count": len(slide.shapes),
                "title_preview": title_text,
                "shapes": shape_infos,
            })
        return results

    def _get_textbox_shape(self, slide_index: int, shape_id: Optional[int] = None, shape_index: Optional[int] = None):
        _validate_slide_index(self.prs, slide_index)
        slide = self.prs.slides[slide_index]

        if shape_id is None and shape_index is None:
            raise ValueError("shape_id 與 shape_index 至少需提供一個")

        if shape_id is not None:
            for idx, shape in enumerate(slide.shapes):
                if getattr(shape, "shape_id", None) == shape_id:
                    if not getattr(shape, "has_text_frame", False):
                        raise ValueError(f"shape_id={shape_id} 不是文字框（has_text_frame=False）")
                    return shape, idx
            raise ValueError(f"找不到 shape_id={shape_id}")

        assert shape_index is not None
        if shape_index < 0 or shape_index >= len(slide.shapes):
            raise IndexError(f"shape_index 超出範圍: {shape_index}, shape_count={len(slide.shapes)}")
        shape = slide.shapes[shape_index]
        if not getattr(shape, "has_text_frame", False):
            raise ValueError(f"shape_index={shape_index} 不是文字框（has_text_frame=False）")
        return shape, shape_index

    def get_textbox_style(self, slide_index: int, shape_id: Optional[int] = None, shape_index: Optional[int] = None) -> Dict[str, Any]:
        shape, resolved_shape_index = self._get_textbox_shape(
            slide_index=slide_index,
            shape_id=shape_id,
            shape_index=shape_index,
        )

        fill_info = {
            "fill_type": "unknown",
            "fill_color": None,
            "fill_transparency": None,
        }
        line_info = {
            "line_style": None,
            "line_color": None,
            "line_width_emu": None,
        }
        notes: List[str] = []

        try:
            fill = shape.fill
            fill_type_raw = getattr(fill, "type", None)
            if fill_type_raw is None:
                fill_info["fill_type"] = "inherit"
            else:
                fill_type_name = str(fill_type_raw).lower()
                if "solid" in fill_type_name:
                    fill_info["fill_type"] = "solid"
                elif "pattern" in fill_type_name:
                    fill_info["fill_type"] = "pattern"
                elif "gradient" in fill_type_name:
                    fill_info["fill_type"] = "gradient"
                elif "picture" in fill_type_name:
                    fill_info["fill_type"] = "picture"
                elif "background" in fill_type_name:
                    fill_info["fill_type"] = "background"
                else:
                    fill_info["fill_type"] = fill_type_name

            fill_info["fill_color"] = _rgb_from_color(getattr(fill, "fore_color", None))
            xml_transparency = _read_fill_transparency_from_xml(shape)
            if xml_transparency is not None:
                fill_info["fill_transparency"] = xml_transparency
            else:
                try:
                    fill_info["fill_transparency"] = float(fill.transparency) if fill.transparency is not None else None
                except Exception:
                    fill_info["fill_transparency"] = None
        except Exception as exc:
            notes.append(f"讀取 fill 資訊失敗: {exc}")

        try:
            line = shape.line
            line_info["line_style"] = _dash_style_to_name(getattr(line, "dash_style", None))
            line_info["line_color"] = _rgb_from_color(getattr(line, "color", None))
            line_info["line_width_emu"] = int(line.width) if line.width is not None else None
        except Exception as exc:
            notes.append(f"讀取 line 資訊失敗: {exc}")

        return {
            "slide_index": slide_index,
            "shape_index": resolved_shape_index,
            "shape_id": getattr(shape, "shape_id", None),
            "name": getattr(shape, "name", None),
            "text_preview": (shape.text_frame.text or "").strip()[:120],
            "fill_type": fill_info["fill_type"],
            "fill_color": fill_info["fill_color"],
            "fill_transparency": fill_info["fill_transparency"],
            "line_style": line_info["line_style"],
            "line_color": line_info["line_color"],
            "line_width_emu": line_info["line_width_emu"],
            "notes": notes,
        }

    def get_slide_textbox_styles(self, slide_index: int) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)
        slide = self.prs.slides[slide_index]
        textbox_styles: List[Dict[str, Any]] = []

        for idx, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            textbox_styles.append(self.get_textbox_style(slide_index=slide_index, shape_index=idx))

        return {
            "slide_index": slide_index,
            "shape_count": len(slide.shapes),
            "textbox_count": len(textbox_styles),
            "textboxes": textbox_styles,
        }

    def set_textbox_style(
            self,
            slide_index: int,
            shape_id: Optional[int] = None,
            shape_index: Optional[int] = None,
            fill_color: Optional[Tuple[int, int, int]] = None,
            fill_transparency: Optional[float] = None,
            line_style: Optional[str] = None,
            line_color: Optional[Tuple[int, int, int]] = None,
            line_width: Optional[int] = None,
        ) -> Dict[str, Any]:
        shape, resolved_shape_index = self._get_textbox_shape(
            slide_index=slide_index,
            shape_id=shape_id,
            shape_index=shape_index,
        )

        if fill_transparency is not None and not (0.0 <= float(fill_transparency) <= 1.0):
            raise ValueError("fill_transparency 必須介於 0.0 ~ 1.0")
        if line_width is not None and line_width < 0:
            raise ValueError("line_width 不可小於 0")

        notes: List[str] = []

        if fill_color is not None or fill_transparency is not None:
            shape.fill.solid()
            if fill_color is not None:
                shape.fill.fore_color.rgb = _rgb_tuple_to_color(fill_color)
            if fill_transparency is not None:
                try:
                    # 優先走 OOXML a:alpha，跨渲染器行為較穩定
                    _write_fill_transparency_to_xml(shape, float(fill_transparency))
                except Exception as exc:
                    notes.append(f"設定 XML fill_transparency 失敗，改用 fallback: {exc}")
                    try:
                        shape.fill.transparency = float(fill_transparency)
                    except Exception as sub_exc:
                        notes.append(f"設定 fallback fill_transparency 失敗: {sub_exc}")

        if line_color is not None:
            shape.line.color.rgb = _rgb_tuple_to_color(line_color)
        if line_width is not None:
            shape.line.width = Emu(line_width)
        if line_style is not None:
            shape.line.dash_style = _name_to_dash_style(line_style)

        result = self.get_textbox_style(slide_index=slide_index, shape_index=resolved_shape_index)
        result["notes"].extend(notes)
        return result

    def delete_textbox(
            self,
            slide_index: int,
            shape_id: Optional[int] = None,
            shape_index: Optional[int] = None,
        ) -> Dict[str, Any]:
        shape, resolved_shape_index = self._get_textbox_shape(
            slide_index=slide_index,
            shape_id=shape_id,
            shape_index=shape_index,
        )

        deleted_shape_id = getattr(shape, "shape_id", None)
        deleted_name = getattr(shape, "name", None)
        deleted_text_preview = (shape.text_frame.text or "").strip()[:120]

        element = shape._element
        parent = element.getparent()
        if parent is None:
            raise RuntimeError("找不到 shape parent，無法刪除文字框")
        parent.remove(element)

        return {
            "slide_index": slide_index,
            "deleted_shape_index": resolved_shape_index,
            "deleted_shape_id": deleted_shape_id,
            "deleted_name": deleted_name,
            "deleted_text_preview": deleted_text_preview,
            "remaining_shape_count": len(self.prs.slides[slide_index].shapes),
        }

    def get_slide_text_fonts(self, slide_index: int) -> Dict[str, Any]:
        """
        讀取指定頁文字來源（文字框與表格 cell）內各 run 的字型資訊。
        """
        _validate_slide_index(self.prs, slide_index)

        slide = self.prs.slides[slide_index]
        font_counter: Dict[str, int] = {}
        unresolved_run_count = 0
        shape_results: List[Dict[str, Any]] = []
        text_shape_count = 0
        table_shape_count = 0
        table_text_cell_count = 0

        def _collect_paragraph_runs(paragraph, paragraph_info: Dict[str, Any]) -> None:
            nonlocal unresolved_run_count
            if len(paragraph.runs) == 0:
                if (paragraph.text or "").strip():
                    unresolved_run_count += 1
                paragraph_info["runs"].append(
                    {
                        "run_index": 0,
                        "text": paragraph.text or "",
                        "font_name": None,
                        "font_size_pt": None,
                        "bold": None,
                        "italic": None,
                    }
                )
                return

            for run_index, run in enumerate(paragraph.runs):
                run_text = run.text or ""
                font_name = run.font.name
                font_size_pt = float(run.font.size.pt) if run.font.size is not None else None
                bold = run.font.bold
                italic = run.font.italic

                if font_name:
                    font_counter[font_name] = font_counter.get(font_name, 0) + 1
                elif run_text.strip():
                    unresolved_run_count += 1

                paragraph_info["runs"].append(
                    {
                        "run_index": run_index,
                        "text": run_text,
                        "font_name": font_name,
                        "font_size_pt": font_size_pt,
                        "bold": bold,
                        "italic": italic,
                    }
                )

        for shape_index, shape in enumerate(slide.shapes):
            has_text_frame = bool(getattr(shape, "has_text_frame", False))
            has_table = bool(getattr(shape, "has_table", False))
            if not has_text_frame and not has_table:
                continue

            shape_info: Dict[str, Any] = {
                "shape_index": shape_index,
                "shape_id": getattr(shape, "shape_id", None),
                "name": getattr(shape, "name", None),
                "text_preview": "",
                "kind": "text_frame" if has_text_frame else "table",
                "paragraphs": [],
            }

            if has_text_frame:
                text_shape_count += 1
                tf = shape.text_frame
                try:
                    shape_info["text_preview"] = (tf.text or "").strip()[:120]
                except Exception:
                    shape_info["text_preview"] = ""

                for paragraph_index, paragraph in enumerate(tf.paragraphs):
                    paragraph_info: Dict[str, Any] = {
                        "paragraph_index": paragraph_index,
                        "text": paragraph.text or "",
                        "runs": [],
                    }
                    _collect_paragraph_runs(paragraph, paragraph_info)
                    shape_info["paragraphs"].append(paragraph_info)

            if has_table:
                table_shape_count += 1
                preview_texts: List[str] = []
                for row_index, row in enumerate(shape.table.rows):
                    for col_index, cell in enumerate(row.cells):
                        cell_text = cell.text or ""
                        if cell_text.strip():
                            table_text_cell_count += 1
                            if len(preview_texts) < 3:
                                preview_texts.append(cell_text.strip())

                        for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs):
                            paragraph_info = {
                                "row_index": row_index,
                                "col_index": col_index,
                                "paragraph_index": paragraph_index,
                                "text": paragraph.text or "",
                                "runs": [],
                            }
                            _collect_paragraph_runs(paragraph, paragraph_info)
                            shape_info["paragraphs"].append(paragraph_info)

                if not shape_info["text_preview"]:
                    shape_info["text_preview"] = " | ".join(preview_texts)[:120]

            shape_results.append(shape_info)

        font_summary = [
            {"font_name": name, "count": count}
            for name, count in sorted(font_counter.items(), key=lambda x: (-x[1], x[0]))
        ]

        notes: List[str] = []
        if unresolved_run_count > 0:
            notes.append("部分文字 run 未設定字型（可能沿用母片/樣式）。")

        return {
            "slide_index": slide_index,
            "shape_count": len(slide.shapes),
            "text_shape_count": text_shape_count,
            "table_shape_count": table_shape_count,
            "table_text_cell_count": table_text_cell_count,
            "detected_font_count": len(font_summary),
            "unresolved_run_count": unresolved_run_count,
            "font_summary": font_summary,
            "shapes": shape_results,
            "notes": notes,
        }

    def scan_presentation_text_fonts(self) -> Dict[str, Any]:
        """
        掃描整份簡報每一頁的文字字型資訊。
        """
        slide_count = len(self.prs.slides)
        slides: List[Dict[str, Any]] = []
        aggregate_counter: Dict[str, int] = {}
        unresolved_run_count = 0

        for slide_index in range(slide_count):
            slide_info = self.get_slide_text_fonts(slide_index)
            slides.append(slide_info)

            unresolved_run_count += slide_info.get("unresolved_run_count", 0)
            for item in slide_info.get("font_summary", []):
                font_name = item.get("font_name")
                count = int(item.get("count", 0))
                if font_name:
                    aggregate_counter[font_name] = aggregate_counter.get(font_name, 0) + count

        font_summary = [
            {"font_name": name, "count": count}
            for name, count in sorted(aggregate_counter.items(), key=lambda x: (-x[1], x[0]))
        ]

        return {
            "file_path": self.file_path,
            "slide_count": slide_count,
            "detected_font_count": len(font_summary),
            "unresolved_run_count": unresolved_run_count,
            "font_summary": font_summary,
            "slides": slides,
        }

    def _get_shape(self, slide_index: int, shape_id: Optional[int] = None, shape_index: Optional[int] = None):
        _validate_slide_index(self.prs, slide_index)
        slide = self.prs.slides[slide_index]

        if shape_id is None and shape_index is None:
            raise ValueError("shape_id 與 shape_index 至少需提供一個")

        if shape_id is not None:
            for idx, shape in enumerate(slide.shapes):
                if getattr(shape, "shape_id", None) == shape_id:
                    return shape, idx
            raise ValueError(f"找不到 shape_id={shape_id}")

        assert shape_index is not None
        if shape_index < 0 or shape_index >= len(slide.shapes):
            raise IndexError(f"shape_index 超出範圍: {shape_index}, shape_count={len(slide.shapes)}")

        return slide.shapes[shape_index], shape_index

    def get_shape_style(self, slide_index: int, shape_id: Optional[int] = None, shape_index: Optional[int] = None) -> Dict[str, Any]:
        shape, resolved_shape_index = self._get_shape(
            slide_index=slide_index,
            shape_id=shape_id,
            shape_index=shape_index,
        )

        fill_info = {
            "fill_type": "unknown",
            "fill_color": None,
            "fill_transparency": None,
        }
        line_info = {
            "line_style": None,
            "line_color": None,
            "line_width_emu": None,
        }
        notes: List[str] = []

        try:
            fill = shape.fill
            fill_type_raw = getattr(fill, "type", None)
            if fill_type_raw is None:
                fill_info["fill_type"] = "inherit"
            else:
                fill_type_name = str(fill_type_raw).lower()
                if "solid" in fill_type_name:
                    fill_info["fill_type"] = "solid"
                elif "pattern" in fill_type_name:
                    fill_info["fill_type"] = "pattern"
                elif "gradient" in fill_type_name:
                    fill_info["fill_type"] = "gradient"
                elif "picture" in fill_type_name:
                    fill_info["fill_type"] = "picture"
                elif "background" in fill_type_name:
                    fill_info["fill_type"] = "background"
                else:
                    fill_info["fill_type"] = fill_type_name

            fill_info["fill_color"] = _rgb_from_color(getattr(fill, "fore_color", None))
            xml_transparency = _read_fill_transparency_from_xml(shape)
            if xml_transparency is not None:
                fill_info["fill_transparency"] = xml_transparency
            else:
                try:
                    fill_info["fill_transparency"] = float(fill.transparency) if fill.transparency is not None else None
                except Exception:
                    fill_info["fill_transparency"] = None
        except Exception as exc:
            notes.append(f"讀取 fill 資訊失敗: {exc}")

        try:
            line = shape.line
            line_info["line_style"] = _dash_style_to_name(getattr(line, "dash_style", None))
            line_info["line_color"] = _rgb_from_color(getattr(line, "color", None))
            line_info["line_width_emu"] = int(line.width) if line.width is not None else None
        except Exception as exc:
            notes.append(f"讀取 line 資訊失敗: {exc}")

        text_preview = None
        try:
            if getattr(shape, "has_text_frame", False):
                text_preview = (shape.text_frame.text or "").strip()[:120]
        except Exception:
            text_preview = None

        return {
            "slide_index": slide_index,
            "shape_index": resolved_shape_index,
            "shape_id": getattr(shape, "shape_id", None),
            "name": getattr(shape, "name", None),
            "shape_type": str(getattr(shape, "shape_type", "")),
            "has_text_frame": bool(getattr(shape, "has_text_frame", False)),
            "text_preview": text_preview,
            "fill_type": fill_info["fill_type"],
            "fill_color": fill_info["fill_color"],
            "fill_transparency": fill_info["fill_transparency"],
            "line_style": line_info["line_style"],
            "line_color": line_info["line_color"],
            "line_width_emu": line_info["line_width_emu"],
            "notes": notes,
        }

    def set_shape_style(
            self,
            slide_index: int,
            shape_id: Optional[int] = None,
            shape_index: Optional[int] = None,
            fill_color: Optional[Tuple[int, int, int]] = None,
            fill_transparency: Optional[float] = None,
            line_style: Optional[str] = None,
            line_color: Optional[Tuple[int, int, int]] = None,
            line_width: Optional[int] = None,
        ) -> Dict[str, Any]:
        shape, resolved_shape_index = self._get_shape(
            slide_index=slide_index,
            shape_id=shape_id,
            shape_index=shape_index,
        )

        if fill_transparency is not None and not (0.0 <= float(fill_transparency) <= 1.0):
            raise ValueError("fill_transparency 必須介於 0.0 ~ 1.0")
        if line_width is not None and line_width < 0:
            raise ValueError("line_width 不可小於 0")

        notes: List[str] = []

        if fill_color is not None or fill_transparency is not None:
            shape.fill.solid()
            if fill_color is not None:
                shape.fill.fore_color.rgb = _rgb_tuple_to_color(fill_color)
            if fill_transparency is not None:
                try:
                    _write_fill_transparency_to_xml(shape, float(fill_transparency))
                except Exception as exc:
                    notes.append(f"設定 XML fill_transparency 失敗，改用 fallback: {exc}")
                    try:
                        shape.fill.transparency = float(fill_transparency)
                    except Exception as sub_exc:
                        notes.append(f"設定 fallback fill_transparency 失敗: {sub_exc}")

        if line_color is not None:
            shape.line.color.rgb = _rgb_tuple_to_color(line_color)
        if line_width is not None:
            shape.line.width = Emu(line_width)
        if line_style is not None:
            shape.line.dash_style = _name_to_dash_style(line_style)

        result = self.get_shape_style(slide_index=slide_index, shape_index=resolved_shape_index)
        result["notes"].extend(notes)
        return result

    def set_shape_fill_transparency(
            self,
            slide_index: int,
            shape_id: Optional[int] = None,
            shape_index: Optional[int] = None,
            fill_transparency: float = 0.4,
            fill_color: Optional[Tuple[int, int, int]] = None,
        ) -> Dict[str, Any]:
        return self.set_shape_style(
            slide_index=slide_index,
            shape_id=shape_id,
            shape_index=shape_index,
            fill_color=fill_color,
            fill_transparency=fill_transparency,
        )

    def delete_slide(self, slide_index: int) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)

        slide_id_list = self.prs.slides._sldIdLst
        slide_id = slide_id_list[slide_index]
        r_id = slide_id.rId
        self.prs.part.drop_rel(r_id)
        del slide_id_list[slide_index]

        return {
            "deleted_slide_index": slide_index,
            "remaining_slide_count": len(self.prs.slides),
        }

    def duplicate_slide(self, slide_index: int) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)

        source_slide = self.prs.slides[slide_index]
        blank_layout = _get_blank_layout(self.prs)
        new_slide = self.prs.slides.add_slide(blank_layout)

        # 複製 shape XML
        for shape in source_slide.shapes:
            new_el = deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

        # 複製背景（盡量）
        try:
            if source_slide.background.fill.type is not None:
                src_fill = source_slide.background.fill
                dst_fill = new_slide.background.fill
                if hasattr(src_fill, "fore_color") and src_fill.fore_color.type is not None:
                    dst_fill.solid()
                    if hasattr(src_fill.fore_color, "rgb") and src_fill.fore_color.rgb:
                        dst_fill.fore_color.rgb = src_fill.fore_color.rgb
        except Exception:
            pass

        return {
            "source_slide_index": slide_index,
            "new_slide_index": len(self.prs.slides) - 1,
            "slide_count": len(self.prs.slides),
        }

    def replace_text(
            self,
            old_text: str,
            new_text: str,
            slide_indices: Optional[List[int]] = None,
            exact_match: bool = False,
            case_sensitive: bool = True,
        ) -> Dict[str, Any]:
        if not old_text:
            raise ValueError("old_text 不可為空")

        if slide_indices is None:
            target_indices = list(range(len(self.prs.slides)))
        else:
            target_indices = slide_indices
            for idx in target_indices:
                _validate_slide_index(self.prs, idx)

        total_replacements = 0
        matched_shapes = []

        def _replace(src: str) -> Tuple[str, int]:
            if exact_match:
                if case_sensitive:
                    if src == old_text:
                        return new_text, 1
                    return src, 0
                else:
                    if src.lower() == old_text.lower():
                        return new_text, 1
                    return src, 0
            else:
                if case_sensitive:
                    count = src.count(old_text)
                    return src.replace(old_text, new_text), count
                else:
                    import re
                    pattern = re.compile(re.escape(old_text), re.IGNORECASE)
                    replaced, count = pattern.subn(new_text, src)
                    return replaced, count

        for slide_index in target_indices:
            slide = self.prs.slides[slide_index]
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue

                changed = False
                shape_replace_count = 0

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original = run.text or ""
                        replaced, cnt = _replace(original)
                        if cnt > 0:
                            run.text = replaced
                            shape_replace_count += cnt
                            changed = True

                    # 如果 paragraph 沒有 runs，直接處理 paragraph.text
                    if len(paragraph.runs) == 0:
                        original = paragraph.text or ""
                        replaced, cnt = _replace(original)
                        if cnt > 0:
                            paragraph.text = replaced
                            shape_replace_count += cnt
                            changed = True

                if changed:
                    total_replacements += shape_replace_count
                    matched_shapes.append({
                        "slide_index": slide_index,
                        "shape_id": getattr(shape, "shape_id", None),
                        "replace_count": shape_replace_count,
                    })

        return {
            "old_text": old_text,
            "new_text": new_text,
            "total_replacements": total_replacements,
            "matched_shapes": matched_shapes,
        }

    def add_bullets(
            self,
            slide_index: int,
            items: List[str],
            left: int,
            top: int,
            width: int,
            height: int,
            font_size: int = 20,
            level: int = 0,
            bold: bool = False,
            font_name: Optional[str] = None,
            font_color: Optional[Tuple[int, int, int]] = None,
        ) -> Dict[str, Any]:
        _validate_slide_index(self.prs, slide_index)

        if not items:
            raise ValueError("items 不可為空")

        slide = self.prs.slides[slide_index]
        shape = slide.shapes.add_textbox(
            Emu(left),
            Emu(top),
            Emu(width),
            Emu(height),
        )

        tf = shape.text_frame
        tf.clear()

        for idx, item in enumerate(items):
            if idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = "" if item is None else str(item)
            p.level = max(0, int(level))
            p.bullet = True

            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                if font_name:
                    run.font.name = font_name
                color = _rgb_tuple_to_color(font_color)
                if color:
                    run.font.color.rgb = color

        return {
            "slide_index": slide_index,
            "shape_id": shape.shape_id,
            "item_count": len(items),
            "left": left,
            "top": top,
            "width": width,
            "height": height,
        }

    def add_title_slide(
            self,
            title: str,
            subtitle: str = "",
        ) -> Dict[str, Any]:
        layout = None

        # 盡量找 Title Slide layout
        for candidate in self.prs.slide_layouts:
            try:
                placeholder_types = []
                for ph in candidate.placeholders:
                    try:
                        placeholder_types.append(ph.placeholder_format.type)
                    except Exception:
                        pass
                if PP_PLACEHOLDER.TITLE in placeholder_types:
                    layout = candidate
                    break
            except Exception:
                continue

        if layout is None:
            layout = self.prs.slide_layouts[0]

        slide = self.prs.slides.add_slide(layout)

        if slide.shapes.title is not None:
            slide.shapes.title.text = title or ""

        # 找副標題 placeholder
        subtitle_set = False
        if subtitle:
            for shape in slide.placeholders:
                try:
                    ph_type = shape.placeholder_format.type
                    if ph_type == PP_PLACEHOLDER.SUBTITLE:
                        shape.text = subtitle
                        subtitle_set = True
                        break
                except Exception:
                    pass

            if not subtitle_set:
                if not subtitle_set and subtitle:
                    tb = slide.shapes.add_textbox(
                        left=Emu(500000),
                        top=Emu(1500000),
                        width=Emu(6000000),
                        height=Emu(800000),
                    )
                    tb.text_frame.text = subtitle

        return {
            "slide_index": len(self.prs.slides) - 1,
            "title": title,
            "subtitle": subtitle,
        }

    def reorder_slides(self, new_order: List[int]) -> Dict[str, Any]:
        slide_count = len(self.prs.slides)

        if len(new_order) != slide_count:
            raise ValueError(f"new_order 長度必須等於目前頁數 {slide_count}")

        if sorted(new_order) != list(range(slide_count)):
            raise ValueError("new_order 必須是 0 到 slide_count-1 的完整排列")

        sld_id_lst = self.prs.slides._sldIdLst
        current_nodes = list(sld_id_lst)

        # 先清空再依新順序放回
        for _ in range(len(sld_id_lst)):
            del sld_id_lst[0]

        for old_index in new_order:
            sld_id_lst.append(current_nodes[old_index])

        return {
            "slide_count": slide_count,
            "new_order": new_order,
        }

def new(
        file_path: str,
        plank_page_num: int = 1,
        plank_page_width: int = 1080,
        plank_page_height: int = 1920,
        dpi: int = DEFAULT_DPI,
    ) -> str:
    """
    建立新的 pptx 檔案

    Args:
        file_path: 輸出檔案路徑
        plank_page_num: 頁面數量
        plank_page_width: 頁面寬度（px）
        plank_page_height: 頁面高度（px）
        dpi: 像素轉尺寸時使用的 DPI

    Returns:
        str: 建立後的檔案路徑
    """
    doc = new_document(
        file_path=file_path,
        plank_page_num=plank_page_num,
        plank_page_width=plank_page_width,
        plank_page_height=plank_page_height,
        dpi=dpi,
    )
    return doc.save(file_path)

def new_document(
        file_path: Optional[str] = None,
        plank_page_num: int = 1,
        plank_page_width: int = 1080,
        plank_page_height: int = 1920,
        dpi: int = DEFAULT_DPI,
    ) -> PPTDocument:
    _ensure_pptx_available()

    if plank_page_num < 1:
        raise ValueError("plank_page_num 至少要 >= 1")
    if plank_page_width <= 0 or plank_page_height <= 0:
        raise ValueError("頁面寬高必須 > 0")

    prs = Presentation()
    prs.slide_width = _px_to_emu(plank_page_width, dpi=dpi)
    prs.slide_height = _px_to_emu(plank_page_height, dpi=dpi)

    _remove_all_slides(prs)

    doc = PPTDocument(prs=prs, file_path=file_path)
    doc.add_blank_slides(plank_page_num)
    return doc

def open_presentation(file_path: str) -> PPTDocument:
    _ensure_pptx_available()

    file_path = _normalize_file_path(file_path)
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"檔案不存在: {file_path}")

    prs = Presentation(file_path)
    return PPTDocument(prs=prs, file_path=file_path)


def save(document: PPTDocument, file_path: Optional[str] = None) -> str:
    return document.save(file_path=file_path)


def add_blank_slide(document: PPTDocument) -> int:
    return document.add_blank_slide()


def add_blank_slides(document: PPTDocument, page_num: int = 1) -> List[int]:
    return document.add_blank_slides(page_num=page_num)


def add_text(
        document: PPTDocument,
        slide_index: int,
        text: str,
        left: int,
        top: int,
        width: int,
        height: int,
        font_size: int = 20,
        bold: bool = False,
        italic: bool = False,
        font_name: Optional[str] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        align: str = "left",
    ) -> Dict[str, Any]:
    return document.add_textbox(
        slide_index=slide_index,
        text=text,
        left=left,
        top=top,
        width=width,
        height=height,
        font_size=font_size,
        bold=bold,
        italic=italic,
        font_name=font_name,
        font_color=font_color,
        align=align,
    )


# ---------------------------------------------------------------------------
# 文字藝術師 / WordArt 類需求（骨架）：規格見 issues/文字藝術師.iss 第一層 API
# 實作「範本複製」或 OOXML 直改時，可斟酌啟用（模組前段 try 區可能已匯入部分）：
# from copy import deepcopy
# from pptx.oxml.ns import qn
# from pptx.oxml import parse_xml
# ---------------------------------------------------------------------------


def add_wordart_like_textbox(
        document: PPTDocument,
        slide_index: int,
        text: str,
        left: int,
        top: int,
        width: int,
        height: int,
        font_size: int = 28,
        bold: bool = True,
        font_name: Optional[str] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        align: str = "center",
    ) -> Dict[str, Any]:
    """
    以一般文字框模擬 WordArt 風格（大字、描邊／底色等由後續參數擴充）。

    實作提示：
    - 底層可呼叫既有 add_text / add_textbox 邏輯並加強字型與外框
    - 非 PowerPoint 原生 WordArt 物件；複雜特效需改 OOXML 或改走 clone_named_shape_from_template
    """
    _ensure_pptx_available()
    _validate_slide_index(document.prs, slide_index)
    result = document.add_textbox(
        slide_index=slide_index,
        text=text,
        left=left,
        top=top,
        width=width,
        height=height,
        font_size=font_size,
        bold=bold,
        italic=False,
        font_name=font_name,
        font_color=font_color,
        align=align,
    )
    result["notes"] = ["以 TextBox 模擬 WordArt；複雜特效請改用範本複製或 OOXML。"]
    return result


def update_wordart_text(
        document: PPTDocument,
        slide_index: int,
        new_text: str,
        shape_name: Optional[str] = None,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    """
    更新現有「類 WordArt」或範本複製 shape 的文字；優先只改 run.text 以保留特效。

    實作提示：
    - 依 shape_name 或 shape_id / shape_index 解析目標 shape（與 delete_textbox 解析策略對齊）
    - 僅在 has_text_frame 時更新；否則 notes 說明並不中斷其他欄位處理（若有多段）
    """
    _ensure_pptx_available()
    _validate_slide_index(document.prs, slide_index)
    slide = document.prs.slides[slide_index]

    resolved_shape_index: Optional[int] = None
    shape = None

    if shape_name is not None:
        target_name = str(shape_name)
        for idx, candidate in enumerate(slide.shapes):
            if getattr(candidate, "name", None) == target_name:
                shape = candidate
                resolved_shape_index = idx
                break
        if shape is None:
            raise ValueError(f"找不到 shape_name={target_name}")
    else:
        shape, resolved_shape_index = document._get_shape(
            slide_index=slide_index,
            shape_id=shape_id,
            shape_index=shape_index,
        )

    notes: List[str] = []
    # 保留原始邊框 XML，避免改字時意外改動 outline。
    original_line_xml = None
    had_explicit_line = False
    try:
        sp_pr = shape._element.spPr
        if sp_pr is not None:
            line_node = sp_pr.find(qn("a:ln"))
            if line_node is not None:
                had_explicit_line = True
                original_line_xml = deepcopy(line_node)
    except Exception as exc:
        notes.append(f"讀取原始邊框設定失敗，將略過邊框保留: {exc}")

    before_text_preview = None
    if getattr(shape, "has_text_frame", False):
        try:
            before_text_preview = (shape.text_frame.text or "").strip()[:120]
        except Exception:
            before_text_preview = None

    if not getattr(shape, "has_text_frame", False):
        notes.append("目標 shape 無 text_frame，未更新文字。")
        return {
            "slide_index": slide_index,
            "shape_index": resolved_shape_index,
            "shape_id": getattr(shape, "shape_id", None),
            "shape_name": getattr(shape, "name", None),
            "updated": False,
            "before_text_preview": before_text_preview,
            "after_text_preview": before_text_preview,
            "notes": notes,
        }

    text_frame = shape.text_frame
    paragraphs = list(text_frame.paragraphs)
    if len(paragraphs) == 0:
        text_frame.text = "" if new_text is None else str(new_text)
        notes.append("原文字框沒有 paragraph，改以 text_frame.text fallback 寫入。")
    else:
        target_text = "" if new_text is None else str(new_text)
        first_paragraph = paragraphs[0]
        first_runs = list(first_paragraph.runs)
        if len(first_runs) > 0:
            first_runs[0].text = target_text
            for run in first_runs[1:]:
                run.text = ""
        else:
            first_paragraph.text = target_text
            notes.append("原第一段沒有 runs，改以 paragraph.text fallback 寫入。")

        for paragraph in paragraphs[1:]:
            for run in paragraph.runs:
                run.text = ""
            if len(paragraph.runs) == 0:
                paragraph.text = ""

    try:
        sp_pr = shape._element.spPr
        if sp_pr is not None:
            current_line_node = sp_pr.find(qn("a:ln"))
            if had_explicit_line:
                if current_line_node is not None:
                    sp_pr.remove(current_line_node)
                if original_line_xml is not None:
                    sp_pr.append(deepcopy(original_line_xml))
            else:
                if current_line_node is not None:
                    sp_pr.remove(current_line_node)
    except Exception as exc:
        notes.append(f"還原原始邊框設定失敗: {exc}")

    after_text_preview = None
    try:
        after_text_preview = (shape.text_frame.text or "").strip()[:120]
    except Exception:
        after_text_preview = None

    return {
        "slide_index": slide_index,
        "shape_index": resolved_shape_index,
        "shape_id": getattr(shape, "shape_id", None),
        "shape_name": getattr(shape, "name", None),
        "updated": True,
        "before_text_preview": before_text_preview,
        "after_text_preview": after_text_preview,
        "notes": notes,
    }


def delete_shape(
        document: PPTDocument,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    """
    刪除投影片上任意一個 shape（含以圖片／群組偽裝者）；本質為從 XML parent 移除元素。

    實作提示：
    - 可參考 issues：`el = shape._element; parent.remove(el)`
    - 與 delete_textbox 差異：不限定文字方塊；需處理找不到 shape 的錯誤訊息
    """
    _ensure_pptx_available()
    shape, resolved_shape_index = document._get_shape(
        slide_index=slide_index,
        shape_id=shape_id,
        shape_index=shape_index,
    )

    deleted_shape_id = getattr(shape, "shape_id", None)
    deleted_name = getattr(shape, "name", None)
    deleted_shape_type = str(getattr(shape, "shape_type", ""))
    deleted_text_preview = None
    if getattr(shape, "has_text_frame", False):
        try:
            deleted_text_preview = (shape.text_frame.text or "").strip()[:120]
        except Exception:
            deleted_text_preview = None

    element = shape._element
    parent = element.getparent()
    if parent is None:
        raise RuntimeError("找不到 shape parent，無法刪除")
    parent.remove(element)

    return {
        "slide_index": slide_index,
        "deleted_shape_index": resolved_shape_index,
        "deleted_shape_id": deleted_shape_id,
        "deleted_name": deleted_name,
        "deleted_shape_type": deleted_shape_type,
        "deleted_text_preview": deleted_text_preview,
        "remaining_shape_count": len(document.prs.slides[slide_index].shapes),
    }


def clone_named_shape_from_template(
        document: PPTDocument,
        slide_index: int,
        shape_name: str,
        new_text: str = "",
        left: Optional[int] = None,
        top: Optional[int] = None,
    ) -> Dict[str, Any]:
    """
    從指定頁依 shape.name 找到範本 shape，複製到同頁（或指定位置），並可選填新文字。

    issues 原簽名以 slide 為單位；此處改為 document + slide_index 以符合 ppt_stdio 慣例。

    實作提示：
    - 複製可 deepcopy(shape._element) 後插入 spTree，並處理 rId / 媒體關聯
    - 若範本在「另一份簡報」，可另增參數 template_path（此處先以註解預留擴充點）
    """
    _ensure_pptx_available()
    _validate_slide_index(document.prs, slide_index)
    if not shape_name:
        raise ValueError("shape_name 不可為空")

    slide = document.prs.slides[slide_index]
    source_shape = None
    source_shape_index = None
    for idx, candidate in enumerate(slide.shapes):
        if getattr(candidate, "name", None) == shape_name:
            source_shape = candidate
            source_shape_index = idx
            break
    if source_shape is None:
        raise ValueError(f"找不到 shape_name={shape_name}")

    before_count = len(slide.shapes)
    new_el = deepcopy(source_shape._element)
    slide.shapes._spTree.insert_element_before(new_el, "p:extLst")

    if len(slide.shapes) <= before_count:
        raise RuntimeError("shape 複製失敗：新增後 shape 數量未增加")

    cloned_shape = slide.shapes[before_count]
    cloned_shape_index = before_count
    notes: List[str] = []

    if left is not None:
        cloned_shape.left = Emu(left)
    if top is not None:
        cloned_shape.top = Emu(top)

    if new_text:
        update_res = update_wordart_text(
            document=document,
            slide_index=slide_index,
            new_text=new_text,
            shape_id=getattr(cloned_shape, "shape_id", None),
        )
        notes.extend(update_res.get("notes", []))

    # TODO: 若要支援「跨簡報範本複製」，可新增 template_path 並補 rels 搬移流程。
    return {
        "slide_index": slide_index,
        "source_shape_index": source_shape_index,
        "source_shape_id": getattr(source_shape, "shape_id", None),
        "source_shape_name": getattr(source_shape, "name", None),
        "cloned_shape_index": cloned_shape_index,
        "cloned_shape_id": getattr(cloned_shape, "shape_id", None),
        "cloned_shape_name": getattr(cloned_shape, "name", None),
        "left": int(cloned_shape.left),
        "top": int(cloned_shape.top),
        "new_text": new_text,
        "notes": notes,
    }


def add_image(
        document: PPTDocument,
        slide_index: int,
        image_path: str,
        left: int,
        top: int,
        width: Optional[int] = None,
        height: Optional[int] = None,
        keep_aspect_ratio: bool = True,
    ) -> Dict[str, Any]:
    return document.add_image(
        slide_index=slide_index,
        image_path=image_path,
        left=left,
        top=top,
        width=width,
        height=height,
        keep_aspect_ratio=keep_aspect_ratio,
    )


def add_table(
        document: PPTDocument,
        slide_index: int,
        rows: int,
        cols: int,
        left: int,
        top: int,
        width: int,
        height: int,
        data: Optional[List[List[Any]]] = None,
        first_row_as_header: bool = False,
        font_size: int = 14,
    ) -> Dict[str, Any]:
    return document.add_table(
        slide_index=slide_index,
        rows=rows,
        cols=cols,
        left=left,
        top=top,
        width=width,
        height=height,
        data=data,
        first_row_as_header=first_row_as_header,
        font_size=font_size,
    )


def add_shape(
        document: PPTDocument,
        slide_index: int,
        shape_type: str,
        left: int,
        top: int,
        width: int,
        height: int,
        text: str = "",
        fill_color: Optional[Tuple[int, int, int]] = None,
        line_color: Optional[Tuple[int, int, int]] = None,
        line_width: Optional[int] = None,
        font_size: int = 18,
        bold: bool = False,
        font_name: Optional[str] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
    ) -> Dict[str, Any]:
    return document.add_shape(
        slide_index=slide_index,
        shape_type=shape_type,
        left=left,
        top=top,
        width=width,
        height=height,
        text=text,
        fill_color=fill_color,
        line_color=line_color,
        line_width=line_width,
        font_size=font_size,
        bold=bold,
        font_name=font_name,
        font_color=font_color,
    )


def add_line(
        document: PPTDocument,
        slide_index: int,
        x1: int,
        y1: int,
        x2: int,
        y2: int,
        line_color: Optional[Tuple[int, int, int]] = None,
        line_width: Optional[int] = None,
    ) -> Dict[str, Any]:
    return document.add_line(
        slide_index=slide_index,
        x1=x1,
        y1=y1,
        x2=x2,
        y2=y2,
        line_color=line_color,
        line_width=line_width,
    )


def add_arrow(
        document: PPTDocument,
        slide_index: int,
        left: int,
        top: int,
        width: int,
        height: int,
        direction: str = "right",
        text: str = "",
        fill_color: Optional[Tuple[int, int, int]] = None,
        line_color: Optional[Tuple[int, int, int]] = None,
        line_width: Optional[int] = None,
        font_size: int = 18,
        bold: bool = False,
        font_name: Optional[str] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
    ) -> Dict[str, Any]:
    return document.add_arrow(
        slide_index=slide_index,
        left=left,
        top=top,
        width=width,
        height=height,
        direction=direction,
        text=text,
        fill_color=fill_color,
        line_color=line_color,
        line_width=line_width,
        font_size=font_size,
        bold=bold,
        font_name=font_name,
        font_color=font_color,
    )


def get_info(document: PPTDocument) -> Dict[str, Any]:
    return document.get_info()


def list_slides(document: PPTDocument) -> List[Dict[str, Any]]:
    return document.list_slides()


def get_textbox_style(
        document: PPTDocument,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    return document.get_textbox_style(
        slide_index=slide_index,
        shape_id=shape_id,
        shape_index=shape_index,
    )


def get_slide_textbox_styles(document: PPTDocument, slide_index: int) -> Dict[str, Any]:
    return document.get_slide_textbox_styles(slide_index=slide_index)


def set_textbox_style(
        document: PPTDocument,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        fill_transparency: Optional[float] = None,
        line_style: Optional[str] = None,
        line_color: Optional[Tuple[int, int, int]] = None,
        line_width: Optional[int] = None,
    ) -> Dict[str, Any]:
    return document.set_textbox_style(
        slide_index=slide_index,
        shape_id=shape_id,
        shape_index=shape_index,
        fill_color=fill_color,
        fill_transparency=fill_transparency,
        line_style=line_style,
        line_color=line_color,
        line_width=line_width,
    )


def delete_textbox(
        document: PPTDocument,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    return document.delete_textbox(
        slide_index=slide_index,
        shape_id=shape_id,
        shape_index=shape_index,
    )


def get_slide_text_fonts(document: PPTDocument, slide_index: int) -> Dict[str, Any]:
    return document.get_slide_text_fonts(slide_index=slide_index)


def scan_presentation_text_fonts(document: PPTDocument) -> Dict[str, Any]:
    return document.scan_presentation_text_fonts()


def set_slide_background_color(
        document: PPTDocument,
        slide_index: int,
        rgb: Tuple[int, int, int],
    ) -> dict:
    return document.set_slide_background_color(slide_index=slide_index, rgb=rgb)


def set_slide_background_image(
        document: PPTDocument,
        slide_index: int,
        image_path: str,
    ) -> dict:
    return document.set_slide_background_image(slide_index=slide_index, image_path=image_path)


# ---------------------------------------------------------------------------
# 佈景主題與背景讀取（骨架）：規格見 issues/布景背景偵測.iss
# ---------------------------------------------------------------------------


def _get_theme_part_info(document: PPTDocument) -> Dict[str, Any]:
    """
    從 python-pptx 的 package / relationship 找出 theme 相關 part 路徑與基本資訊。

    實作提示：
    - 可從 document.prs.part.package、slide master rels、presentation rels 尋找 theme 關聯
    - 關鍵欄位至少覆蓋：has_theme、theme_part_name、slide_master_count、notes
    - theme_part_name 格式期望類似：/ppt/theme/theme1.xml
    - 若某關聯失敗或不存在，改寫入 notes，不直接中斷整體流程
    """
    notes: List[str] = []
    theme_part_name: Optional[str] = None

    try:
        slide_master_count = len(document.prs.slide_masters)
    except Exception:
        slide_master_count = 0
        notes.append("無法取得 slide master 數量。")

    rel_sources: List[Tuple[str, Any]] = []
    try:
        rel_sources.append(("presentation", document.prs.part.rels))
    except Exception as exc:
        notes.append(f"無法讀取 presentation rels: {exc}")

    try:
        for idx, master in enumerate(document.prs.slide_masters):
            try:
                rel_sources.append((f"slide_master_{idx}", master.part.rels))
            except Exception as exc:
                notes.append(f"無法讀取 slide_master_{idx} rels: {exc}")
    except Exception as exc:
        notes.append(f"無法走訪 slide masters: {exc}")

    for source_name, rels in rel_sources:
        try:
            rel_values = list(rels.values())
        except Exception as exc:
            notes.append(f"{source_name} rels 讀取失敗: {exc}")
            continue

        for rel in rel_values:
            reltype = str(getattr(rel, "reltype", ""))
            if not reltype.endswith("/theme"):
                continue

            try:
                partname = str(rel.target_part.partname)
                if partname:
                    theme_part_name = partname if partname.startswith("/") else f"/{partname}"
                    break
            except Exception:
                try:
                    target_ref = str(getattr(rel, "target_ref", "")).strip()
                    if target_ref:
                        normalized = target_ref if target_ref.startswith("/") else f"/{target_ref}"
                        theme_part_name = normalized
                        break
                except Exception as exc:
                    notes.append(f"{source_name} theme 關聯解析失敗: {exc}")

        if theme_part_name:
            break

    if not theme_part_name:
        notes.append("找不到 theme relationship，可能為無主題或無法由目前關聯追溯。")

    return {
        "has_theme": bool(theme_part_name),
        "theme_part_name": theme_part_name,
        "slide_master_count": slide_master_count,
        "notes": notes,
    }


def _parse_theme_xml(theme_xml_bytes: bytes) -> Dict[str, Any]:
    """
    解析 theme XML（a:theme / clrScheme / fontScheme 等），抽出 color_scheme、font_scheme。

    實作提示：
    - 使用 ElementTree 與 OOXML 命名空間對應表
    - 需支援抽取 theme_name、color_scheme、font_scheme
    - 單一節點失敗不應讓整體解析中斷，可寫入 notes
    """
    from xml.etree import ElementTree as ET

    notes: List[str] = []
    theme_name: Optional[str] = None
    color_scheme: Dict[str, Any] = {}
    font_scheme: Dict[str, Any] = {}

    def _local(tag: str) -> str:
        return tag.split("}", 1)[-1]

    def _parse_font_group(group_elem: Optional[ET.Element], ns: Dict[str, str]) -> Dict[str, Any]:
        parsed: Dict[str, Any] = {"latin": None, "ea": None, "cs": None, "scripts": {}}
        if group_elem is None:
            return parsed

        for key in ("latin", "ea", "cs"):
            node = group_elem.find(f"a:{key}", ns)
            if node is not None:
                parsed[key] = node.attrib.get("typeface")

        for node in group_elem.findall("a:font", ns):
            script = node.attrib.get("script")
            typeface = node.attrib.get("typeface")
            if script and typeface:
                parsed["scripts"][script] = typeface
        return parsed

    try:
        root = ET.fromstring(theme_xml_bytes)
    except Exception as exc:
        return {
            "theme_name": None,
            "color_scheme": {},
            "font_scheme": {},
            "notes": [f"theme XML 解析失敗: {exc}"],
        }

    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    theme_name = root.attrib.get("name")

    theme_elements = root.find("a:themeElements", ns)
    if theme_elements is None:
        notes.append("themeElements 不存在。")
        return {
            "theme_name": theme_name,
            "color_scheme": color_scheme,
            "font_scheme": font_scheme,
            "notes": notes,
        }

    clr_scheme = theme_elements.find("a:clrScheme", ns)
    if clr_scheme is None:
        notes.append("clrScheme 不存在。")
    else:
        color_scheme["name"] = clr_scheme.attrib.get("name")
        values: Dict[str, Optional[str]] = {}
        for node in list(clr_scheme):
            key = _local(node.tag)
            srgb = node.find(".//a:srgbClr", ns)
            sys_clr = node.find(".//a:sysClr", ns)
            scheme_clr = node.find(".//a:schemeClr", ns)

            val: Optional[str] = None
            if srgb is not None:
                val = srgb.attrib.get("val")
            elif sys_clr is not None:
                val = sys_clr.attrib.get("lastClr") or sys_clr.attrib.get("val")
            elif scheme_clr is not None:
                scheme_val = scheme_clr.attrib.get("val")
                val = f"scheme:{scheme_val}" if scheme_val else "scheme"

            values[key] = val
        color_scheme["values"] = values

    fs = theme_elements.find("a:fontScheme", ns)
    if fs is None:
        notes.append("fontScheme 不存在。")
    else:
        font_scheme["name"] = fs.attrib.get("name")
        font_scheme["major"] = _parse_font_group(fs.find("a:majorFont", ns), ns)
        font_scheme["minor"] = _parse_font_group(fs.find("a:minorFont", ns), ns)

    return {
        "theme_name": theme_name,
        "color_scheme": color_scheme,
        "font_scheme": font_scheme,
        "notes": notes,
    }


def _get_slide_background_xml_info(document: PPTDocument, slide_index: int) -> Dict[str, Any]:
    """
    從指定頁的 slide XML / cSld / bg 等讀取背景相關原始線索（供 get_slide_background_info 彙整）。

    實作提示：
    - 區分 bgPr（背景樣式）與繼承母片的情況
    - 可搭配 document.prs.slides[slide_index].element
    - 嘗試辨識 solidFill / blipFill / gradFill / noFill / bgRef 等線索
    - 回傳原始線索時，保留 source 與 notes，避免過度推論
    """
    from xml.etree import ElementTree as ET

    _validate_slide_index(document.prs, slide_index)

    notes: List[str] = []
    result: Dict[str, Any] = {
        "slide_index": slide_index,
        "background_type": "unknown",
        "source": "slide_xml",
        "color_rgb": None,
        "image_ref": None,
        "image_path_hint": None,
        "by_shape_detection": False,
        "notes": notes,
    }

    ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    }
    rel_embed_key = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"

    def _hex_to_rgb(value: Optional[str]) -> Optional[List[int]]:
        if not value:
            return None
        v = value.strip().lstrip("#")
        if len(v) != 6:
            return None
        try:
            return [int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16)]
        except Exception:
            return None

    slide = document.prs.slides[slide_index]
    try:
        root = ET.fromstring(slide.part.blob)
    except Exception as exc:
        notes.append(f"slide XML 解析失敗: {exc}")
        return result

    bg = root.find("./p:cSld/p:bg", ns)
    if bg is None:
        result["background_type"] = "inherit"
        result["source"] = "slide_master"
        notes.append("未找到 cSld/bg，判斷為沿用母片背景。")
        return result

    bg_pr = bg.find("p:bgPr", ns)
    bg_ref = bg.find("p:bgRef", ns)

    if bg_pr is None and bg_ref is not None:
        result["background_type"] = "inherit"
        result["source"] = "slide_bgRef"
        notes.append("找到 bgRef，背景沿用母片/主題參照。")
        return result

    if bg_pr is None:
        notes.append("找到 bg 但缺少 bgPr/bgRef，無法判定背景類型。")
        return result

    solid_fill = bg_pr.find("a:solidFill", ns)
    if solid_fill is not None:
        srgb = solid_fill.find("a:srgbClr", ns)
        sys_clr = solid_fill.find("a:sysClr", ns)
        scheme_clr = solid_fill.find("a:schemeClr", ns)

        rgb = None
        if srgb is not None:
            rgb = _hex_to_rgb(srgb.attrib.get("val"))
        elif sys_clr is not None:
            rgb = _hex_to_rgb(sys_clr.attrib.get("lastClr") or sys_clr.attrib.get("val"))
        elif scheme_clr is not None:
            notes.append(f"solidFill 使用 schemeClr={scheme_clr.attrib.get('val')}，無法直接換算 RGB。")

        result["background_type"] = "solid"
        result["source"] = "slide_bgPr_solidFill"
        result["color_rgb"] = rgb
        if rgb is None:
            notes.append("solidFill 存在，但未取得可用 RGB。")
        return result

    blip_fill = bg_pr.find("a:blipFill", ns)
    if blip_fill is not None:
        result["background_type"] = "picture"
        result["source"] = "slide_bgPr_blipFill"

        blip = blip_fill.find("a:blip", ns)
        image_ref = None
        if blip is not None:
            image_ref = blip.attrib.get(rel_embed_key)
        result["image_ref"] = image_ref

        if image_ref:
            try:
                rel = slide.part.rels[image_ref]
                partname = str(rel.target_part.partname)
                result["image_path_hint"] = partname if partname.startswith("/") else f"/{partname}"
            except Exception as exc:
                notes.append(f"圖片 rId={image_ref} 無法解析到目標 part: {exc}")
        else:
            notes.append("blipFill 存在，但找不到 embed rId。")
        return result

    if bg_ref is not None:
        result["background_type"] = "inherit"
        result["source"] = "slide_bgRef"
        notes.append("背景設定以 bgRef 參照母片/主題。")
        return result

    if bg_pr.find("a:gradFill", ns) is not None:
        notes.append("偵測到 gradFill，當前版本未細分，回傳 unknown。")
    elif bg_pr.find("a:noFill", ns) is not None:
        notes.append("偵測到 noFill，背景可能仍由母片呈現，回傳 unknown。")
    else:
        notes.append("bgPr 存在但非 solid/picture，回傳 unknown。")

    return result


def _detect_full_slide_picture_shape(document: PPTDocument, slide_index: int) -> Dict[str, Any]:
    """
    偵測是否為「滿版圖片 shape 模擬背景」：位置約 (0,0)、尺寸接近整張投影片。

    實作提示：
    - 比對 shape 類型是否為圖片、left/top/width/height 與投影片寬高（EMU）
    - 回傳是否命中、候選 shape 索引、image_ref/image_path_hint、notes
    - 容差建議使用比例閾值（例如寬高落在 98%~102%）
    """
    _validate_slide_index(document.prs, slide_index)

    notes: List[str] = []
    slide = document.prs.slides[slide_index]
    slide_width = int(document.prs.slide_width)
    slide_height = int(document.prs.slide_height)
    tolerance = 0.02

    best_candidate: Optional[Dict[str, Any]] = None

    for shape_index, shape in enumerate(slide.shapes):
        if not hasattr(shape, "image"):
            continue

        try:
            left = int(shape.left)
            top = int(shape.top)
            width = int(shape.width)
            height = int(shape.height)
        except Exception:
            continue

        x_ok = abs(left) <= max(1, int(slide_width * tolerance))
        y_ok = abs(top) <= max(1, int(slide_height * tolerance))
        w_ok = abs(width - slide_width) <= max(1, int(slide_width * tolerance))
        h_ok = abs(height - slide_height) <= max(1, int(slide_height * tolerance))
        if not (x_ok and y_ok and w_ok and h_ok):
            continue

        score = (
            abs(left) / max(1, slide_width)
            + abs(top) / max(1, slide_height)
            + abs(width - slide_width) / max(1, slide_width)
            + abs(height - slide_height) / max(1, slide_height)
        )

        image_ref = None
        image_path_hint = None

        try:
            image_path_hint = getattr(shape.image, "filename", None)
        except Exception:
            image_path_hint = None

        try:
            blips = shape._element.xpath(
                ".//a:blip",
                namespaces={
                    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
                },
            )
            if blips:
                image_ref = blips[0].get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
        except Exception:
            image_ref = None

        if image_ref and not image_path_hint:
            try:
                rel = slide.part.rels[image_ref]
                partname = str(rel.target_part.partname)
                image_path_hint = partname if partname.startswith("/") else f"/{partname}"
            except Exception:
                pass

        candidate = {
            "matched": True,
            "shape_index": shape_index,
            "shape_id": getattr(shape, "shape_id", None),
            "image_ref": image_ref,
            "image_path_hint": image_path_hint,
            "by_shape_detection": True,
            "score": score,
            "notes": [],
        }

        if best_candidate is None or score < best_candidate["score"]:
            best_candidate = candidate

    if best_candidate is None:
        notes.append("未偵測到滿版圖片 shape。")
        return {
            "matched": False,
            "shape_index": None,
            "shape_id": None,
            "image_ref": None,
            "image_path_hint": None,
            "by_shape_detection": False,
            "notes": notes,
        }

    notes.append(
        f"偵測到候選滿版圖片 shape(index={best_candidate['shape_index']}, id={best_candidate['shape_id']})。"
    )
    best_candidate["notes"] = notes
    best_candidate.pop("score", None)
    return best_candidate


def get_presentation_theme_info(document: PPTDocument) -> Dict[str, Any]:
    """
    讀取簡報整體 theme / 佈景主題資訊。

    預期欄位（issues 規格）：
    - file_path, slide_count, has_theme, theme_part_name, theme_name,
      color_scheme, font_scheme, slide_master_count, notes

    注意：
    - 優先讀取真實 theme XML，不足處才以 notes 說明
    - 單一欄位讀不到不可使整體失敗
    """
    import zipfile

    _ensure_pptx_available()

    try:
        slide_count = len(document.prs.slides)
    except Exception:
        slide_count = 0

    try:
        slide_master_count = len(document.prs.slide_masters)
    except Exception:
        slide_master_count = 0

    result: Dict[str, Any] = {
        "file_path": document.file_path,
        "slide_count": slide_count,
        "has_theme": False,
        "theme_part_name": None,
        "theme_name": None,
        "color_scheme": {},
        "font_scheme": {},
        "slide_master_count": slide_master_count,
        "notes": [],
    }

    try:
        part_info = _get_theme_part_info(document)
    except Exception as exc:
        result["notes"].append(f"_get_theme_part_info 失敗: {exc}")
        return result

    result["has_theme"] = bool(part_info.get("has_theme"))
    result["theme_part_name"] = part_info.get("theme_part_name")
    if part_info.get("slide_master_count") is not None:
        result["slide_master_count"] = part_info.get("slide_master_count")
    result["notes"].extend(part_info.get("notes", []))

    theme_part_name = result.get("theme_part_name")
    file_path = result.get("file_path")

    if not theme_part_name:
        return result

    if not file_path:
        result["notes"].append("document.file_path 為空，無法從 pptx zip 讀取 theme XML。")
        return result

    if not os.path.exists(file_path):
        result["notes"].append(f"找不到 PPTX 檔案，無法讀取 theme XML: {file_path}")
        return result

    zip_name = str(theme_part_name).lstrip("/")
    try:
        with zipfile.ZipFile(file_path, "r") as zf:
            if zip_name not in zf.namelist():
                result["notes"].append(f"在壓縮包中找不到 theme part: {zip_name}")
                return result
            theme_xml_bytes = zf.read(zip_name)
    except Exception as exc:
        result["notes"].append(f"讀取 theme XML 失敗: {exc}")
        return result

    parsed = _parse_theme_xml(theme_xml_bytes)
    result["theme_name"] = parsed.get("theme_name")
    result["color_scheme"] = parsed.get("color_scheme", {})
    result["font_scheme"] = parsed.get("font_scheme", {})
    result["notes"].extend(parsed.get("notes", []))
    return result


def get_slide_background_info(document: PPTDocument, slide_index: int) -> Dict[str, Any]:
    """
    讀取單頁投影片背景資訊。

    background_type 可能值：inherit, solid, picture, shape_picture_simulated, unknown

    預期欄位（issues 規格）：
    - slide_index, background_type, source, color_rgb, image_ref, image_path_hint,
      by_shape_detection, notes

    判斷策略：
    - 優先使用 slide XML 背景訊號（避免只靠 shape 猜測）
    - 若偵測到滿版圖片 shape，再標示 shape_picture_simulated
    - 無法確定時回傳 unknown 並在 notes 說明依據不足
    """
    _ensure_pptx_available()
    _validate_slide_index(document.prs, slide_index)

    result: Dict[str, Any] = {
        "slide_index": slide_index,
        "background_type": "unknown",
        "source": "unresolved",
        "color_rgb": None,
        "image_ref": None,
        "image_path_hint": None,
        "by_shape_detection": False,
        "notes": [],
    }

    xml_info = _get_slide_background_xml_info(document, slide_index)
    for key in ("background_type", "source", "color_rgb", "image_ref", "image_path_hint", "by_shape_detection"):
        if key in xml_info:
            result[key] = xml_info[key]
    result["notes"].extend(xml_info.get("notes", []))

    shape_info = _detect_full_slide_picture_shape(document, slide_index)
    result["notes"].extend(shape_info.get("notes", []))

    if shape_info.get("matched"):
        if result["background_type"] in ("inherit", "unknown"):
            result["background_type"] = "shape_picture_simulated"
            result["source"] = "shape_detection"
            result["by_shape_detection"] = True
        if not result.get("image_ref"):
            result["image_ref"] = shape_info.get("image_ref")
        if not result.get("image_path_hint"):
            result["image_path_hint"] = shape_info.get("image_path_hint")

    valid_types = {"inherit", "solid", "picture", "shape_picture_simulated", "unknown"}
    if result["background_type"] not in valid_types:
        result["notes"].append(f"偵測到未定義背景類型: {result['background_type']}，已回退為 unknown。")
        result["background_type"] = "unknown"

    return result


def scan_presentation_backgrounds(document: PPTDocument) -> Dict[str, Any]:
    """
    掃描整份簡報每一頁背景。

    回傳：file_path, slide_count, theme_info（get_presentation_theme_info 結果）,
    slides（每頁為 get_slide_background_info 結果）
    """
    _ensure_pptx_available()

    slide_count = len(document.prs.slides)
    theme_info = get_presentation_theme_info(document)
    slides: List[Dict[str, Any]] = []

    for i in range(slide_count):
        try:
            slides.append(get_slide_background_info(document, i))
        except Exception as exc:
            slides.append(
                {
                    "slide_index": i,
                    "background_type": "unknown",
                    "source": "scan_exception",
                    "color_rgb": None,
                    "image_ref": None,
                    "image_path_hint": None,
                    "by_shape_detection": False,
                    "notes": [f"掃描失敗: {exc}"],
                }
            )

    return {
        "file_path": document.file_path,
        "slide_count": slide_count,
        "theme_info": theme_info,
        "slides": slides,
    }


def get_shape_style(
        document: PPTDocument,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    return document.get_shape_style(
        slide_index=slide_index,
        shape_id=shape_id,
        shape_index=shape_index,
    )


def set_shape_style(
        document: PPTDocument,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        fill_transparency: Optional[float] = None,
        line_style: Optional[str] = None,
        line_color: Optional[Tuple[int, int, int]] = None,
        line_width: Optional[int] = None,
    ) -> Dict[str, Any]:
    return document.set_shape_style(
        slide_index=slide_index,
        shape_id=shape_id,
        shape_index=shape_index,
        fill_color=fill_color,
        fill_transparency=fill_transparency,
        line_style=line_style,
        line_color=line_color,
        line_width=line_width,
    )


def set_shape_fill_transparency(
        document: PPTDocument,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        fill_transparency: float = 0.4,
        fill_color: Optional[Tuple[int, int, int]] = None,
    ) -> Dict[str, Any]:
    return document.set_shape_fill_transparency(
        slide_index=slide_index,
        shape_id=shape_id,
        shape_index=shape_index,
        fill_transparency=fill_transparency,
        fill_color=fill_color,
    )


def delete_slide(document: PPTDocument, slide_index: int) -> Dict[str, Any]:
    return document.delete_slide(slide_index)


def duplicate_slide(document: PPTDocument, slide_index: int) -> Dict[str, Any]:
    return document.duplicate_slide(slide_index)


def replace_text(
        document: PPTDocument,
        old_text: str,
        new_text: str,
        slide_indices: Optional[List[int]] = None,
        exact_match: bool = False,
        case_sensitive: bool = True,
    ) -> Dict[str, Any]:
    return document.replace_text(
        old_text=old_text,
        new_text=new_text,
        slide_indices=slide_indices,
        exact_match=exact_match,
        case_sensitive=case_sensitive,
    )


def add_bullets(
        document: PPTDocument,
        slide_index: int,
        items: List[str],
        left: int,
        top: int,
        width: int,
        height: int,
        font_size: int = 20,
        level: int = 0,
        bold: bool = False,
        font_name: Optional[str] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
    ) -> Dict[str, Any]:
    return document.add_bullets(
        slide_index=slide_index,
        items=items,
        left=left,
        top=top,
        width=width,
        height=height,
        font_size=font_size,
        level=level,
        bold=bold,
        font_name=font_name,
        font_color=font_color,
    )


def add_title_slide(
        document: PPTDocument,
        title: str,
        subtitle: str = "",
    ) -> Dict[str, Any]:
    return document.add_title_slide(title=title, subtitle=subtitle)


def reorder_slides(document: PPTDocument, new_order: List[int]) -> Dict[str, Any]:
    return document.reorder_slides(new_order)


def set_slides_background_color(
        document: PPTDocument,
        slide_indices: List[int],
        rgb: Tuple[int, int, int],
    ) -> List[Dict[str, Any]]:
    return document.set_slides_background_color(slide_indices=slide_indices, rgb=rgb)


def set_slides_background_image(
        document: PPTDocument,
        slide_indices: List[int],
        image_path: str,
    ) -> List[Dict[str, Any]]:
    return document.set_slides_background_image(slide_indices=slide_indices, image_path=image_path)


def _find_libreoffice_executable(explicit_path: Optional[str] = None) -> str:
    """
    尋找 LibreOffice / soffice 可執行檔
    """
    candidates = []

    if explicit_path:
        candidates.append(explicit_path)

    # PATH 中常見名稱
    for name in ["soffice", "libreoffice", "soffice.exe"]:
        found = shutil.which(name)
        if found:
            candidates.append(found)

    # Windows 常見安裝路徑
    candidates.extend([
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ])

    for path in candidates:
        if path and os.path.exists(path):
            return os.path.abspath(path)

    raise FileNotFoundError("找不到 LibreOffice/soffice，可明確傳入 libreoffice_path")


def convert_pptx_to_pdf(
        pptx_path: str,
        output_dir: Optional[str] = None,
        libreoffice_path: Optional[str] = None,
        overwrite: bool = True,
    ) -> str:
    """
    用 LibreOffice headless 將 PPTX 轉成 PDF
    """
    pptx_path = os.path.abspath(pptx_path)
    if not os.path.exists(pptx_path):
        raise FileNotFoundError(f"找不到 PPTX: {pptx_path}")

    soffice = _find_libreoffice_executable(libreoffice_path)

    if output_dir is None:
        output_dir = os.path.dirname(pptx_path) or os.getcwd()
    output_dir = os.path.abspath(output_dir)
    os.makedirs(output_dir, exist_ok=True)

    expected_pdf = os.path.join(
        output_dir,
        f"{Path(pptx_path).stem}.pdf"
    )

    if os.path.exists(expected_pdf) and overwrite:
        os.remove(expected_pdf)

    cmd = [
        soffice,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        pptx_path,
    ]

    proc = subprocess.run(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        check=False,
    )

    if proc.returncode != 0:
        raise RuntimeError(
            f"LibreOffice 轉 PDF 失敗\n"
            f"cmd={' '.join(cmd)}\n"
            f"stdout={proc.stdout}\n"
            f"stderr={proc.stderr}"
        )

    if not os.path.exists(expected_pdf):
        raise FileNotFoundError(
            f"LibreOffice 指令執行完畢，但找不到輸出的 PDF: {expected_pdf}\n"
            f"stdout={proc.stdout}\n"
            f"stderr={proc.stderr}"
        )

    return expected_pdf


def render_slide_to_image(
        pptx_path: str,
        slide_index: int,
        output_path: str,
        dpi: int = 150,
        libreoffice_path: Optional[str] = None,
        temp_dir: Optional[str] = None,
    ) -> Dict[str, Any]:
    """
    將指定頁投影片輸出成圖片
    - 先轉 PDF
    - 再用 PyMuPDF 將指定頁 rasterize 成 png
    """
    if slide_index < 0:
        raise ValueError("slide_index 必須 >= 0")

    output_path = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    with tempfile.TemporaryDirectory(dir=temp_dir) as tmpdir:
        pdf_path = convert_pptx_to_pdf(
            pptx_path=pptx_path,
            output_dir=tmpdir,
            libreoffice_path=libreoffice_path,
            overwrite=True,
        )

        doc = fitz.open(pdf_path)
        try:
            if slide_index >= len(doc):
                raise IndexError(f"slide_index 超出範圍: {slide_index}, 總頁數={len(doc)}")

            page = doc[slide_index]
            scale = dpi / 72.0
            matrix = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            pix.save(output_path)

            return {
                "pptx_path": os.path.abspath(pptx_path),
                "pdf_path": pdf_path,
                "slide_index": slide_index,
                "page_count": len(doc),
                "dpi": dpi,
                "output_path": output_path,
                "width": pix.width,
                "height": pix.height,
            }
        finally:
            doc.close()


def render_slides_to_grid_image(
        pptx_path: str,
        slide_indices: List[int],
        output_path: str,
        cols: int = 2,
        dpi: int = 150,
        libreoffice_path: Optional[str] = None,
        temp_dir: Optional[str] = None,
        add_page_title: bool = True,
        figure_title: Optional[str] = None,
    ) -> Dict[str, Any]:
    """
    將多頁投影片輸出並拼成一張 grid 圖
    """
    if not slide_indices:
        raise ValueError("slide_indices 不可為空")
    if cols < 1:
        raise ValueError("cols 必須 >= 1")

    _activate_matplotlib_chinese_font()

    output_path = os.path.abspath(output_path)
    os.makedirs(os.path.dirname(output_path), exist_ok=True)

    with tempfile.TemporaryDirectory(dir=temp_dir) as tmpdir:
        pdf_path = convert_pptx_to_pdf(
            pptx_path=pptx_path,
            output_dir=tmpdir,
            libreoffice_path=libreoffice_path,
            overwrite=True,
        )

        doc = fitz.open(pdf_path)
        image_paths = []
        try:
            page_count = len(doc)

            for slide_index in slide_indices:
                if slide_index < 0 or slide_index >= page_count:
                    raise IndexError(f"slide_index 超出範圍: {slide_index}, 總頁數={page_count}")

            # 先逐頁輸出暫存 png
            for slide_index in slide_indices:
                page = doc[slide_index]
                scale = dpi / 72.0
                matrix = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=matrix, alpha=False)

                img_path = os.path.join(tmpdir, f"slide_{slide_index}.png")
                pix.save(img_path)
                image_paths.append(img_path)
        finally:
            doc.close()

        rows = math.ceil(len(image_paths) / cols)
        fig, axes = plt.subplots(rows, cols, figsize=(cols * 4.5, rows * 3.2), squeeze=False)

        for ax in axes.flat:
            ax.axis("off")

        for i, img_path in enumerate(image_paths):
            r = i // cols
            c = i % cols
            ax = axes[r][c]
            img = mpimg.imread(img_path)
            ax.imshow(img)
            ax.axis("off")
            if add_page_title:
                ax.set_title(f"Slide {slide_indices[i]}", fontsize=10)

        if figure_title:
            fig.suptitle(figure_title)

        plt.tight_layout()
        fig.savefig(output_path, dpi=dpi, bbox_inches="tight")
        plt.close(fig)

        return {
            "pptx_path": os.path.abspath(pptx_path),
            "slide_indices": slide_indices,
            "count": len(slide_indices),
            "cols": cols,
            "rows": rows,
            "dpi": dpi,
            "output_path": output_path,
        }


if __name__ == "__main__":
    # 簡單測試
    out = new("test_output/demo_new_ppt", plank_page_num=2, plank_page_width=1080, plank_page_height=1920)
    doc = open_presentation(out)

    add_text(
        document=doc,
        slide_index=0,
        text="Hello PPT Core",
        left=1000000,
        top=500000,
        width=5000000,
        height=800000,
        font_size=24,
        bold=True,
    )

    add_table(
        document=doc,
        slide_index=1,
        rows=3,
        cols=3,
        left=800000,
        top=800000,
        width=6000000,
        height=2000000,
        data=[
            ["欄位A", "欄位B", "欄位C"],
            [1, 2, 3],
            [4, 5, 6],
        ],
        first_row_as_header=True,
    )

    add_title_slide(doc, "主標題", "副標題")
    add_bullets(
        document=doc,
        slide_index=0,
        items=["第一點", "第二點", "第三點"],
        left=800000,
        top=1200000,
        width=5000000,
        height=2000000,
        font_size=24,
    )
    replace_text(doc, "第二點", "第二點（已更新）")
    duplicate_slide(doc, 0)
    reorder_slides(doc, [1, 0, 2, 3])

    save(doc)
    print(get_info(doc))
