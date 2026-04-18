# -*- coding: utf-8 -*-
"""
PPT 表格進階操作（供 ppt_stdio.PPTDocument 委派）。

設計說明（與 issues/操作表格的功能.iss 對齊）：
- 列／欄「結構」增刪：python-pptx 無穩定高階 API 時，以 rebuild_table_with_modified_structure
  讀取舊表文字與部分樣式後刪除舊 shape、再以 add_table 重建並回寫。
- 儲存格邊框：優先使用 python-pptx 的 cell.border_*（_LineFormat）；若環境不支援則回傳 notes，不默默失敗。
- 本模組不 import ppt_stdio，避免循環匯入；僅假設 document 具 prs 與 _get_shape。
"""

from __future__ import annotations

from typing import Any, Dict, List, Optional, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# 不同 python-pptx 版本：舊版為 MSO_ANCHORING；0.6.23+ 為 MSO_VERTICAL_ANCHOR，並提供別名 MSO_ANCHOR
try:
    from pptx.enum.text import MSO_ANCHORING  # type: ignore
except ImportError:
    try:
        from pptx.enum.text import MSO_ANCHOR as MSO_ANCHORING  # type: ignore
    except ImportError:
        from pptx.enum.text import MSO_VERTICAL_ANCHOR as MSO_ANCHORING  # type: ignore


def _validate_slide(document: Any, slide_index: int) -> None:
    prs = document.prs
    if slide_index < 0 or slide_index >= len(prs.slides):
        raise IndexError(f"slide_index 超出範圍: {slide_index}, 總頁數={len(prs.slides)}")


def _rgb_tuple_to_color(rgb: Optional[Tuple[int, int, int]]) -> Optional[RGBColor]:
    if rgb is None:
        return None
    if len(rgb) != 3:
        raise ValueError("rgb 顏色格式必須是 (R, G, B)")
    r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    for v in (r, g, b):
        if v < 0 or v > 255:
            raise ValueError("rgb 每個值都必須介於 0~255")
    return RGBColor(r, g, b)


def _rgb_from_color_obj(color_obj: Any) -> Optional[List[int]]:
    if color_obj is None:
        return None
    try:
        rgb = color_obj.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    hex_value = str(rgb).strip().lstrip("#")
    if len(hex_value) != 6:
        return None
    try:
        return [int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16)]
    except Exception:
        return None


def _parse_h_align(name: Optional[str]) -> Optional[Any]:
    if name is None:
        return None
    key = str(name).strip().lower()
    mapping = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    if key not in mapping:
        raise ValueError("h_align 必須是 left/center/right/justify")
    return mapping[key]


def _parse_v_align(name: Optional[str]) -> Optional[Any]:
    if name is None:
        return None
    key = str(name).strip().lower()
    mapping = {
        "top": MSO_ANCHORING.TOP,
        "middle": MSO_ANCHORING.MIDDLE,
        "bottom": MSO_ANCHORING.BOTTOM,
    }
    if key not in mapping:
        raise ValueError("v_align 必須是 top/middle/bottom")
    return mapping[key]


def _get_table_shape(document: Any, slide_index: int, shape_id: Optional[int], shape_index: Optional[int]):
    _validate_slide(document, slide_index)
    shape, resolved_index = document._get_shape(
        slide_index=slide_index,
        shape_id=shape_id,
        shape_index=shape_index,
    )
    if not getattr(shape, "has_table", False):
        raise ValueError(
            f"指定 shape 不是表格：shape_id={getattr(shape, 'shape_id', None)}, shape_index={resolved_index}"
        )
    return shape, shape.table, resolved_index


def _validate_cell(table: Any, row_idx: int, col_idx: int) -> None:
    nrows = len(table.rows)
    ncols = len(table.columns)
    if row_idx < 0 or row_idx >= nrows:
        raise IndexError(f"row_idx 超出範圍: {row_idx}, rows={nrows}")
    if col_idx < 0 or col_idx >= ncols:
        raise IndexError(f"col_idx 超出範圍: {col_idx}, cols={ncols}")


def _validate_row(table: Any, row_idx: int) -> None:
    nrows = len(table.rows)
    if row_idx < 0 or row_idx >= nrows:
        raise IndexError(f"row_idx 超出範圍: {row_idx}, rows={nrows}")


def _validate_col(table: Any, col_idx: int) -> None:
    ncols = len(table.columns)
    if col_idx < 0 or col_idx >= ncols:
        raise IndexError(f"col_idx 超出範圍: {col_idx}, cols={ncols}")


def _serialize_cell(cell: Any) -> Dict[str, Any]:
    notes: List[str] = []
    text = ""
    try:
        text = cell.text_frame.text or ""
    except Exception as exc:
        notes.append(f"讀取儲存格文字失敗: {exc}")

    fill_rgb: Optional[List[int]] = None
    try:
        fill = cell.fill
        fill_type = str(getattr(fill, "type", "") or "").lower()
        if "solid" in fill_type or getattr(fill, "fore_color", None) is not None:
            fill_rgb = _rgb_from_color_obj(getattr(fill, "fore_color", None))
    except Exception as exc:
        notes.append(f"讀取儲存格底色失敗: {exc}")

    h_align = None
    v_align = None
    try:
        p = cell.text_frame.paragraphs[0]
        al = getattr(p, "alignment", None)
        if al is not None:
            h_align = str(al).lower()
    except Exception as exc:
        notes.append(f"讀取水平對齊失敗: {exc}")

    try:
        va = getattr(cell, "vertical_anchor", None)
        if va is not None:
            v_align = str(va).lower()
    except Exception as exc:
        notes.append(f"讀取垂直對齊失敗: {exc}")

    borders: Dict[str, Any] = {}
    for side in ("top", "bottom", "left", "right"):
        attr = f"border_{side}"
        try:
            line = getattr(cell, attr, None)
            if line is None:
                continue
            borders[side] = {
                "color": _rgb_from_color_obj(getattr(line, "color", None)),
                "width_pt": float(line.width.pt) if getattr(line, "width", None) is not None else None,
            }
        except Exception as exc:
            notes.append(f"讀取邊框 {side} 失敗: {exc}")

    return {
        "text": text,
        "fill_rgb": fill_rgb,
        "h_align": h_align,
        "v_align": v_align,
        "borders": borders,
        "notes": notes,
    }


def _serialize_table(shape: Any, table: Any, slide_index: int, shape_index: int) -> Dict[str, Any]:
    nrows = len(table.rows)
    ncols = len(table.columns)
    cells: List[List[Dict[str, Any]]] = []
    for r in range(nrows):
        row_cells: List[Dict[str, Any]] = []
        for c in range(ncols):
            row_cells.append(_serialize_cell(table.cell(r, c)))
        cells.append(row_cells)

    return {
        "slide_index": slide_index,
        "shape_index": shape_index,
        "shape_id": getattr(shape, "shape_id", None),
        "rows": nrows,
        "cols": ncols,
        "left": int(shape.left),
        "top": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
        "cells": cells,
    }


def _apply_cell_text_style(
        cell: Any,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
    ) -> List[str]:
    notes: List[str] = []
    tf = cell.text_frame
    if len(tf.paragraphs) == 0:
        tf.text = ""
    p = tf.paragraphs[0]
    if len(p.runs) == 0:
        p.text = p.text or ""
    for run in p.runs:
        if font_size is not None:
            run.font.size = Pt(int(font_size))
        if bold is not None:
            run.font.bold = bool(bold)
        if italic is not None:
            run.font.italic = bool(italic)
        if font_name is not None:
            run.font.name = str(font_name)
        color = _rgb_tuple_to_color(font_color)
        if color is not None:
            run.font.color.rgb = color
    return notes


def _apply_cell_fill(cell: Any, fill_color: Optional[Tuple[int, int, int]]) -> List[str]:
    notes: List[str] = []
    if fill_color is None:
        return notes
    try:
        cell.fill.solid()
        c = _rgb_tuple_to_color(fill_color)
        if c is not None:
            cell.fill.fore_color.rgb = c
    except Exception as exc:
        notes.append(f"設定儲存格底色失敗: {exc}")
    return notes


def _apply_cell_alignment(cell: Any, h_align: Optional[str], v_align: Optional[str]) -> List[str]:
    notes: List[str] = []
    ha = _parse_h_align(h_align) if h_align is not None else None
    va = _parse_v_align(v_align) if v_align is not None else None
    try:
        if ha is not None:
            for p in cell.text_frame.paragraphs:
                p.alignment = ha
    except Exception as exc:
        notes.append(f"設定水平對齊失敗: {exc}")
    try:
        if va is not None:
            cell.vertical_anchor = va
    except Exception as exc:
        notes.append(f"設定垂直對齊失敗: {exc}")
    return notes


def _apply_cell_border_style(
        cell: Any,
        border_color: Optional[Tuple[int, int, int]] = None,
        border_width: Optional[float] = None,
        border_style: Optional[str] = None,
        scope: str = "all",
    ) -> List[str]:
    """
    scope: all | outer | inner
    inner：僅儲存格內側（實作上以四邊皆設相同樣式近似；細緻內框需 OOXML 再強化）。
    """
    notes: List[str] = []
    if border_color is None and border_width is None and border_style is None:
        return notes

    if border_style is not None and str(border_style).strip().lower() not in ("solid", "none", ""):
        notes.append("border_style 目前僅完整驗證 solid；其他樣式將嘗試套用 dash 可能因版本而失敗。")

    sides = ["top", "bottom", "left", "right"]
    if scope == "outer":
        notes.append("outer 模式仍以四邊線條實作；若需僅外框請後續改 OOXML tcPr。")

    color = _rgb_tuple_to_color(border_color) if border_color is not None else None
    width_pt = float(border_width) if border_width is not None else 1.0

    for side in sides:
        line = getattr(cell, f"border_{side}", None)
        if line is None:
            continue
        try:
            if color is not None:
                line.color.rgb = color
            line.width = Pt(width_pt)
        except Exception as exc:
            notes.append(f"設定邊框 {side} 失敗: {exc}")
    return notes


def list_slide_tables(document: Any, slide_index: int) -> Dict[str, Any]:
    """列出指定頁所有表格摘要。"""
    _validate_slide(document, slide_index)
    slide = document.prs.slides[slide_index]
    tables: List[Dict[str, Any]] = []
    for idx, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        tables.append({
            "shape_index": idx,
            "shape_id": getattr(shape, "shape_id", None),
            "name": getattr(shape, "name", None),
            "rows": len(table.rows),
            "cols": len(table.columns),
            "left": int(shape.left),
            "top": int(shape.top),
            "width": int(shape.width),
            "height": int(shape.height),
        })
    return {"slide_index": slide_index, "tables": tables, "count": len(tables)}


def get_table_detail(
        document: Any,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    """取得指定表格詳細資訊（含儲存格序列化）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    detail = _serialize_table(shape, table, slide_index, resolved_index)
    return {"table": detail, "notes": []}


def update_table_cell(
        document: Any,
        slide_index: int,
        row_idx: int,
        col_idx: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        text: Optional[str] = None,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        h_align: Optional[str] = None,
        v_align: Optional[str] = None,
        border_color: Optional[Tuple[int, int, int]] = None,
        border_width: Optional[float] = None,
        border_style: Optional[str] = None,
        clear_text: bool = False,
    ) -> Dict[str, Any]:
    """更新單一儲存格文字與／或樣式。clear_text=True 時清空文字。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_cell(table, row_idx, col_idx)
    cell = table.cell(row_idx, col_idx)
    notes: List[str] = []

    if clear_text:
        cell.text = ""
        notes.append("已清空儲存格文字。")
    elif text is not None:
        cell.text = str(text)

    notes.extend(_apply_cell_fill(cell, fill_color))
    notes.extend(_apply_cell_text_style(cell, font_name, font_size, bold, italic, font_color))
    notes.extend(_apply_cell_alignment(cell, h_align, v_align))
    notes.extend(_apply_cell_border_style(cell, border_color, border_width, border_style, scope="all"))

    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "row_idx": row_idx,
        "col_idx": col_idx,
        "notes": notes,
    }


def set_table_cell_style(
        document: Any,
        slide_index: int,
        row_idx: int,
        col_idx: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        h_align: Optional[str] = None,
        v_align: Optional[str] = None,
        border_color: Optional[Tuple[int, int, int]] = None,
        border_width: Optional[float] = None,
        border_style: Optional[str] = None,
    ) -> Dict[str, Any]:
    """僅設定儲存格樣式，不改文字。"""
    return update_table_cell(
        document=document,
        slide_index=slide_index,
        row_idx=row_idx,
        col_idx=col_idx,
        shape_id=shape_id,
        shape_index=shape_index,
        text=None,
        font_name=font_name,
        font_size=font_size,
        bold=bold,
        italic=italic,
        font_color=font_color,
        fill_color=fill_color,
        h_align=h_align,
        v_align=v_align,
        border_color=border_color,
        border_width=border_width,
        border_style=border_style,
        clear_text=False,
    )


def update_table_row(
        document: Any,
        slide_index: int,
        row_idx: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        row_text: Optional[str] = None,
        cell_texts: Optional[List[str]] = None,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        h_align: Optional[str] = None,
        v_align: Optional[str] = None,
    ) -> Dict[str, Any]:
    """更新整列文字與／或樣式。cell_texts 若提供，長度須等於欄數。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_row(table, row_idx)
    ncols = len(table.columns)
    notes: List[str] = []

    if cell_texts is not None:
        if len(cell_texts) != ncols:
            raise ValueError(f"cell_texts 長度必須等於欄數 {ncols}")
        texts = [str(x) for x in cell_texts]
    elif row_text is not None:
        texts = [str(row_text)] * ncols
    else:
        texts = [None] * ncols

    for c in range(ncols):
        t = texts[c]
        update_table_cell(
            document=document,
            slide_index=slide_index,
            row_idx=row_idx,
            col_idx=c,
            shape_id=getattr(shape, "shape_id", None),
            shape_index=resolved_index,
            text=t,
            font_name=font_name,
            font_size=font_size,
            bold=bold,
            italic=italic,
            font_color=font_color,
            fill_color=fill_color,
            h_align=h_align,
            v_align=v_align,
        )

    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "row_idx": row_idx,
        "notes": notes,
    }


def update_table_column(
        document: Any,
        slide_index: int,
        col_idx: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        column_text: Optional[str] = None,
        cell_texts: Optional[List[str]] = None,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        h_align: Optional[str] = None,
        v_align: Optional[str] = None,
    ) -> Dict[str, Any]:
    """更新整欄文字與／或樣式。cell_texts 若提供，長度須等於列數。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_col(table, col_idx)
    nrows = len(table.rows)
    notes: List[str] = []

    if cell_texts is not None:
        if len(cell_texts) != nrows:
            raise ValueError(f"cell_texts 長度必須等於列數 {nrows}")
        texts = [str(x) for x in cell_texts]
    elif column_text is not None:
        texts = [str(column_text)] * nrows
    else:
        texts = [None] * nrows

    for r in range(nrows):
        t = texts[r]
        update_table_cell(
            document=document,
            slide_index=slide_index,
            row_idx=r,
            col_idx=col_idx,
            shape_id=getattr(shape, "shape_id", None),
            shape_index=resolved_index,
            text=t,
            font_name=font_name,
            font_size=font_size,
            bold=bold,
            italic=italic,
            font_color=font_color,
            fill_color=fill_color,
            h_align=h_align,
            v_align=v_align,
        )

    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "col_idx": col_idx,
        "notes": notes,
    }


def set_table_row_style(
        document: Any,
        slide_index: int,
        row_idx: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        h_align: Optional[str] = None,
        v_align: Optional[str] = None,
        border_color: Optional[Tuple[int, int, int]] = None,
        border_width: Optional[float] = None,
        border_style: Optional[str] = None,
    ) -> Dict[str, Any]:
    """設定整列樣式（不含文字變更）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_row(table, row_idx)
    for c in range(len(table.columns)):
        set_table_cell_style(
            document=document,
            slide_index=slide_index,
            row_idx=row_idx,
            col_idx=c,
            shape_id=getattr(shape, "shape_id", None),
            shape_index=resolved_index,
            font_name=font_name,
            font_size=font_size,
            bold=bold,
            italic=italic,
            font_color=font_color,
            fill_color=fill_color,
            h_align=h_align,
            v_align=v_align,
            border_color=border_color,
            border_width=border_width,
            border_style=border_style,
        )
    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "row_idx": row_idx,
        "notes": [],
    }


def set_table_column_style(
        document: Any,
        slide_index: int,
        col_idx: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        fill_color: Optional[Tuple[int, int, int]] = None,
        h_align: Optional[str] = None,
        v_align: Optional[str] = None,
        border_color: Optional[Tuple[int, int, int]] = None,
        border_width: Optional[float] = None,
        border_style: Optional[str] = None,
    ) -> Dict[str, Any]:
    """設定整欄樣式（不含文字變更）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_col(table, col_idx)
    for r in range(len(table.rows)):
        set_table_cell_style(
            document=document,
            slide_index=slide_index,
            row_idx=r,
            col_idx=col_idx,
            shape_id=getattr(shape, "shape_id", None),
            shape_index=resolved_index,
            font_name=font_name,
            font_size=font_size,
            bold=bold,
            italic=italic,
            font_color=font_color,
            fill_color=fill_color,
            h_align=h_align,
            v_align=v_align,
            border_color=border_color,
            border_width=border_width,
            border_style=border_style,
        )
    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "col_idx": col_idx,
        "notes": [],
    }


def set_table_row_height(
        document: Any,
        slide_index: int,
        row_idx: int,
        height_emu: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    """設定列高（EMU）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_row(table, row_idx)
    if height_emu <= 0:
        raise ValueError("height_emu 必須 > 0")
    table.rows[row_idx].height = Emu(int(height_emu))
    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "row_idx": row_idx,
        "height_emu": int(height_emu),
        "notes": [],
    }


def set_table_column_width(
        document: Any,
        slide_index: int,
        col_idx: int,
        width_emu: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    """設定欄寬（EMU）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_col(table, col_idx)
    if width_emu <= 0:
        raise ValueError("width_emu 必須 > 0")
    table.columns[col_idx].width = Emu(int(width_emu))
    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "col_idx": col_idx,
        "width_emu": int(width_emu),
        "notes": [],
    }


def distribute_table_column_widths(
        document: Any,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        column_indices: Optional[List[int]] = None,
    ) -> Dict[str, Any]:
    """
    平均分配欄寬。column_indices 為 None 時重分所有欄；否則僅重分指定欄，寬度來自表格總寬扣除未列欄目前寬度。
    """
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    ncols = len(table.columns)
    total = int(shape.width)
    if column_indices is None:
        idxs = list(range(ncols))
    else:
        idxs = sorted({int(i) for i in column_indices})
        for i in idxs:
            _validate_col(table, i)

    other_sum = sum(int(table.columns[i].width) for i in range(ncols) if i not in idxs)
    remaining = total - other_sum
    if remaining <= 0:
        raise ValueError("剩餘可分配寬度不足，請先縮小未列欄或使用較大表格寬度")
    each = int(remaining // len(idxs)) if idxs else 0
    if each <= 0:
        raise ValueError("計算後每欄寬度無效")

    for i in idxs:
        table.columns[i].width = Emu(each)

    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "column_indices": idxs,
        "assigned_width_emu": each,
        "notes": [],
    }


def distribute_table_row_heights(
        document: Any,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        row_indices: Optional[List[int]] = None,
    ) -> Dict[str, Any]:
    """平均分配列高。row_indices 為 None 時重分所有列。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    nrows = len(table.rows)
    total = int(shape.height)
    if row_indices is None:
        idxs = list(range(nrows))
    else:
        idxs = sorted({int(i) for i in row_indices})
        for i in idxs:
            _validate_row(table, i)

    other_sum = sum(int(table.rows[i].height) for i in range(nrows) if i not in idxs)
    remaining = total - other_sum
    if remaining <= 0:
        raise ValueError("剩餘可分配高度不足，請先縮小未列列或使用較大表格高度")
    each = int(remaining // len(idxs)) if idxs else 0
    if each <= 0:
        raise ValueError("計算後每列高度無效")

    for i in idxs:
        table.rows[i].height = Emu(each)

    return {
        "slide_index": slide_index,
        "shape_index": resolved_index,
        "shape_id": getattr(shape, "shape_id", None),
        "row_indices": idxs,
        "assigned_height_emu": each,
        "notes": [],
    }


def delete_table(
        document: Any,
        slide_index: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
    ) -> Dict[str, Any]:
    """刪除整個表格 shape。"""
    from dataProcess.ppt_stdio import delete_shape

    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    sid = getattr(shape, "shape_id", None)
    result = delete_shape(
        document=document,
        slide_index=slide_index,
        shape_id=sid,
        shape_index=resolved_index,
    )
    return {
        "deleted": True,
        "slide_index": slide_index,
        "shape_id": sid,
        "result": result,
        "notes": [],
    }


def rebuild_table_with_modified_structure(
        document: Any,
        slide_index: int,
        new_rows: int,
        new_cols: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        first_row_as_header: bool = False,
        font_size: int = 14,
    ) -> Dict[str, Any]:
    """
    以新列欄數重建表格：讀取舊儲存格文字（樣式簡化不回寫），刪除舊 shape，再 add_table 並填入可對應範圍文字。
    """
    if new_rows < 1 or new_cols < 1:
        raise ValueError("new_rows 與 new_cols 都必須 >= 1")

    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    old_rows = len(table.rows)
    old_cols = len(table.columns)
    geom = (int(shape.left), int(shape.top), int(shape.width), int(shape.height))
    left, top, width, height = geom

    matrix: List[List[str]] = []
    for r in range(old_rows):
        row_vals: List[str] = []
        for c in range(old_cols):
            try:
                row_vals.append(table.cell(r, c).text_frame.text or "")
            except Exception:
                row_vals.append("")
        matrix.append(row_vals)

    delete_table(document, slide_index, shape_id=getattr(shape, "shape_id", None), shape_index=resolved_index)

    data: List[List[str]] = []
    for r in range(new_rows):
        row_vals = []
        for c in range(new_cols):
            if r < len(matrix) and c < len(matrix[r]):
                row_vals.append(matrix[r][c])
            else:
                row_vals.append("")
        data.append(row_vals)

    created = document.add_table(
        slide_index=slide_index,
        rows=new_rows,
        cols=new_cols,
        left=left,
        top=top,
        width=width,
        height=height,
        data=data,
        first_row_as_header=first_row_as_header,
        font_size=font_size,
    )

    notes = [
        "rebuild_table_with_modified_structure：已保留可對應範圍之文字；細部儲存格樣式未完整回寫（可後續擴充）。",
    ]
    return {
        "slide_index": slide_index,
        "old_rows": old_rows,
        "old_cols": old_cols,
        "new_rows": new_rows,
        "new_cols": new_cols,
        "created": created,
        "notes": notes,
    }


def insert_table_row(
        document: Any,
        slide_index: int,
        insert_before: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        first_row_as_header: bool = False,
        font_size: int = 14,
    ) -> Dict[str, Any]:
    """插入一列（以重建表格實作）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    nrows = len(table.rows)
    if insert_before < 0 or insert_before > nrows:
        raise IndexError(f"insert_before 超出範圍: {insert_before}, rows={nrows}")

    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)

    matrix: List[List[str]] = []
    for r in range(nrows):
        row_vals = []
        for c in range(len(table.columns)):
            row_vals.append(table.cell(r, c).text_frame.text or "")
        matrix.append(row_vals)

    new_matrix = matrix[:insert_before] + [[""] * len(matrix[0])] + matrix[insert_before:]

    delete_table(document, slide_index, shape_id=getattr(shape, "shape_id", None), shape_index=resolved_index)

    new_rows = len(new_matrix)
    new_cols = len(new_matrix[0]) if new_matrix else 1

    created = document.add_table(
        slide_index=slide_index,
        rows=new_rows,
        cols=new_cols,
        left=left,
        top=top,
        width=width,
        height=height,
        data=new_matrix,
        first_row_as_header=first_row_as_header,
        font_size=font_size,
    )
    return {
        "slide_index": slide_index,
        "insert_before": insert_before,
        "created": created,
        "notes": ["insert_table_row：以重建表格實作。"],
    }


def delete_table_row(
        document: Any,
        slide_index: int,
        row_idx: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        first_row_as_header: bool = False,
        font_size: int = 14,
    ) -> Dict[str, Any]:
    """刪除一列（以重建表格實作）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_row(table, row_idx)
    nrows = len(table.rows)
    ncols = len(table.columns)
    if nrows <= 1:
        raise ValueError("表格至少需保留一列")

    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)

    matrix: List[List[str]] = []
    for r in range(nrows):
        if r == row_idx:
            continue
        row_vals = []
        for c in range(ncols):
            row_vals.append(table.cell(r, c).text_frame.text or "")
        matrix.append(row_vals)

    delete_table(document, slide_index, shape_id=getattr(shape, "shape_id", None), shape_index=resolved_index)

    created = document.add_table(
        slide_index=slide_index,
        rows=len(matrix),
        cols=ncols,
        left=left,
        top=top,
        width=width,
        height=height,
        data=matrix,
        first_row_as_header=first_row_as_header,
        font_size=font_size,
    )
    return {"slide_index": slide_index, "deleted_row": row_idx, "created": created, "notes": ["delete_table_row：以重建表格實作。"]}


def insert_table_column(
        document: Any,
        slide_index: int,
        insert_before: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        first_row_as_header: bool = False,
        font_size: int = 14,
    ) -> Dict[str, Any]:
    """插入一欄（以重建表格實作）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    nrows = len(table.rows)
    ncols = len(table.columns)
    if insert_before < 0 or insert_before > ncols:
        raise IndexError(f"insert_before 超出範圍: {insert_before}, cols={ncols}")

    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)

    matrix: List[List[str]] = []
    for r in range(nrows):
        row_vals = []
        for c in range(ncols):
            row_vals.append(table.cell(r, c).text_frame.text or "")
        matrix.append(row_vals)

    for r in range(nrows):
        matrix[r] = matrix[r][:insert_before] + [""] + matrix[r][insert_before:]

    delete_table(document, slide_index, shape_id=getattr(shape, "shape_id", None), shape_index=resolved_index)

    created = document.add_table(
        slide_index=slide_index,
        rows=nrows,
        cols=len(matrix[0]),
        left=left,
        top=top,
        width=width,
        height=height,
        data=matrix,
        first_row_as_header=first_row_as_header,
        font_size=font_size,
    )
    return {"slide_index": slide_index, "insert_before": insert_before, "created": created, "notes": ["insert_table_column：以重建表格實作。"]}


def delete_table_column(
        document: Any,
        slide_index: int,
        col_idx: int,
        shape_id: Optional[int] = None,
        shape_index: Optional[int] = None,
        first_row_as_header: bool = False,
        font_size: int = 14,
    ) -> Dict[str, Any]:
    """刪除一欄（以重建表格實作）。"""
    shape, table, resolved_index = _get_table_shape(document, slide_index, shape_id, shape_index)
    _validate_col(table, col_idx)
    nrows = len(table.rows)
    ncols = len(table.columns)
    if ncols <= 1:
        raise ValueError("表格至少需保留一欄")

    left = int(shape.left)
    top = int(shape.top)
    width = int(shape.width)
    height = int(shape.height)

    matrix: List[List[str]] = []
    for r in range(nrows):
        row_vals = []
        for c in range(ncols):
            if c == col_idx:
                continue
            row_vals.append(table.cell(r, c).text_frame.text or "")
        matrix.append(row_vals)

    delete_table(document, slide_index, shape_id=getattr(shape, "shape_id", None), shape_index=resolved_index)

    created = document.add_table(
        slide_index=slide_index,
        rows=nrows,
        cols=ncols - 1,
        left=left,
        top=top,
        width=width,
        height=height,
        data=matrix,
        first_row_as_header=first_row_as_header,
        font_size=font_size,
    )
    return {"slide_index": slide_index, "deleted_col": col_idx, "created": created, "notes": ["delete_table_column：以重建表格實作。"]}
